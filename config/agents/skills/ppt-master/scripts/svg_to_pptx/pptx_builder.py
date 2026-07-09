"""Core PPTX assembly: create_pptx_with_native_svg."""

from __future__ import annotations

import hashlib
import json
import mimetypes
import os
import re
import posixpath
import shutil
import stat
import subprocess
import tempfile
import uuid
import zipfile
from concurrent.futures import ProcessPoolExecutor, as_completed
from datetime import datetime, timezone
from pathlib import Path
from typing import Any
from xml.sax.saxutils import escape

from pptx import Presentation
from pptx.util import Emu

from .drawingml_converter import convert_svg_to_slide_shapes
from .pptx_dimensions import (
    CANVAS_FORMATS,
    get_slide_dimensions, get_pixel_dimensions,
    get_viewbox_dimensions, detect_format_from_svg,
)
from .pptx_media import (
    PNG_RENDERER,
    get_png_renderer_info, convert_svg_to_png, convert_svg_to_png_cached,
)
from .pptx_notes import (
    markdown_to_plain_text,
    create_notes_master_rels_xml,
    create_notes_master_xml,
    create_notes_slide_xml,
    create_notes_slide_rels_xml,
)
from .pptx_narration import (
    AUDIO_CONTENT_TYPES,
    AUDIO_REL_TYPE,
    IMAGE_REL_TYPE,
    MEDIA_REL_TYPE,
    TRANSPARENT_PNG_BYTES,
    apply_recorded_timing,
    inject_narration,
    next_shape_id,
    probe_audio_duration,
)
from .pptx_slide_xml import (
    ANIMATIONS_AVAILABLE, TRANSITIONS,
    create_slide_xml_with_svg, create_slide_rels_xml,
)

# Re-import create_transition_xml only if available
try:
    from pptx_animations import (
        create_transition_xml,
        create_sequence_timing_xml,
        pick_animation_effect,
    )
except ImportError:
    create_transition_xml = None
    create_sequence_timing_xml = None
    pick_animation_effect = None


def _append_relationship(
    rels_path: Path,
    rel_type: str,
    target: str,
) -> str:
    """Append a relationship entry with the next available rId."""
    with open(rels_path, 'r', encoding='utf-8') as f:
        rels_content = f.read()

    rid_numbers = [int(match) for match in re.findall(r'Id="rId(\d+)"', rels_content)]
    next_rid = f'rId{max(rid_numbers, default=0) + 1}'
    rel_xml = (
        f'  <Relationship Id="{next_rid}" '
        f'Type="{rel_type}" Target="{target}"/>'
    )
    rels_content = rels_content.replace(
        '</Relationships>', rel_xml + '\n</Relationships>',
    )

    with open(rels_path, 'w', encoding='utf-8') as f:
        f.write(rels_content)

    return next_rid


def _find_relationship_id(
    rels_path: Path,
    rel_type: str,
    target: str,
) -> str | None:
    """Find an existing relationship id by type and target."""
    if not rels_path.exists():
        return None
    rels_content = rels_path.read_text(encoding='utf-8')
    pattern = (
        r'<Relationship\b[^>]*\bId="([^"]+)"[^>]*'
        rf'\bType="{re.escape(rel_type)}"[^>]*'
        rf'\bTarget="{re.escape(target)}"[^>]*/>'
    )
    match = re.search(pattern, rels_content)
    return match.group(1) if match else None


def _add_default_content_type(content_types: str, extension: str, content_type: str) -> str:
    """Add a Default content type if it is not already present."""
    ext = extension.lstrip(".")
    if f'Extension="{ext}"' in content_types:
        return content_types
    entry = f'  <Default Extension="{ext}" ContentType="{content_type}"/>'
    override_pos = content_types.find('<Override ')
    if override_pos >= 0:
        return content_types[:override_pos] + entry + '\n' + content_types[override_pos:]
    return content_types.replace('</Types>', entry + '\n</Types>')


_IMAGE_CONTENT_TYPES = {
    'png': 'image/png',
    'jpg': 'image/jpeg',
    'jpeg': 'image/jpeg',
    'gif': 'image/gif',
    'webp': 'image/webp',
    'svg': 'image/svg+xml',
    'bmp': 'image/bmp',
    'emf': 'image/x-emf',
    'tif': 'image/tiff',
    'tiff': 'image/tiff',
    'wmf': 'image/x-wmf',
}


def _content_type_for_extension(ext: str) -> str:
    clean = ext.lower().lstrip('.')
    content_type = _IMAGE_CONTENT_TYPES.get(clean) or mimetypes.guess_type(f'x.{clean}')[0]
    if not content_type:
        raise ValueError(f"Unknown media content type for extension: {ext}")
    return content_type


def _as_dict(value: Any) -> dict[str, Any]:
    return value if isinstance(value, dict) else {}


def _create_writable_work_dir(output_path: Path) -> Path:
    """Create a real writable work directory for PPTX assembly."""
    parents = [output_path.parent, Path.cwd(), Path(tempfile.gettempdir())]
    seen: set[str] = set()
    errors: list[str] = []

    for parent in parents:
        parent = parent if str(parent) else Path(".")
        try:
            key = str(parent.resolve())
        except OSError:
            key = str(parent.absolute())
        if key in seen:
            continue
        seen.add(key)

        try:
            parent.mkdir(parents=True, exist_ok=True)
        except OSError as exc:
            errors.append(f"{parent}: cannot create parent ({exc})")
            continue

        for _ in range(3):
            work_dir = parent / f".pptx-build-{os.getpid()}-{uuid.uuid4().hex}"
            try:
                work_dir.mkdir(mode=0o700)
                probe_path = work_dir / ".write-probe"
                probe_path.write_text("ok", encoding="utf-8")
                probe_path.unlink()
                return work_dir
            except OSError as exc:
                errors.append(f"{work_dir}: {exc}")
                shutil.rmtree(work_dir, ignore_errors=True)

    details = "\n  - ".join(errors) if errors else "no candidate directories available"
    raise PermissionError(
        "Unable to create a writable PPTX work directory. "
        "Set the output path to a writable project directory or adjust sandbox permissions. "
        f"Tried:\n  - {details}"
    )


def _relax_output_permissions(output_path: Path) -> list[str]:
    """Make exported files readable outside the sandbox owner where possible."""
    warnings: list[str] = []

    try:
        current_mode = output_path.stat().st_mode
        readable_mode = (
            current_mode
            | stat.S_IRUSR
            | stat.S_IWUSR
            | stat.S_IRGRP
            | stat.S_IROTH
        )
        os.chmod(output_path, readable_mode)
    except OSError as exc:
        warnings.append(f"chmod skipped for {output_path}: {exc}")

    if os.name != 'nt':
        return warnings

    # Windows ACLs can remain sandbox-only even when the file mode looks sane.
    # Grant the built-in Users SID read access; the SID avoids localization
    # issues on non-English Windows installations.
    try:
        result = subprocess.run(
            ['icacls', str(output_path), '/grant', '*S-1-5-32-545:R'],
            capture_output=True,
            text=True,
            check=False,
        )
    except OSError as exc:
        warnings.append(f"icacls skipped for {output_path}: {exc}")
    else:
        if result.returncode != 0:
            message = (result.stderr or result.stdout or '').strip()
            details = f": {message}" if message else ''
            warnings.append(f"icacls failed for {output_path}{details}")

    return warnings


_NOTES_MASTER_REL_TYPE = (
    'http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesMaster'
)


def _ensure_notes_master(extract_dir: Path) -> None:
    """Create notesMaster parts and wire them into the presentation package."""
    ppt_dir = extract_dir / 'ppt'
    notes_masters_dir = ppt_dir / 'notesMasters'
    notes_masters_dir.mkdir(exist_ok=True)

    notes_master_path = notes_masters_dir / 'notesMaster1.xml'
    if not notes_master_path.exists():
        notes_master_path.write_text(create_notes_master_xml(), encoding='utf-8')

    theme_dir = ppt_dir / 'theme'
    theme_dir.mkdir(exist_ok=True)
    theme1_path = theme_dir / 'theme1.xml'
    theme2_path = theme_dir / 'theme2.xml'
    if not theme2_path.exists():
        if theme1_path.exists():
            shutil.copy2(theme1_path, theme2_path)
        else:
            raise RuntimeError('Cannot create notes theme: ppt/theme/theme1.xml is missing')

    notes_master_rels_dir = notes_masters_dir / '_rels'
    notes_master_rels_dir.mkdir(exist_ok=True)
    notes_master_rels_path = notes_master_rels_dir / 'notesMaster1.xml.rels'
    if not notes_master_rels_path.exists():
        notes_master_rels_path.write_text(
            create_notes_master_rels_xml(),
            encoding='utf-8',
        )

    presentation_rels_path = ppt_dir / '_rels' / 'presentation.xml.rels'
    notes_master_rid = _find_relationship_id(
        presentation_rels_path,
        _NOTES_MASTER_REL_TYPE,
        'notesMasters/notesMaster1.xml',
    )
    if notes_master_rid is None:
        notes_master_rid = _append_relationship(
            presentation_rels_path,
            _NOTES_MASTER_REL_TYPE,
            'notesMasters/notesMaster1.xml',
        )

    presentation_path = ppt_dir / 'presentation.xml'
    presentation_xml = presentation_path.read_text(encoding='utf-8')
    if '<p:notesMasterIdLst>' in presentation_xml:
        return
    notes_master_lst = (
        f'<p:notesMasterIdLst><p:notesMasterId r:id="{notes_master_rid}"/>'
        '</p:notesMasterIdLst>'
    )
    if '</p:sldMasterIdLst>' not in presentation_xml:
        raise RuntimeError('presentation.xml is missing p:sldMasterIdLst')
    presentation_xml = presentation_xml.replace(
        '</p:sldMasterIdLst>',
        '</p:sldMasterIdLst>' + notes_master_lst,
        1,
    )
    presentation_path.write_text(presentation_xml, encoding='utf-8')


def _to_float(value: Any, default: float) -> float:
    if value is None:
        return default
    try:
        number = float(value)
    except (TypeError, ValueError):
        return default
    return number if number >= 0 else default


def _slide_config(animation_config: dict[str, Any] | None, svg_stem: str) -> dict[str, Any]:
    if not animation_config:
        return {}
    slides = _as_dict(animation_config.get('slides'))
    return _as_dict(slides.get(svg_stem))


def _slide_transition_settings(
    slide_cfg: dict[str, Any],
    transition: str | None,
    duration: float,
    auto_advance: float | None,
    cli_overrides: dict[str, bool],
) -> tuple[str | None, float, float | None]:
    trans_cfg = _as_dict(slide_cfg.get('transition'))
    effect = transition
    if not cli_overrides.get('transition') and 'effect' in trans_cfg:
        cfg_effect = str(trans_cfg.get('effect'))
        effect = None if cfg_effect == 'none' else cfg_effect
    if not cli_overrides.get('transition_duration'):
        duration = _to_float(trans_cfg.get('duration'), duration)
    if not cli_overrides.get('auto_advance') and 'auto_advance' in trans_cfg:
        auto_advance = _to_float(trans_cfg.get('auto_advance'), auto_advance or 0)
    return effect, duration, auto_advance


def _slide_animation_settings(
    slide_cfg: dict[str, Any],
    animation: str | None,
    duration: float,
    stagger: float,
    trigger: str,
    cli_overrides: dict[str, bool],
) -> tuple[str | None, float, float, str]:
    anim_cfg = _as_dict(slide_cfg.get('animation'))
    effect = animation
    if not cli_overrides.get('animation') and 'effect' in anim_cfg:
        cfg_effect = str(anim_cfg.get('effect'))
        effect = None if cfg_effect == 'none' else cfg_effect
    if not cli_overrides.get('animation_duration'):
        duration = _to_float(anim_cfg.get('duration'), duration)
    if not cli_overrides.get('animation_stagger'):
        stagger = _to_float(anim_cfg.get('stagger'), stagger)
    if not cli_overrides.get('animation_trigger') and anim_cfg.get('trigger'):
        trigger = str(anim_cfg.get('trigger'))
    return effect, duration, stagger, trigger


def _build_sequence_targets(
    anim_targets: list[tuple[int, str]],
    slide_cfg: dict[str, Any],
    animation: str,
    duration: float,
    stagger: float,
    mixed_animation_offset: int,
) -> tuple[list[tuple[int, int, str, float]], int]:
    groups_cfg = _as_dict(slide_cfg.get('groups'))
    ordered: list[tuple[int, int, int, str, dict[str, Any]]] = []
    for idx, (sid, svg_id) in enumerate(anim_targets):
        group_cfg = _as_dict(groups_cfg.get(svg_id))
        if str(group_cfg.get('effect', '')).lower() == 'none':
            continue
        order_value = group_cfg.get('order')
        try:
            order = int(order_value)
            has_order = 0
        except (TypeError, ValueError):
            order = idx
            has_order = 1
        group_entry = dict(group_cfg)
        group_entry['_shape_id'] = sid
        ordered.append((has_order, order, idx, svg_id, group_entry))

    ordered.sort(key=lambda item: (item[0], item[1], item[2]))

    seq_targets: list[tuple[int, int, str, float]] = []
    for seq_idx, (_has_order, _order, _original_idx, _svg_id, group_cfg) in enumerate(ordered):
        shape_id = int(group_cfg['_shape_id'])
        raw_effect = group_cfg.get('effect')
        if raw_effect in ('auto', 'mixed', 'random'):
            effect = pick_animation_effect(
                str(raw_effect), seq_idx, mixed_animation_offset, group_id=_svg_id,
            )
        else:
            effect = str(raw_effect or pick_animation_effect(
                animation, seq_idx, mixed_animation_offset, group_id=_svg_id,
            ))
        item_duration = _to_float(group_cfg.get('duration'), duration)
        delay_seconds = _to_float(
            group_cfg.get('delay'),
            0 if seq_idx == 0 else stagger,
        )
        seq_targets.append((shape_id, int(delay_seconds * 1000), effect, item_duration))

    mixed_count = 0
    if animation == 'mixed':
        mixed_count = sum(1 for _target in seq_targets[1:])
    elif animation == 'auto':
        # 'auto' accumulates a cross-slide offset so the image pool and the
        # unmatched-id fallback rotate as the deck advances. Single-effect
        # semantic matches (title→fade, chart→wipe etc.) are unaffected
        # because they ignore the offset.
        mixed_count = len(seq_targets)
    return seq_targets, mixed_count


def _prerender_legacy_pngs(
    svg_files: list[Path],
    media_dir: Path,
    pixel_width: int,
    pixel_height: int,
    cache_dir: Path | None,
    workers: int,
    verbose: bool,
) -> dict[int, bool]:
    """Render every SVG→PNG into media_dir in parallel.

    Returns {1-based slide index: success}. Falls back to sequential when
    workers<=1 or len(svg_files)<=2.
    """
    results: dict[int, bool] = {}
    targets: list[tuple[int, Path, Path]] = [
        (i, svg, media_dir / f'image{i}.png')
        for i, svg in enumerate(svg_files, 1)
    ]

    if workers <= 1 or len(targets) <= 2:
        for i, svg, png in targets:
            ok = convert_svg_to_png_cached(svg, png, pixel_width, pixel_height, cache_dir)
            results[i] = ok
            if verbose:
                tag = 'cached/ok' if ok else 'failed'
                print(f"  [PNG {i}/{len(targets)}] {svg.name} - {tag}")
        return results

    with ProcessPoolExecutor(max_workers=workers) as pool:
        future_map = {
            pool.submit(
                convert_svg_to_png_cached,
                svg, png, pixel_width, pixel_height, cache_dir,
            ): (i, svg)
            for i, svg, png in targets
        }
        done = 0
        for future in as_completed(future_map):
            i, svg = future_map[future]
            try:
                ok = future.result()
            except Exception as exc:
                ok = False
                if verbose:
                    print(f"  [PNG] {svg.name} - worker error: {exc}")
            results[i] = ok
            done += 1
            if verbose:
                tag = 'cached/ok' if ok else 'failed'
                print(f"  [PNG {done}/{len(targets)}] {svg.name} - {tag}")

    return results


_REL_TARGET_RE = re.compile(r'<Relationship\b[^/]*?/>', re.DOTALL)
_TARGET_ATTR_RE = re.compile(r'Target="([^"]+)"')
_TARGET_MODE_EXT_RE = re.compile(r'TargetMode="External"')


def _verify_internal_rels_targets(extract_dir: Path) -> list[str]:
    """Return a list of dangling internal Targets across every .rels in the package.

    Each entry is formatted as "<rels-path> -> <missing-target>". An empty list
    means every internal Target resolves to a real file in the package.
    """
    problems: list[str] = []
    for rels_path in extract_dir.rglob('*.rels'):
        rels_rel = rels_path.relative_to(extract_dir).as_posix()
        # `_rels/foo.xml.rels` lives one level below its referent's directory;
        # Targets resolve relative to the parent of that `_rels` folder.
        base_dir = posixpath.dirname(posixpath.dirname(rels_rel))
        content = rels_path.read_text(encoding='utf-8')
        for match in _REL_TARGET_RE.finditer(content):
            element = match.group(0)
            if _TARGET_MODE_EXT_RE.search(element):
                continue
            target_match = _TARGET_ATTR_RE.search(element)
            if not target_match:
                continue
            target = target_match.group(1)
            if target.startswith(('http://', 'https://', 'mailto:')):
                continue
            resolved = posixpath.normpath(posixpath.join(base_dir, target)) if base_dir else posixpath.normpath(target)
            if not (extract_dir / resolved).exists():
                problems.append(f'{rels_rel} -> {resolved}')
    return problems


def _presentation_format(width: float, height: float) -> str:
    """Map the slide aspect ratio to PowerPoint's PresentationFormat label.
    Non-standard ratios (square, portrait, banner crops) report 'Custom'.
    """
    if width <= 0 or height <= 0:
        return 'Custom'
    ratio = width / height
    for target, label in (
        (4 / 3, 'On-screen Show (4:3)'),
        (16 / 9, 'On-screen Show (16:9)'),
        (16 / 10, 'On-screen Show (16:10)'),
    ):
        if abs(ratio - target) < 0.02:
            return label
    return 'Custom'


def _stamp_docprops(
    extract_dir: Path,
    slide_count: int,
    pres_format: str,
    meta: dict[str, Any] | None = None,
) -> None:
    """Overwrite the misleading python-pptx default metadata with accurate
    values. Factual fields (slide count, export timestamp, presentation format,
    application) are always machine-derived. Authored fields — including the
    title — come solely from an optional per-project ``metadata.json``
    (``meta``); whatever it omits stays blank. ``lastModifiedBy`` follows
    ``creator`` rather than ever carrying the base template's author or a tool
    name. No field is guessed from slide content: a blank title is preferable
    to an unreliable heuristic pick.
    """
    meta = meta or {}

    def field(key: str, default: str = '') -> str:
        value = meta.get(key)
        return value.strip() if isinstance(value, str) and value.strip() else default

    title = field('title')
    creator = field('creator')

    now = datetime.now(timezone.utc).strftime('%Y-%m-%dT%H:%M:%SZ')

    core_path = extract_dir / 'docProps' / 'core.xml'
    if core_path.exists():
        core_path.write_text(
            "<?xml version='1.0' encoding='UTF-8' standalone='yes'?>\n"
            '<cp:coreProperties '
            'xmlns:cp="http://schemas.openxmlformats.org/package/2006/metadata/core-properties" '
            'xmlns:dc="http://purl.org/dc/elements/1.1/" '
            'xmlns:dcterms="http://purl.org/dc/terms/" '
            'xmlns:dcmitype="http://purl.org/dc/dcmitype/" '
            'xmlns:xsi="http://www.w3.org/2001/XMLSchema-instance">'
            f'<dc:title>{escape(title)}</dc:title>'
            f'<dc:subject>{escape(field("subject"))}</dc:subject>'
            f'<dc:creator>{escape(creator)}</dc:creator>'
            f'<cp:keywords>{escape(field("keywords"))}</cp:keywords>'
            f'<dc:description>{escape(field("description"))}</dc:description>'
            f'<dc:language>{escape(field("language"))}</dc:language>'
            f'<cp:lastModifiedBy>{escape(creator)}</cp:lastModifiedBy>'
            '<cp:revision>1</cp:revision>'
            f'<dcterms:created xsi:type="dcterms:W3CDTF">{now}</dcterms:created>'
            f'<dcterms:modified xsi:type="dcterms:W3CDTF">{now}</dcterms:modified>'
            f'<cp:category>{escape(field("category"))}</cp:category>'
            f'<cp:contentStatus>{escape(field("contentStatus"))}</cp:contentStatus>'
            '</cp:coreProperties>',
            encoding='utf-8',
        )

    app_path = extract_dir / 'docProps' / 'app.xml'
    if app_path.exists():
        app = app_path.read_text(encoding='utf-8')
        app = re.sub(r'<Slides>.*?</Slides>', f'<Slides>{slide_count}</Slides>', app)
        app = re.sub(
            r'<Company>.*?</Company>',
            f'<Company>{escape(field("company"))}</Company>',
            app,
        )
        app = re.sub(
            r'<Manager>.*?</Manager>',
            f'<Manager>{escape(field("manager"))}</Manager>',
            app,
        )
        app = re.sub(
            r'<Application>.*?</Application>',
            '<Application>Microsoft Office PowerPoint</Application>',
            app,
        )
        app = re.sub(
            r'<PresentationFormat>.*?</PresentationFormat>',
            f'<PresentationFormat>{escape(pres_format)}</PresentationFormat>',
            app,
        )
        app_path.write_text(app, encoding='utf-8')


def create_pptx_with_native_svg(
    svg_files: list[Path],
    output_path: Path,
    canvas_format: str | None = None,
    verbose: bool = True,
    transition: str | None = 'fade',
    transition_duration: float = 0.5,
    auto_advance: float | None = None,
    use_compat_mode: bool = True,
    notes: dict[str, str] | None = None,
    enable_notes: bool = True,
    use_native_shapes: bool = False,
    animation: str | None = None,
    animation_duration: float = 0.4,
    animation_stagger: float = 0.5,
    animation_trigger: str = 'after-previous',
    animation_config: dict[str, Any] | None = None,
    animation_cli_overrides: dict[str, bool] | None = None,
    narration_audio: dict[str, Path] | None = None,
    use_narration_timings: bool = False,
    narration_padding: float = 0.5,
    cache_dir: Path | None = None,
    workers: int | None = None,
    merge_paragraphs: bool = True,
    conversion_trace_path: Path | None = None,
    doc_metadata: dict[str, Any] | None = None,
) -> bool:
    """Create a PPTX file with native SVG.

    Args:
        svg_files: List of SVG files.
        output_path: Output PPTX path.
        canvas_format: Canvas format key.
        verbose: Whether to output detailed information.
        transition: Transition effect name.
        transition_duration: Transition duration in seconds.
        auto_advance: Auto-advance interval in seconds.
        use_compat_mode: Use Office compatibility mode (PNG + SVG dual format).
        notes: Notes dict, key is SVG stem, value is notes content.
        enable_notes: Whether to enable notes embedding.
        use_native_shapes: Convert SVG to native DrawingML shapes.
        animation: Per-element entrance animation mode (single effect name,
            'mixed', 'random', or None to disable). Native shapes mode only.
        animation_duration: Per-element entrance duration in seconds.
        animation_stagger: Delay between elements in ``after-previous``
            trigger mode (seconds). Ignored otherwise.
        animation_trigger: PowerPoint Start mode — ``'after-previous'`` (default),
            ``'on-click'``, or ``'with-previous'``.
        animation_config: Optional sidecar overrides loaded from animations.json.
        animation_cli_overrides: Flags indicating explicit CLI overrides.
        narration_audio: Optional dict mapping SVG stem to narration audio file.
        use_narration_timings: Whether to set slide auto-advance from audio duration.
        narration_padding: Extra seconds added after each narration before advancing.
        conversion_trace_path: Optional JSON path for native conversion diagnostics.

    Returns:
        Whether all slides were successfully created.
    """
    if not svg_files:
        print("Error: No SVG files found")
        return False

    # Native shapes mode takes priority over compat mode
    if use_native_shapes:
        use_compat_mode = False

    # Check compatibility mode dependencies
    renderer_name, renderer_status, renderer_hint = get_png_renderer_info()
    if not use_native_shapes and use_compat_mode and PNG_RENDERER is None:
        print("Warning: No PNG rendering library installed, cannot use compatibility mode")
        print(f"  {renderer_hint}")
        print("  Will use pure SVG mode (may not display in Office LTSC 2021 and similar versions)")
        use_compat_mode = False

    # Auto-detect canvas format or get dimensions from viewBox
    custom_pixels: tuple[int, int] | None = None
    if canvas_format is None:
        canvas_format = detect_format_from_svg(svg_files[0])
        if canvas_format and verbose:
            format_name = CANVAS_FORMATS.get(canvas_format, {}).get('name', canvas_format)
            print(f"  Detected canvas format: {format_name}")

    if canvas_format is None:
        custom_pixels = get_viewbox_dimensions(svg_files[0])
        if custom_pixels and verbose:
            print(f"  Using SVG viewBox dimensions: {custom_pixels[0]} x {custom_pixels[1]} px")

    if canvas_format is None and custom_pixels is None:
        canvas_format = 'ppt169'
        if verbose:
            print(f"  Using default format: PPT 16:9")

    width_emu, height_emu = get_slide_dimensions(canvas_format or 'ppt169', custom_pixels)
    pixel_width, pixel_height = get_pixel_dimensions(canvas_format or 'ppt169', custom_pixels)

    if verbose:
        print(f"  Slide dimensions: {pixel_width} x {pixel_height} px")
        print(f"  SVG file count: {len(svg_files)}")
        if use_native_shapes:
            print(f"  Mode: Native DrawingML shapes (directly editable)")
        elif use_compat_mode:
            print(f"  Compatibility mode: Enabled (PNG + SVG dual format)")
            print(f"  PNG renderer: {renderer_name} {renderer_status}")
        else:
            print(f"  Compatibility mode: Disabled (pure SVG)")
        if transition:
            trans_name = TRANSITIONS.get(transition, {}).get('name', transition) if TRANSITIONS else transition
            print(f"  Transition effect: {trans_name}")
        if enable_notes and notes:
            print(f"  Speaker notes: {len(notes)} page(s)")
        elif enable_notes:
            print(f"  Speaker notes: Enabled (no notes files found)")
        else:
            print(f"  Speaker notes: Disabled")
        print()

    animation_cli_overrides = animation_cli_overrides or {}

    temp_dir = _create_writable_work_dir(output_path)

    try:
        # Create base PPTX with python-pptx
        prs = Presentation()
        prs.slide_width = width_emu
        prs.slide_height = height_emu

        blank_layout = prs.slide_layouts[6]
        for _ in svg_files:
            prs.slides.add_slide(blank_layout)

        base_pptx = temp_dir / 'base.pptx'
        prs.save(str(base_pptx))

        # Extract PPTX
        extract_dir = temp_dir / 'pptx_content'
        with zipfile.ZipFile(base_pptx, 'r') as zf:
            zf.extractall(extract_dir)

        media_dir = extract_dir / 'ppt' / 'media'
        media_dir.mkdir(exist_ok=True)

        prerender_results: dict[int, bool] | None = None
        if not use_native_shapes and use_compat_mode and PNG_RENDERER is not None:
            if workers is None:
                resolved_workers = min(os.cpu_count() or 2, len(svg_files), 8)
            else:
                resolved_workers = max(0, workers)
            if verbose:
                cache_label = str(cache_dir) if cache_dir else 'disabled'
                mode = f'parallel x{resolved_workers}' if resolved_workers > 1 else 'sequential'
                print(f"  Pre-rendering PNGs ({mode}, cache: {cache_label})")
            prerender_results = _prerender_legacy_pngs(
                svg_files, media_dir, pixel_width, pixel_height,
                cache_dir, resolved_workers, verbose,
            )
            if verbose:
                print()

        success_count = 0
        has_any_image = False
        media_cache: dict[tuple[str, str], str] = {}
        image_exts_used: set[str] = set()
        notes_slides_created: set[int] = set()
        narration_slides_created: set[int] = set()
        audio_exts_used: set[str] = set()
        mixed_animation_offset = 0
        conversion_trace: list[dict[str, Any]] | None = [] if conversion_trace_path else None

        for i, svg_path in enumerate(svg_files, 1):
            slide_num = i

            try:
                # ---- Native shapes mode ----
                if use_native_shapes:
                    slide_cfg = _slide_config(animation_config, svg_path.stem)
                    slide_xml, media_files_dict, rel_entries, anim_targets = (
                        convert_svg_to_slide_shapes(
                            svg_path, slide_num=slide_num, verbose=verbose,
                            merge_paragraphs=merge_paragraphs,
                            trace_out=conversion_trace,
                        )
                    )
                    slide_transition, slide_transition_duration, slide_auto_advance = (
                        _slide_transition_settings(
                            slide_cfg,
                            transition,
                            transition_duration,
                            auto_advance,
                            animation_cli_overrides,
                        )
                    )
                    (
                        slide_animation,
                        slide_animation_duration,
                        slide_animation_stagger,
                        slide_animation_trigger,
                    ) = _slide_animation_settings(
                        slide_cfg,
                        animation,
                        animation_duration,
                        animation_stagger,
                        animation_trigger,
                        animation_cli_overrides,
                    )

                    # Order matters: OOXML schema requires <p:transition>
                    # to precede <p:timing> inside <p:sld>. Both use the same
                    # </p:sld> string-replace anchor, so transition must be
                    # injected first and timing second.
                    if slide_transition and ANIMATIONS_AVAILABLE and create_transition_xml:
                        transition_xml = '\n' + create_transition_xml(
                            effect=slide_transition,
                            duration=slide_transition_duration,
                            advance_after=slide_auto_advance,
                        )
                        slide_xml = slide_xml.replace(
                            '</p:sld>',
                            transition_xml + '\n</p:sld>',
                        )

                    if (slide_animation and slide_animation != 'none'
                            and create_sequence_timing_xml
                            and pick_animation_effect
                            and anim_targets):
                        seq_targets, mixed_count = _build_sequence_targets(
                            anim_targets,
                            slide_cfg,
                            slide_animation,
                            slide_animation_duration,
                            slide_animation_stagger,
                            mixed_animation_offset,
                        )
                        if slide_animation in ('mixed', 'auto'):
                            mixed_animation_offset += mixed_count
                        timing_xml = '\n' + create_sequence_timing_xml(
                            seq_targets, duration=slide_animation_duration,
                            trigger=slide_animation_trigger,
                        )
                        slide_xml = slide_xml.replace(
                            '</p:sld>',
                            timing_xml + '\n</p:sld>',
                        )

                    # Write slide XML
                    slide_xml_path = extract_dir / 'ppt' / 'slides' / f'slide{slide_num}.xml'
                    with open(slide_xml_path, 'w', encoding='utf-8') as f:
                        f.write(slide_xml)

                    # Write media files
                    media_name_map: dict[str, str] = {}
                    for media_name, media_data in media_files_dict.items():
                        ext = media_name.rsplit('.', 1)[-1].lower()
                        media_hash = hashlib.sha256(media_data).hexdigest()
                        cache_key = (ext, media_hash)
                        cached_name = media_cache.get(cache_key)

                        if cached_name is None:
                            cached_name = f'image_{media_hash[:16]}.{ext}'
                            media_cache[cache_key] = cached_name
                            with open(media_dir / cached_name, 'wb') as f:
                                f.write(media_data)

                        media_name_map[media_name] = cached_name

                    for rel in rel_entries:
                        target = rel.get('target', '')
                        if not target.startswith('../media/'):
                            continue
                        media_name = target.split('../media/', 1)[1]
                        mapped_name = media_name_map.get(media_name)
                        if mapped_name:
                            rel['target'] = f'../media/{mapped_name}'

                    # Build relationships XML
                    rels_dir = extract_dir / 'ppt' / 'slides' / '_rels'
                    rels_dir.mkdir(exist_ok=True)
                    rels_path = rels_dir / f'slide{slide_num}.xml.rels'

                    extra_rels = ''
                    for rel in rel_entries:
                        extra_rels += (
                            f'\n  <Relationship Id="{rel["id"]}" '
                            f'Type="{rel["type"]}" Target="{rel["target"]}"/>'
                        )

                    rels_xml = f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout" Target="../slideLayouts/slideLayout1.xml"/>{extra_rels}
</Relationships>'''
                    with open(rels_path, 'w', encoding='utf-8') as f:
                        f.write(rels_xml)

                    # Track image formats for Content_Types
                    for media_name in media_name_map.values():
                        ext = media_name.rsplit('.', 1)[-1].lower()
                        _content_type_for_extension(ext)
                        image_exts_used.add(ext)
                        has_any_image = True

                # ---- Legacy SVG embedding mode ----
                else:
                    slide_cfg = _slide_config(animation_config, svg_path.stem)
                    slide_transition, slide_transition_duration, slide_auto_advance = (
                        _slide_transition_settings(
                            slide_cfg,
                            transition,
                            transition_duration,
                            auto_advance,
                            animation_cli_overrides,
                        )
                    )
                    svg_filename = f'image{i}.svg'
                    png_filename = f'image{i}.png'
                    png_rid = 'rId2'
                    svg_rid = 'rId3' if use_compat_mode else 'rId2'

                    shutil.copy(svg_path, media_dir / svg_filename)

                    slide_has_png = False
                    if use_compat_mode:
                        if prerender_results is not None:
                            png_success = prerender_results.get(i, False)
                        else:
                            png_path = media_dir / png_filename
                            png_success = convert_svg_to_png(
                                svg_path, png_path,
                                width=pixel_width, height=pixel_height,
                            )
                        if png_success:
                            slide_has_png = True
                            has_any_image = True
                            image_exts_used.add('png')
                        else:
                            if verbose:
                                print(f"  [{i}/{len(svg_files)}] {svg_path.name} - PNG generation failed, using pure SVG")
                            svg_rid = 'rId2'

                    slide_xml_path = extract_dir / 'ppt' / 'slides' / f'slide{slide_num}.xml'
                    slide_xml = create_slide_xml_with_svg(
                        slide_num,
                        png_rid=png_rid, svg_rid=svg_rid,
                        width_emu=width_emu, height_emu=height_emu,
                        transition=slide_transition,
                        transition_duration=slide_transition_duration,
                        auto_advance=slide_auto_advance,
                        use_compat_mode=(use_compat_mode and slide_has_png),
                    )
                    with open(slide_xml_path, 'w', encoding='utf-8') as f:
                        f.write(slide_xml)

                    rels_dir = extract_dir / 'ppt' / 'slides' / '_rels'
                    rels_dir.mkdir(exist_ok=True)
                    rels_path = rels_dir / f'slide{slide_num}.xml.rels'
                    rels_xml = create_slide_rels_xml(
                        png_rid=png_rid, png_filename=png_filename,
                        svg_rid=svg_rid, svg_filename=svg_filename,
                        use_compat_mode=(use_compat_mode and slide_has_png),
                    )
                    with open(rels_path, 'w', encoding='utf-8') as f:
                        f.write(rels_xml)

                # --- Process notes (shared between native and legacy mode) ---
                notes_content = ''
                if enable_notes:
                    svg_stem = svg_path.stem
                    notes_content = notes.get(svg_stem, '') if notes else ''
                    notes_text = markdown_to_plain_text(notes_content) if notes_content else ''
                    if notes_text:
                        _ensure_notes_master(extract_dir)

                        notes_slides_dir = extract_dir / 'ppt' / 'notesSlides'
                        notes_slides_dir.mkdir(exist_ok=True)

                        notes_xml_path = notes_slides_dir / f'notesSlide{slide_num}.xml'
                        notes_xml = create_notes_slide_xml(slide_num, notes_text)
                        with open(notes_xml_path, 'w', encoding='utf-8') as f:
                            f.write(notes_xml)

                        notes_rels_dir = notes_slides_dir / '_rels'
                        notes_rels_dir.mkdir(exist_ok=True)
                        notes_rels_path = notes_rels_dir / f'notesSlide{slide_num}.xml.rels'
                        notes_rels_xml = create_notes_slide_rels_xml(slide_num)
                        with open(notes_rels_path, 'w', encoding='utf-8') as f:
                            f.write(notes_rels_xml)

                        _append_relationship(
                            rels_path,
                            'http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesSlide',
                            f'../notesSlides/notesSlide{slide_num}.xml',
                        )
                        notes_slides_created.add(slide_num)

                # --- Process narration audio (shared between native and legacy mode) ---
                svg_stem = svg_path.stem
                audio_path = narration_audio.get(svg_stem) if narration_audio else None
                if audio_path:
                    slide_xml_path = extract_dir / 'ppt' / 'slides' / f'slide{slide_num}.xml'
                    rels_path = extract_dir / 'ppt' / 'slides' / '_rels' / f'slide{slide_num}.xml.rels'

                    ext = audio_path.suffix.lower()
                    media_name = f'narration{slide_num}{ext}'
                    shutil.copy2(audio_path, media_dir / media_name)
                    audio_exts_used.add(ext)

                    poster_name = 'narration_poster.png'
                    poster_path = media_dir / poster_name
                    if not poster_path.exists():
                        poster_path.write_bytes(TRANSPARENT_PNG_BYTES)
                    has_any_image = True
                    image_exts_used.add('png')

                    media_rid = _append_relationship(
                        rels_path,
                        MEDIA_REL_TYPE,
                        f'../media/{media_name}',
                    )
                    audio_rid = _append_relationship(
                        rels_path,
                        AUDIO_REL_TYPE,
                        f'../media/{media_name}',
                    )
                    poster_rid = _append_relationship(
                        rels_path,
                        IMAGE_REL_TYPE,
                        f'../media/{poster_name}',
                    )

                    slide_xml = slide_xml_path.read_text(encoding='utf-8')
                    narration_shape_id = next_shape_id(slide_xml)
                    slide_xml = inject_narration(
                        slide_xml,
                        shape_id=narration_shape_id,
                        shape_name=media_name,
                        audio_rid=audio_rid,
                        media_rid=media_rid,
                        poster_rid=poster_rid,
                    )

                    if use_narration_timings:
                        duration = probe_audio_duration(audio_path)
                        if duration is None:
                            raise RuntimeError(
                                f"Unable to read narration duration with ffprobe: {audio_path}"
                            )
                        slide_xml = apply_recorded_timing(
                            slide_xml,
                            advance_after=duration + narration_padding,
                            transition_duration=slide_transition_duration,
                            transition_effect=slide_transition or 'fade',
                        )
                    slide_xml_path.write_text(slide_xml, encoding='utf-8')
                    narration_slides_created.add(slide_num)

                if verbose:
                    if use_native_shapes:
                        mode_str = " (Native)"
                    elif use_compat_mode and not use_native_shapes:
                        mode_str = " (PNG+SVG)" if has_any_image else " (SVG)"
                    else:
                        mode_str = " (SVG)"
                    has_notes = slide_num in notes_slides_created
                    notes_str = " +notes" if has_notes else ""
                    narration_str = " +narration" if slide_num in narration_slides_created else ""
                    print(f"  [{i}/{len(svg_files)}] {svg_path.name}{mode_str}{notes_str}{narration_str}")

                success_count += 1

            except Exception as e:
                if verbose:
                    print(f"  [{i}/{len(svg_files)}] {svg_path.name} - Error: {e}")
                if use_native_shapes:
                    raise

        # Update [Content_Types].xml
        content_types_path = extract_dir / '[Content_Types].xml'
        with open(content_types_path, 'r', encoding='utf-8') as f:
            content_types = f.read()

        if not use_native_shapes:
            content_types = _add_default_content_type(content_types, 'svg', 'image/svg+xml')
        for ext in sorted(image_exts_used):
            content_types = _add_default_content_type(
                content_types,
                ext,
                _content_type_for_extension(ext),
            )
        with open(content_types_path, 'w', encoding='utf-8') as f:
            f.write(content_types)

        if audio_exts_used:
            for ext in sorted(audio_exts_used):
                content_type = AUDIO_CONTENT_TYPES.get(ext)
                if content_type:
                    content_types = _add_default_content_type(content_types, ext, content_type)
            if 'Extension="png"' not in content_types:
                content_types = _add_default_content_type(content_types, 'png', 'image/png')
            with open(content_types_path, 'w', encoding='utf-8') as f:
                f.write(content_types)

        # Add notes master / slides content types
        if enable_notes and notes_slides_created:
            notes_theme_override = (
                '  <Override PartName="/ppt/theme/theme2.xml" '
                'ContentType="application/vnd.openxmlformats-officedocument.theme+xml"/>'
            )
            if notes_theme_override not in content_types:
                content_types = content_types.replace(
                    '</Types>',
                    notes_theme_override + '\n</Types>',
                )
            notes_master_override = (
                '  <Override PartName="/ppt/notesMasters/notesMaster1.xml" '
                'ContentType="application/vnd.openxmlformats-officedocument.presentationml.notesMaster+xml"/>'
            )
            if notes_master_override not in content_types:
                content_types = content_types.replace(
                    '</Types>',
                    notes_master_override + '\n</Types>',
                )
            for i in sorted(notes_slides_created):
                override = (
                    f'  <Override PartName="/ppt/notesSlides/notesSlide{i}.xml" '
                    f'ContentType="application/vnd.openxmlformats-officedocument.presentationml.notesSlide+xml"/>'
                )
                if override not in content_types:
                    content_types = content_types.replace('</Types>', override + '\n</Types>')
            with open(content_types_path, 'w', encoding='utf-8') as f:
                f.write(content_types)

        rels_problems = _verify_internal_rels_targets(extract_dir)
        if rels_problems:
            details = '\n'.join(f'  - {p}' for p in rels_problems)
            raise RuntimeError(
                'PPTX package contains dangling internal relationship targets; '
                'PowerPoint will report the file as corrupt:\n' + details
            )

        # Replace the python-pptx base-template metadata (stale "Steve Canny"
        # author, 2013 dates, "generated using python-pptx", Slides=0) with
        # accurate, tool-neutral document properties.
        pres_format = _presentation_format(width_emu, height_emu)
        _stamp_docprops(extract_dir, len(svg_files), pres_format, doc_metadata)

        # Repackage PPTX to a temporary file first. The public output path is
        # replaced only after every slide and relationship has succeeded.
        temp_output_path = temp_dir / 'result.pptx'
        with zipfile.ZipFile(temp_output_path, 'w', zipfile.ZIP_DEFLATED) as zf:
            for file_path in extract_dir.rglob('*'):
                if file_path.is_file():
                    arcname = file_path.relative_to(extract_dir)
                    zf.write(file_path, arcname)
        shutil.move(str(temp_output_path), str(output_path))
        permission_warnings = _relax_output_permissions(output_path)

        if conversion_trace_path and conversion_trace is not None:
            conversion_trace_path.parent.mkdir(parents=True, exist_ok=True)
            payload = {
                'output': str(output_path),
                'slide_count': len(svg_files),
                'slides': conversion_trace,
            }
            conversion_trace_path.write_text(
                json.dumps(payload, ensure_ascii=False, indent=2),
                encoding='utf-8',
            )

        if verbose:
            print()
            print(f"[Done] Saved: {output_path}")
            for warning in permission_warnings:
                print(f"  [warn] {warning}")
            if conversion_trace_path and conversion_trace is not None:
                print(f"  Trace: {conversion_trace_path}")
            print(f"  Succeeded: {success_count}, Failed: {len(svg_files) - success_count}")
            if use_compat_mode and has_any_image:
                print(f"  Mode: Office compatibility mode (supports all Office versions)")
                if PNG_RENDERER == 'svglib' and renderer_hint:
                    print(f"  [Tip] {renderer_hint}")

        return success_count == len(svg_files)

    finally:
        shutil.rmtree(temp_dir, ignore_errors=True)
