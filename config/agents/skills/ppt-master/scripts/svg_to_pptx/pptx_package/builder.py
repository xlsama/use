"""Core PPTX assembly: create_pptx_with_native_svg."""

from __future__ import annotations

import hashlib
import json
import math
import mimetypes
import os
import posixpath
import random
import re
import shutil
import stat
import subprocess
import tempfile
import uuid
import zipfile
from concurrent.futures import ProcessPoolExecutor, as_completed
from dataclasses import asdict, dataclass
from datetime import datetime, timezone
from pathlib import Path
from typing import Any
from urllib.parse import urlsplit
from xml.etree import ElementTree as ET
from xml.sax.saxutils import escape

from pptx import Presentation
from pptx.util import Emu

from pptx_transitions import (
    TRANSITIONS,
    create_transition_xml,
    normalize_transition_effect,
    set_directory_use_timings,
    validate_generated_transition_xml,
    validate_pptx_transition_package,
    validate_seconds,
)
from pptx_animations import (
    animation_seconds_to_milliseconds,
    create_sequence_timing_xml,
    normalize_animation_effect,
    normalize_animation_trigger,
    pick_animation_effect,
    validate_generated_animation_xml,
    validate_pptx_animation_package,
)

from ..drawingml.converter import convert_svg_to_slide_shapes
from ..drawingml.theme_colors import (
    ThemeColorSpec,
    apply_theme_color_spec,
    rewrite_chart_accent_colors,
)
from ..drawingml.theme_fonts import (
    MasterTextStyleSpec,
    ThemeFontSpec,
    apply_master_text_style_spec,
    apply_theme_font_spec,
)
from ..drawingml.utils import EMU_PER_PX
from ..semantic_markers import (
    chrome_token_from_markers,
    page_layout_name_from_svg,
)
from .dimensions import (
    CANVAS_FORMATS,
    get_slide_dimensions, get_pixel_dimensions,
    get_viewbox_dimensions, detect_format_from_svg,
)
from .media import (
    PNG_RENDERER,
    get_png_renderer_info, convert_svg_to_png, convert_svg_to_png_cached,
)
from .notes import (
    markdown_to_plain_text,
    create_notes_master_rels_xml,
    create_notes_master_xml,
    create_notes_slide_xml,
    create_notes_slide_rels_xml,
)
from .narration import (
    AUDIO_CONTENT_TYPES,
    AUDIO_REL_TYPE,
    AUDIO_MARKER_PNG_BYTES,
    IMAGE_REL_TYPE,
    MEDIA_REL_TYPE,
    apply_recorded_timing,
    inject_narration,
    next_shape_id,
    probe_audio_duration,
)
from .slide_xml import (
    create_slide_xml_with_svg, create_slide_rels_xml,
)
from .template_structure import (
    NativeStructureContract,
    OOXML_UINT32_MAX,
    TEMPLATE_PLACEHOLDER_TYPES,
    TemplateElementSpec,
    TemplateSlideSpec,
    TemplateStructureError,
    is_proxy_placeholder,
    match_native_placeholders,
    parse_preserve_slides,
    parse_template_slides,
    template_placeholder_bindings,
)
from .template_validation import validate_pptx_template_package

SLIDE_LAYOUT_REL_TYPE = (
    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout"
)
SLIDE_MASTER_REL_TYPE = (
    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideMaster"
)
PML_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
DML_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
REL_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
P14_NS = "http://schemas.microsoft.com/office/powerpoint/2010/main"

for _prefix, _uri in (("p", PML_NS), ("a", DML_NS), ("r", REL_NS), ("p14", P14_NS)):
    try:
        ET.register_namespace(_prefix, _uri)
    except (ValueError, AttributeError):
        pass


@dataclass(frozen=True)
class PptxStructureContext:
    """Resolved base package structure reused when slide XML is regenerated."""

    slide_layout_targets: dict[int, str]
    slide_master_parts: dict[int, str]

    def slide_layout_target(self, slide_num: int) -> str:
        """Return the slide layout target for a generated slide."""
        try:
            return self.slide_layout_targets[slide_num]
        except KeyError as exc:
            raise RuntimeError(
                f"Missing slide layout relationship for generated slide {slide_num}"
            ) from exc

    def slide_master_part(self, slide_num: int) -> str:
        """Return the slide master package part for a generated slide."""
        try:
            return self.slide_master_parts[slide_num]
        except KeyError as exc:
            raise RuntimeError(
                f"Missing slide master relationship for generated slide {slide_num}"
            ) from exc


@dataclass
class _TemplateRuntimeSlide:
    """Parsed slide package state used by explicit Layout structure export."""

    spec: TemplateSlideSpec
    slide_path: Path
    rels_path: Path
    tree: ET.ElementTree
    root: ET.Element
    rels: dict[str, dict[str, str]]
    shapes: dict[str, ET.Element]
    shape_ids_by_svg_id: dict[str, list[str]]


def _relationship_attrs(elem: ET.Element) -> dict[str, str]:
    return {key.rsplit("}", 1)[-1]: value for key, value in elem.attrib.items()}


def _resolve_package_target(source_part: str, target: str) -> str:
    """Resolve a relationship target relative to a package part path."""
    return posixpath.normpath(posixpath.join(posixpath.dirname(source_part), target))


def _relationships_path_for_part(extract_dir: Path, part_name: str) -> Path:
    """Return the package relationship sidecar path for a part name."""
    path = Path(part_name)
    return extract_dir / path.parent / "_rels" / f"{path.name}.rels"


def _find_relationship_target(
    rels_path: Path,
    rel_type: str,
) -> str | None:
    """Find the first relationship target for a relationship type."""
    if not rels_path.exists():
        return None
    root = ET.parse(rels_path).getroot()
    for elem in root:
        attrs = _relationship_attrs(elem)
        if attrs.get("Type") == rel_type:
            return attrs.get("Target")
    return None


def _read_relationships(rels_path: Path) -> dict[str, dict[str, str]]:
    """Return relationship attributes keyed by rId."""
    if not rels_path.exists():
        return {}
    root = ET.parse(rels_path).getroot()
    rels: dict[str, dict[str, str]] = {}
    for elem in root:
        attrs = _relationship_attrs(elem)
        rel_id = attrs.get("Id")
        if rel_id:
            rels[rel_id] = attrs
    return rels


def _find_relationship_id(
    rels_path: Path,
    rel_type: str,
    target: str,
) -> str | None:
    """Find an existing relationship by type and target."""
    for rel_id, attrs in _read_relationships(rels_path).items():
        if attrs.get("Type") == rel_type and attrs.get("Target") == target:
            return rel_id
    return None


def _read_slide_layout_targets(extract_dir: Path, slide_count: int) -> PptxStructureContext:
    """Read the actual layout relationship target for every generated slide."""
    slide_layout_targets: dict[int, str] = {}
    slide_master_parts: dict[int, str] = {}
    rels_dir = extract_dir / "ppt" / "slides" / "_rels"
    for slide_num in range(1, slide_count + 1):
        rels_path = rels_dir / f"slide{slide_num}.xml.rels"
        if not rels_path.exists():
            raise RuntimeError(f"Missing slide relationship file: {rels_path}")
        target = _find_relationship_target(rels_path, SLIDE_LAYOUT_REL_TYPE)
        if not target:
            raise RuntimeError(f"Slide {slide_num} has no slide layout relationship")
        slide_layout_targets[slide_num] = target

        slide_part = f"ppt/slides/slide{slide_num}.xml"
        layout_part = _resolve_package_target(slide_part, target)
        layout_rels_path = _relationships_path_for_part(extract_dir, layout_part)
        master_target = _find_relationship_target(layout_rels_path, SLIDE_MASTER_REL_TYPE)
        if not master_target:
            raise RuntimeError(
                f"Slide {slide_num} layout has no slide master relationship"
            )
        slide_master_parts[slide_num] = _resolve_package_target(layout_part, master_target)
    return PptxStructureContext(
        slide_layout_targets=slide_layout_targets,
        slide_master_parts=slide_master_parts,
    )


_SLIDE_BACKGROUND_RE = re.compile(
    r"(?P<prefix><p:cSld\b[^>]*>\s*)"
    r"(?P<bg><p:bg\b.*?</p:bg>)"
    r"(?P<suffix>\s*<p:spTree\b)",
    re.DOTALL,
)


def _extract_slide_background_xml(slide_xml: str) -> str | None:
    """Return the slide-level p:bg XML when it directly precedes spTree."""
    match = _SLIDE_BACKGROUND_RE.search(slide_xml)
    return match.group("bg") if match else None


def _remove_slide_background_xml(slide_xml: str) -> str:
    """Remove a promoted slide-level p:bg from cSld."""
    return _SLIDE_BACKGROUND_RE.sub(r"\g<prefix>\g<suffix>", slide_xml, count=1)


def _put_background_on_part(part_xml: str, background_xml: str) -> str | None:
    """Replace or insert p:bg before a slide/master/layout spTree.

    Returns None when the part carries a p:bg the canonical pattern cannot
    replace; inserting there would leave two p:bg children under p:cSld.
    """
    match = _SLIDE_BACKGROUND_RE.search(part_xml)
    if match:
        return (
            part_xml[:match.start("bg")]
            + background_xml
            + part_xml[match.end("bg"):]
        )
    if "<p:bg" in part_xml:
        return None

    cslide_match = re.search(r"(<p:cSld\b[^>]*>)", part_xml)
    if not cslide_match:
        raise RuntimeError("PPTX slide/master/layout part has no p:cSld element")
    return (
        part_xml[:cslide_match.end()]
        + background_xml
        + part_xml[cslide_match.end():]
    )


def _dominant_variant(
    values_by_slide: dict[int, Any],
) -> tuple[Any | None, list[int]]:
    """Return the most common value and its slides, or None on a tie."""
    slides_by_value: dict[Any, list[int]] = {}
    for slide_num, value in sorted(values_by_slide.items()):
        slides_by_value.setdefault(value, []).append(slide_num)
    if not slides_by_value:
        return None, []
    best_count = max(len(slides) for slides in slides_by_value.values())
    dominant = [
        (value, slides)
        for value, slides in slides_by_value.items()
        if len(slides) == best_count
    ]
    if len(dominant) != 1:
        return None, []
    return dominant[0]


def _is_strict_majority(subset_size: int, total: int) -> bool:
    return subset_size >= 2 and subset_size * 2 > total


def _promote_common_slide_backgrounds_to_masters(
    extract_dir: Path,
    structure: PptxStructureContext,
    slide_count: int,
    *,
    verbose: bool = False,
) -> int:
    """Promote the majority slide background to its shared slide master.

    Every slide in the master group must carry an explicit background —
    a slide without one would start inheriting the promoted master fill.
    Minority slides keep their own slide-level background, which always
    overrides the master fill.
    """
    slides_by_master: dict[str, list[int]] = {}
    for slide_num in range(1, slide_count + 1):
        master_part = structure.slide_master_part(slide_num)
        slides_by_master.setdefault(master_part, []).append(slide_num)

    promoted = 0
    for master_part, slide_nums in slides_by_master.items():
        slide_backgrounds: dict[int, str] = {}
        for slide_num in slide_nums:
            slide_path = extract_dir / "ppt" / "slides" / f"slide{slide_num}.xml"
            slide_xml = slide_path.read_text(encoding="utf-8")
            background_xml = _extract_slide_background_xml(slide_xml)
            if not background_xml:
                slide_backgrounds = {}
                break
            slide_backgrounds[slide_num] = background_xml

        if not slide_backgrounds:
            continue
        background_xml, dominant_slides = _dominant_variant(slide_backgrounds)
        if background_xml is None:
            continue
        if not _is_strict_majority(len(dominant_slides), len(slide_nums)):
            continue

        master_path = extract_dir / master_part
        master_xml = master_path.read_text(encoding="utf-8")
        promoted_master_xml = _put_background_on_part(master_xml, background_xml)
        if promoted_master_xml is None:
            continue
        master_path.write_text(promoted_master_xml, encoding="utf-8")

        for slide_num in dominant_slides:
            slide_path = extract_dir / "ppt" / "slides" / f"slide{slide_num}.xml"
            slide_xml = slide_path.read_text(encoding="utf-8")
            slide_path.write_text(
                _remove_slide_background_xml(slide_xml),
                encoding="utf-8",
            )
            promoted += 1

    if verbose and promoted:
        print(f"  Baseline master background: promoted {promoted} slide background(s)")
    return promoted


_CHROME_TRACE_TOKENS = (
    "logo",
    "footer",
    "header",
    "watermark",
    "chrome",
    "pagenumber",
    "slidenumber",
    "pagenum",
    "slidenum",
)
_TOP_LEVEL_SHAPE_TAGS = {
    f"{{{PML_NS}}}sp",
    f"{{{PML_NS}}}grpSp",
    f"{{{PML_NS}}}pic",
    f"{{{PML_NS}}}cxnSp",
    f"{{{PML_NS}}}graphicFrame",
}
_REL_ATTRS = {
    f"{{{REL_NS}}}embed",
    f"{{{REL_NS}}}link",
    f"{{{REL_NS}}}id",
}


def _chrome_token_from_svg_id(svg_id: str | None) -> str | None:
    """Return the baseline chrome token encoded in a source SVG id."""
    if not svg_id:
        return None
    lower = svg_id.lower()
    compact = re.sub(r"[-_\s]+", "", lower)
    if compact in _CHROME_TRACE_TOKENS:
        return compact
    split_tokens = {token for token in re.split(r"[-_\s]+", lower) if token}
    for token in _CHROME_TRACE_TOKENS:
        if token in split_tokens:
            return token
    return None


def _trace_chrome_shape_ids(
    trace: dict[str, Any] | None,
) -> dict[str, list[str]]:
    """Map chrome token to generated top-level shape ids for one slide."""
    result: dict[str, list[str]] = {}
    if not trace:
        return result
    for event in trace.get("events", []):
        if event.get("decision") != "native":
            continue
        semantic_role = event.get("data-pptx-role")
        placeholder = event.get("data-pptx-placeholder")
        has_explicit_semantics = (
            semantic_role is not None or placeholder is not None
        )
        token = (
            chrome_token_from_markers(semantic_role, placeholder)
            if has_explicit_semantics
            else _chrome_token_from_svg_id(event.get("id"))
        )
        shape_id = event.get("shape_id")
        if token and shape_id is not None:
            shape_ids = result.setdefault(token, [])
            normalized_shape_id = str(shape_id)
            if normalized_shape_id not in shape_ids:
                shape_ids.append(normalized_shape_id)
    return result


def _trace_native_shape_ids(
    trace: dict[str, Any] | None,
) -> dict[str, list[str]]:
    """Map every traced SVG id to generated top-level shape ids."""
    result: dict[str, list[str]] = {}
    if not trace:
        return result
    for event in trace.get("events", []):
        if event.get("decision") != "native":
            continue
        svg_id = event.get("id")
        shape_id = event.get("shape_id")
        if not svg_id or shape_id is None:
            continue
        shape_ids = result.setdefault(str(svg_id), [])
        normalized = str(shape_id)
        if normalized not in shape_ids:
            shape_ids.append(normalized)
    return result


def _shape_id(elem: ET.Element) -> str | None:
    for cnv in elem.iter(f"{{{PML_NS}}}cNvPr"):
        return cnv.attrib.get("id")
    return None


def _set_shape_name(elem: ET.Element, name: str) -> None:
    """Give one top-level shape a deterministic read-back identity."""
    for cnv in elem.iter(f"{{{PML_NS}}}cNvPr"):
        cnv.set("name", name)
        return
    raise TemplateStructureError(
        f"Cannot name structured shape {name!r}: p:cNvPr is missing"
    )


def _top_level_shape_name_roster(root: ET.Element) -> tuple[str, ...]:
    """Return the exact visible top-level shape-name sequence for read-back."""
    sp_tree = root.find(f".//{{{PML_NS}}}cSld/{{{PML_NS}}}spTree")
    if sp_tree is None:
        raise TemplateStructureError("Structured part has no p:cSld/p:spTree")
    names: list[str] = []
    for child in sp_tree:
        if child.tag not in _TOP_LEVEL_SHAPE_TAGS:
            continue
        c_nv_pr = next(child.iter(f"{{{PML_NS}}}cNvPr"), None)
        name = c_nv_pr.get("name") if c_nv_pr is not None else None
        if not name:
            raise TemplateStructureError(
                "Structured part contains a top-level shape without a name"
            )
        names.append(name)
    return tuple(names)


def _top_level_shapes_by_id(root: ET.Element) -> dict[str, ET.Element]:
    sp_tree = root.find(f".//{{{PML_NS}}}cSld/{{{PML_NS}}}spTree")
    if sp_tree is None:
        return {}
    shapes: dict[str, ET.Element] = {}
    for child in list(sp_tree):
        if child.tag not in _TOP_LEVEL_SHAPE_TAGS:
            continue
        shape_id = _shape_id(child)
        if shape_id:
            shapes[shape_id] = child
    return shapes


def _timing_shape_ids(root: ET.Element) -> set[str]:
    """Return slide-local shape ids referenced by animation timing."""
    return {
        elem.attrib["spid"]
        for elem in root.findall(f".//{{{PML_NS}}}timing//{{{PML_NS}}}spTgt")
        if elem.attrib.get("spid")
    }


def _relationship_ids_in_shape(elem: ET.Element) -> set[str]:
    rel_ids: set[str] = set()
    for node in elem.iter():
        for attr_name, value in node.attrib.items():
            if attr_name in _REL_ATTRS and value:
                rel_ids.add(value)
    return rel_ids


def _shape_relationships_supported(
    elem: ET.Element,
    rels: dict[str, dict[str, str]],
) -> bool:
    """Only image relationships are safe to copy into a slide master here."""
    for rel_id in _relationship_ids_in_shape(elem):
        attrs = rels.get(rel_id)
        if not attrs:
            return False
        if attrs.get("TargetMode"):
            return False
        if attrs.get("Type") != IMAGE_REL_TYPE:
            return False
    return True


def _canonical_shape_xml(
    elem: ET.Element,
    rels: dict[str, dict[str, str]],
) -> bytes:
    """Canonicalize ids and relationship ids for cross-slide equality."""
    clone = ET.fromstring(ET.tostring(elem, encoding="utf-8"))
    for cnv in clone.iter(f"{{{PML_NS}}}cNvPr"):
        cnv.set("id", "ID")
        # Generated names include the slide-local shape id (for example,
        # ``Image 2`` versus ``Image 8``) but do not affect rendering.
        if "name" in cnv.attrib:
            cnv.set("name", "NAME")
    for fld in clone.iter(f"{{{DML_NS}}}fld"):
        # The literal inside a slide-number field is a per-slide render
        # cache that PowerPoint recomputes from the slide position.
        if fld.attrib.get("type") == "slidenum":
            cached = fld.find(f"{{{DML_NS}}}t")
            if cached is not None:
                cached.text = ""
    for node in clone.iter():
        for attr_name, value in list(node.attrib.items()):
            if attr_name not in _REL_ATTRS:
                continue
            attrs = rels.get(value, {})
            node.set(
                attr_name,
                f"{attrs.get('Type', '')}|{attrs.get('Target', '')}",
            )
    return ET.tostring(clone, encoding="utf-8")


def _ensure_relationship(
    rels_path: Path,
    rel_type: str,
    target: str,
) -> str:
    existing = _find_relationship_id(rels_path, rel_type, target)
    if existing:
        return existing
    return _append_relationship(rels_path, rel_type, target)


def _copy_shape_relationships_to_part(
    elem: ET.Element,
    slide_rels: dict[str, dict[str, str]],
    target_rels_path: Path,
) -> ET.Element:
    """Clone a shape and retarget supported relationship ids to another part."""
    clone = ET.fromstring(ET.tostring(elem, encoding="utf-8"))
    for node in clone.iter():
        for attr_name, value in list(node.attrib.items()):
            if attr_name not in _REL_ATTRS:
                continue
            rel = slide_rels.get(value)
            if not rel:
                raise RuntimeError(f"Missing slide relationship for {value}")
            new_rid = _ensure_relationship(
                target_rels_path,
                rel["Type"],
                rel["Target"],
            )
            node.set(attr_name, new_rid)
    return clone


def _copy_shape_relationships_to_master(
    elem: ET.Element,
    slide_rels: dict[str, dict[str, str]],
    master_rels_path: Path,
) -> ET.Element:
    """Clone a shape and retarget supported relationship ids to the master."""
    return _copy_shape_relationships_to_part(elem, slide_rels, master_rels_path)


def _next_master_shape_id(master_xml: str) -> int:
    ids = [
        int(match)
        for match in re.findall(r"<p:cNvPr\b[^>]*\bid=\"(\d+)\"", master_xml)
    ]
    return max(ids, default=1) + 1


def _renumber_shape_ids(elem: ET.Element, start_id: int) -> None:
    next_id = start_id
    for cnv in elem.iter(f"{{{PML_NS}}}cNvPr"):
        cnv.set("id", str(next_id))
        next_id += 1


def _append_shape_to_master(master_path: Path, elem: ET.Element) -> None:
    master_xml = master_path.read_text(encoding="utf-8")
    _renumber_shape_ids(elem, _next_master_shape_id(master_xml))
    shape_xml = ET.tostring(elem, encoding="unicode")
    if "</p:spTree>" not in master_xml:
        raise RuntimeError(f"Slide master has no p:spTree: {master_path}")
    master_path.write_text(
        master_xml.replace("</p:spTree>", f"{shape_xml}\n</p:spTree>", 1),
        encoding="utf-8",
    )


def _append_shape_to_part(part_path: Path, elem: ET.Element) -> None:
    """Append a top-level shape to a master/layout spTree with fresh ids."""
    tree = ET.parse(part_path)
    root = tree.getroot()
    sp_tree = root.find(f".//{{{PML_NS}}}cSld/{{{PML_NS}}}spTree")
    if sp_tree is None:
        raise RuntimeError(f"PPTX part has no p:spTree: {part_path}")
    existing_ids = [
        int(cnv.attrib["id"])
        for cnv in root.iter(f"{{{PML_NS}}}cNvPr")
        if cnv.attrib.get("id", "").isdigit()
    ]
    clone = ET.fromstring(ET.tostring(elem, encoding="utf-8"))
    _renumber_shape_ids(clone, max(existing_ids, default=1) + 1)
    sp_tree.append(clone)
    _write_xml_tree(part_path, tree)


def _write_xml_tree(path: Path, tree: ET.ElementTree) -> None:
    tree.write(path, encoding="utf-8", xml_declaration=True)


COVER_LAYOUT_NAME = "Cover"


def _next_layout_part_number(extract_dir: Path) -> int:
    layouts_dir = extract_dir / "ppt" / "slideLayouts"
    numbers = [
        int(match.group(1))
        for path in layouts_dir.glob("slideLayout*.xml")
        if (match := re.fullmatch(r"slideLayout(\d+)\.xml", path.name))
    ]
    return max(numbers, default=0) + 1


def _next_slide_layout_id(extract_dir: Path) -> int:
    """Return a package-wide unused id for a new sldLayoutId entry."""
    ids: list[int] = []
    masters_dir = extract_dir / "ppt" / "slideMasters"
    for master_path in sorted(masters_dir.glob("slideMaster*.xml")):
        master_root = ET.parse(master_path).getroot()
        ids.extend(
            int(entry.attrib["id"])
            for entry in master_root.findall(
                f"{{{PML_NS}}}sldLayoutIdLst/{{{PML_NS}}}sldLayoutId"
            )
            if entry.attrib.get("id", "").isdigit()
        )
    presentation_path = extract_dir / "ppt" / "presentation.xml"
    if presentation_path.exists():
        ids.extend(
            int(value)
            for value in re.findall(
                r'\bid="(\d{9,})"', presentation_path.read_text(encoding="utf-8")
            )
        )
    next_id = max([*ids, 2147483648]) + 1
    if next_id > OOXML_UINT32_MAX:
        raise TemplateStructureError(
            "Cannot register another Slide Layout because the OOXML UInt32 "
            "identifier range is exhausted"
        )
    return next_id


def _create_cover_layout(extract_dir: Path, master_part: str, base_layout_part: str) -> str:
    """Clone a layout into a Cover layout that hides master shapes.

    Returns the new layout target relative to slide parts.
    """
    base_layout_path = extract_dir / base_layout_part
    layout_xml = base_layout_path.read_text(encoding="utf-8")

    root_match = re.search(r"<p:sldLayout\b[^>]*>", layout_xml)
    if not root_match:
        raise RuntimeError(f"Slide layout has no p:sldLayout root: {base_layout_part}")
    root_tag = root_match.group(0)
    if "showMasterSp=" in root_tag:
        new_root_tag = re.sub(r'showMasterSp="[^"]*"', 'showMasterSp="0"', root_tag)
    else:
        new_root_tag = root_tag[:-1] + ' showMasterSp="0">'
    layout_xml = layout_xml.replace(root_tag, new_root_tag, 1)
    layout_xml = re.sub(
        r"(<p:cSld\b[^>]*?)\s+name=\"[^\"]*\"",
        rf'\g<1> name="{COVER_LAYOUT_NAME}"',
        layout_xml,
        count=1,
    )

    layout_num = _next_layout_part_number(extract_dir)
    new_layout_part = f"ppt/slideLayouts/slideLayout{layout_num}.xml"
    new_layout_path = extract_dir / new_layout_part
    new_layout_path.write_text(layout_xml, encoding="utf-8")

    base_rels_path = _relationships_path_for_part(extract_dir, base_layout_part)
    new_rels_path = _relationships_path_for_part(extract_dir, new_layout_part)
    new_rels_path.parent.mkdir(exist_ok=True)
    new_rels_path.write_text(
        base_rels_path.read_text(encoding="utf-8"), encoding="utf-8"
    )

    master_path = extract_dir / master_part
    master_rels_path = _relationships_path_for_part(extract_dir, master_part)
    layout_target = posixpath.relpath(
        new_layout_part, posixpath.dirname(master_part)
    )
    rel_id = _append_relationship(master_rels_path, SLIDE_LAYOUT_REL_TYPE, layout_target)

    master_xml = master_path.read_text(encoding="utf-8")
    layout_id = _next_slide_layout_id(extract_dir)
    entry = f'<p:sldLayoutId id="{layout_id}" r:id="{rel_id}"/>'
    if "</p:sldLayoutIdLst>" not in master_xml:
        raise RuntimeError(f"Slide master has no sldLayoutIdLst: {master_part}")
    master_path.write_text(
        master_xml.replace("</p:sldLayoutIdLst>", f"{entry}</p:sldLayoutIdLst>", 1),
        encoding="utf-8",
    )

    content_types_path = extract_dir / "[Content_Types].xml"
    content_types_path.write_text(
        _add_content_type_override(
            content_types_path.read_text(encoding="utf-8"),
            new_layout_part,
            "application/vnd.openxmlformats-officedocument.presentationml.slideLayout+xml",
        ),
        encoding="utf-8",
    )
    return posixpath.relpath(new_layout_part, "ppt/slides")


def _create_custom_layout(
    extract_dir: Path,
    master_part: str,
    base_layout_part: str,
    layout_name: str,
    *,
    show_master_shapes: bool = True,
) -> tuple[str, str]:
    """Clone a clean custom layout and register it under its slide master.

    Returns ``(slide_relationship_target, package_part)``.
    """
    base_layout_path = extract_dir / base_layout_part
    tree = ET.parse(base_layout_path)
    root = tree.getroot()
    root.set("type", "cust")
    root.set("preserve", "1")
    root.set("showMasterSp", "1" if show_master_shapes else "0")

    c_sld = root.find(f"{{{PML_NS}}}cSld")
    if c_sld is None:
        raise RuntimeError(f"Slide layout has no p:cSld: {base_layout_part}")
    c_sld.set("name", layout_name)
    sp_tree = c_sld.find(f"{{{PML_NS}}}spTree")
    if sp_tree is None:
        raise RuntimeError(f"Slide layout has no p:spTree: {base_layout_part}")
    for child in list(sp_tree):
        if child.tag in _TOP_LEVEL_SHAPE_TAGS:
            sp_tree.remove(child)

    layout_num = _next_layout_part_number(extract_dir)
    new_layout_part = f"ppt/slideLayouts/slideLayout{layout_num}.xml"
    new_layout_path = extract_dir / new_layout_part
    _write_xml_tree(new_layout_path, tree)

    base_rels_path = _relationships_path_for_part(extract_dir, base_layout_part)
    new_rels_path = _relationships_path_for_part(extract_dir, new_layout_part)
    new_rels_path.parent.mkdir(exist_ok=True)
    new_rels_path.write_text(
        base_rels_path.read_text(encoding="utf-8"),
        encoding="utf-8",
    )
    master_target = posixpath.relpath(
        master_part,
        posixpath.dirname(new_layout_part),
    )
    rels_content = new_rels_path.read_text(encoding="utf-8")
    master_rel_ids = [
        rel_id
        for rel_id, attrs in _read_relationships(new_rels_path).items()
        if attrs.get("Type") == SLIDE_MASTER_REL_TYPE
    ]
    if len(master_rel_ids) != 1:
        raise RuntimeError(
            f"Cloned slide layout must have one Master relationship: {new_layout_part}"
        )
    master_rel_id = master_rel_ids[0]
    master_rel_pattern = re.compile(
        rf'(<Relationship\b[^>]*\bId="{re.escape(master_rel_id)}"'
        rf'[^>]*\bTarget=")[^"]*(")'
    )
    rels_content, replaced = master_rel_pattern.subn(
        rf"\g<1>{master_target}\g<2>",
        rels_content,
        count=1,
    )
    if replaced != 1:
        raise RuntimeError(
            f"Could not retarget cloned Layout to Master {master_part}"
        )
    new_rels_path.write_text(rels_content, encoding="utf-8")

    master_path = extract_dir / master_part
    master_rels_path = _relationships_path_for_part(extract_dir, master_part)
    layout_target = posixpath.relpath(new_layout_part, posixpath.dirname(master_part))
    rel_id = _append_relationship(master_rels_path, SLIDE_LAYOUT_REL_TYPE, layout_target)
    master_tree = ET.parse(master_path)
    master_root = master_tree.getroot()
    layout_list = master_root.find(f"{{{PML_NS}}}sldLayoutIdLst")
    if layout_list is None:
        raise RuntimeError(f"Slide master has no sldLayoutIdLst: {master_part}")
    layout_id = _next_slide_layout_id(extract_dir)
    ET.SubElement(
        layout_list,
        f"{{{PML_NS}}}sldLayoutId",
        {"id": str(layout_id), f"{{{REL_NS}}}id": rel_id},
    )
    _write_xml_tree(master_path, master_tree)

    content_types_path = extract_dir / "[Content_Types].xml"
    content_types_path.write_text(
        _add_content_type_override(
            content_types_path.read_text(encoding="utf-8"),
            new_layout_part,
            "application/vnd.openxmlformats-officedocument.presentationml.slideLayout+xml",
        ),
        encoding="utf-8",
    )
    return posixpath.relpath(new_layout_part, "ppt/slides"), new_layout_part


def _set_master_picker_name(master_path: Path, master_name: str) -> None:
    """Set the visible PowerPoint Master name on its common slide data."""
    tree = ET.parse(master_path)
    c_sld = tree.getroot().find(f"{{{PML_NS}}}cSld")
    if c_sld is None:
        raise RuntimeError(f"Slide master has no p:cSld: {master_path}")
    c_sld.set("name", master_name)
    _write_xml_tree(master_path, tree)


def _next_master_part_number(extract_dir: Path) -> int:
    numbers = [
        int(match.group(1))
        for path in (extract_dir / "ppt" / "slideMasters").glob("slideMaster*.xml")
        if (match := re.fullmatch(r"slideMaster(\d+)\.xml", path.name))
    ]
    return max(numbers, default=0) + 1


def _clone_structured_master(
    extract_dir: Path,
    source_master_part: str,
    master_name: str,
) -> str:
    """Clone a clean Master part and register it with the Presentation."""
    master_num = _next_master_part_number(extract_dir)
    master_part = f"ppt/slideMasters/slideMaster{master_num}.xml"
    master_path = extract_dir / master_part
    source_master_path = extract_dir / source_master_part
    shutil.copyfile(source_master_path, master_path)

    tree = ET.parse(master_path)
    root = tree.getroot()
    c_sld = root.find(f"{{{PML_NS}}}cSld")
    if c_sld is None:
        raise RuntimeError(f"Slide master has no p:cSld: {source_master_part}")
    c_sld.set("name", master_name)
    layout_list = root.find(f"{{{PML_NS}}}sldLayoutIdLst")
    if layout_list is None:
        raise RuntimeError(
            f"Slide master has no p:sldLayoutIdLst: {source_master_part}"
        )
    for entry in list(layout_list):
        layout_list.remove(entry)
    _write_xml_tree(master_path, tree)

    source_rels = _relationships_path_for_part(extract_dir, source_master_part)
    master_rels = _relationships_path_for_part(extract_dir, master_part)
    master_rels.parent.mkdir(parents=True, exist_ok=True)
    shutil.copyfile(source_rels, master_rels)
    for rel_id, attrs in tuple(_read_relationships(master_rels).items()):
        if attrs.get("Type") == SLIDE_LAYOUT_REL_TYPE:
            _remove_relationship(master_rels, rel_id)

    presentation_rels = extract_dir / "ppt" / "_rels" / "presentation.xml.rels"
    relationship_target = posixpath.relpath(master_part, "ppt")
    relationship_id = _append_relationship(
        presentation_rels,
        SLIDE_MASTER_REL_TYPE,
        relationship_target,
    )
    presentation_path = extract_dir / "ppt" / "presentation.xml"
    presentation_xml = presentation_path.read_text(encoding="utf-8")
    master_ids = [
        int(value)
        for value in re.findall(r'<p:sldMasterId\b[^>]*\bid="(\d+)"', presentation_xml)
    ]
    master_id = max(master_ids, default=(1 << 31) - 1) + 1
    if master_id > OOXML_UINT32_MAX:
        raise TemplateStructureError("Presentation Master id exceeds OOXML UInt32")
    entry = f'<p:sldMasterId id="{master_id}" r:id="{relationship_id}"/>'
    if "</p:sldMasterIdLst>" not in presentation_xml:
        raise RuntimeError("presentation.xml has no p:sldMasterIdLst")
    presentation_path.write_text(
        presentation_xml.replace(
            "</p:sldMasterIdLst>",
            f"{entry}</p:sldMasterIdLst>",
            1,
        ),
        encoding="utf-8",
    )

    content_types_path = extract_dir / "[Content_Types].xml"
    content_types_path.write_text(
        _add_content_type_override(
            content_types_path.read_text(encoding="utf-8"),
            master_part,
            "application/vnd.openxmlformats-officedocument.presentationml.slideMaster+xml",
        ),
        encoding="utf-8",
    )
    return master_part


def _assign_structured_masters(
    extract_dir: Path,
    structure: PptxStructureContext,
    specs: list[TemplateSlideSpec],
) -> dict[str, str]:
    """Create one registered package Master for each explicit SVG Master key."""
    if not specs:
        raise TemplateStructureError("Structured export requires at least one slide")
    source_master = structure.slide_master_part(specs[0].slide_num)
    masters: dict[str, str] = {}
    for spec in specs:
        existing_part = masters.get(spec.master_key)
        if existing_part is None:
            if not masters:
                existing_part = source_master
                _set_master_picker_name(
                    extract_dir / existing_part,
                    spec.master_name,
                )
            else:
                existing_part = _clone_structured_master(
                    extract_dir,
                    source_master,
                    spec.master_name,
                )
            masters[spec.master_key] = existing_part
        structure.slide_master_parts[spec.slide_num] = existing_part
    return masters


def _clear_master_placeholder_shapes(master_path: Path) -> None:
    """Remove base-package placeholders before installing structured content."""
    tree = ET.parse(master_path)
    root = tree.getroot()
    sp_tree = root.find(f".//{{{PML_NS}}}cSld/{{{PML_NS}}}spTree")
    if sp_tree is None:
        raise RuntimeError(f"Slide master has no p:spTree: {master_path}")
    for child in list(sp_tree):
        if child.find(f".//{{{PML_NS}}}ph") is not None:
            sp_tree.remove(child)
    _write_xml_tree(master_path, tree)


def _set_slide_layout_target(rels_path: Path, target: str) -> None:
    """Point a slide's layout relationship at a different layout part."""
    content = rels_path.read_text(encoding="utf-8")
    rel_id = None
    for existing_id, attrs in _read_relationships(rels_path).items():
        if attrs.get("Type") == SLIDE_LAYOUT_REL_TYPE:
            rel_id = existing_id
            break
    if rel_id is None:
        raise RuntimeError(f"No slide layout relationship in {rels_path}")
    pattern = re.compile(
        rf'(<Relationship\b[^>]*\bId="{re.escape(rel_id)}"[^>]*\bTarget=")[^"]*(")'
    )
    new_content, replaced = pattern.subn(rf"\g<1>{target}\g<2>", content, count=1)
    if not replaced:
        raise RuntimeError(f"Could not retarget layout relationship in {rels_path}")
    rels_path.write_text(new_content, encoding="utf-8")


_BASELINE_LAYOUT_ROLE_TOKENS = (
    ("Cover", frozenset({"cover", "frontcover"}), ("封面",)),
    (
        "Agenda",
        frozenset({"agenda", "contents", "outline", "toc"}),
        ("目录", "议程"),
    ),
    (
        "Section",
        frozenset({"chapter", "divider", "section", "transition"}),
        ("章节", "过渡页"),
    ),
    (
        "Closing",
        frozenset({"closing", "end", "ending", "qa", "thankyou", "thanks"}),
        ("封底", "结束", "结尾", "结语", "致谢", "谢谢"),
    ),
)


def _baseline_layout_role(svg_path: Path) -> str:
    """Use explicit page semantics, then fall back to conservative filename roles."""
    semantic_role = page_layout_name_from_svg(svg_path)
    if semantic_role:
        return semantic_role
    stem = svg_path.stem.casefold()
    tokens = {
        token
        for token in re.split(r"[^0-9a-z]+", stem)
        if token and not token.isdigit()
    }
    for role, english_tokens, cjk_tokens in _BASELINE_LAYOUT_ROLE_TOKENS:
        if tokens.intersection(english_tokens) or any(
            token in stem for token in cjk_tokens
        ):
            return role
    return "Content"


def _layout_identity(layout_path: Path) -> tuple[str, bool]:
    """Return a layout's picker name and master-shape visibility."""
    root = ET.parse(layout_path).getroot()
    c_sld = root.find(f"{{{PML_NS}}}cSld")
    name = c_sld.attrib.get("name", "") if c_sld is not None else ""
    return name, root.attrib.get("showMasterSp", "1") != "0"


def _shared_explicit_slide_background(
    extract_dir: Path,
    slide_nums: list[int],
) -> str | None:
    """Return one exact background only when every family slide carries it."""
    backgrounds: list[str] = []
    for slide_num in slide_nums:
        slide_path = extract_dir / "ppt" / "slides" / f"slide{slide_num}.xml"
        background = _extract_slide_background_xml(
            slide_path.read_text(encoding="utf-8")
        )
        if background is None:
            return None
        backgrounds.append(background)
    if not backgrounds or len(set(backgrounds)) != 1:
        return None
    return backgrounds[0]


def _extract_baseline_layout_families(
    extract_dir: Path,
    structure: PptxStructureContext,
    svg_files: list[Path],
    *,
    verbose: bool = False,
) -> int:
    """Build conservative post-generation layout families for a free SVG deck."""
    families: dict[tuple[str, str, bool], list[int]] = {}
    base_layouts: dict[tuple[str, str, bool], str] = {}
    for slide_num, svg_path in enumerate(svg_files, 1):
        rels_path = (
            extract_dir
            / "ppt"
            / "slides"
            / "_rels"
            / f"slide{slide_num}.xml.rels"
        )
        layout_target = _find_relationship_target(rels_path, SLIDE_LAYOUT_REL_TYPE)
        if not layout_target:
            raise RuntimeError(f"Slide {slide_num} has no slide layout relationship")
        layout_part = _resolve_package_target(
            f"ppt/slides/slide{slide_num}.xml", layout_target
        )
        layout_name, show_master_shapes = _layout_identity(extract_dir / layout_part)
        role = _baseline_layout_role(svg_path)
        if role == "Content" and layout_name == COVER_LAYOUT_NAME:
            role = COVER_LAYOUT_NAME
        key = (
            structure.slide_master_part(slide_num),
            role,
            show_master_shapes,
        )
        families.setdefault(key, []).append(slide_num)
        base_layouts.setdefault(key, layout_part)

    created = 0
    lifted_backgrounds = 0
    role_counts: dict[tuple[str, str], int] = {}
    for key, slide_nums in families.items():
        master_part, role, show_master_shapes = key
        role_key = (master_part, role)
        role_counts[role_key] = role_counts.get(role_key, 0) + 1
        variant = role_counts[role_key]
        layout_name = role if variant == 1 else f"{role} {variant}"
        layout_target, layout_part = _create_custom_layout(
            extract_dir,
            master_part,
            base_layouts[key],
            layout_name,
            show_master_shapes=show_master_shapes,
        )

        background_xml = _shared_explicit_slide_background(
            extract_dir, slide_nums
        )
        if background_xml is not None:
            layout_path = extract_dir / layout_part
            layout_xml = layout_path.read_text(encoding="utf-8")
            updated_layout_xml = _put_background_on_part(layout_xml, background_xml)
            if updated_layout_xml is not None:
                layout_path.write_text(updated_layout_xml, encoding="utf-8")
                for slide_num in slide_nums:
                    slide_path = (
                        extract_dir / "ppt" / "slides" / f"slide{slide_num}.xml"
                    )
                    slide_path.write_text(
                        _remove_slide_background_xml(
                            slide_path.read_text(encoding="utf-8")
                        ),
                        encoding="utf-8",
                    )
                lifted_backgrounds += len(slide_nums)

        for slide_num in slide_nums:
            rels_path = (
                extract_dir
                / "ppt"
                / "slides"
                / "_rels"
                / f"slide{slide_num}.xml.rels"
            )
            _set_slide_layout_target(rels_path, layout_target)
        created += 1

    if verbose and created:
        print(
            "  Baseline layout families: "
            f"created {created} reusable layout(s), "
            f"lifted {lifted_backgrounds} slide background(s)"
        )
    return created


def _promote_common_chrome_shapes_to_layouts(
    extract_dir: Path,
    slide_count: int,
    conversion_traces: list[dict[str, Any]] | None,
    *,
    verbose: bool = False,
) -> int:
    """Promote exact leading chrome shared by every slide in one layout."""
    if not conversion_traces:
        return 0
    trace_by_slide = {
        int(trace.get("slide_num", 0)): trace
        for trace in conversion_traces
        if trace.get("slide_num") is not None
    }
    if len(trace_by_slide) < slide_count:
        return 0

    slides_by_layout: dict[str, list[int]] = {}
    for slide_num in range(1, slide_count + 1):
        rels_path = (
            extract_dir
            / "ppt"
            / "slides"
            / "_rels"
            / f"slide{slide_num}.xml.rels"
        )
        layout_target = _find_relationship_target(rels_path, SLIDE_LAYOUT_REL_TYPE)
        if not layout_target:
            raise RuntimeError(f"Slide {slide_num} has no slide layout relationship")
        layout_part = _resolve_package_target(
            f"ppt/slides/slide{slide_num}.xml",
            layout_target,
        )
        slides_by_layout.setdefault(layout_part, []).append(slide_num)

    promoted = 0
    promoted_roles = 0
    for layout_part, slide_nums in slides_by_layout.items():
        if len(slide_nums) < 2:
            continue
        slide_state: dict[int, dict[str, Any]] = {}
        for slide_num in slide_nums:
            slide_path = extract_dir / "ppt" / "slides" / f"slide{slide_num}.xml"
            rels_path = (
                extract_dir
                / "ppt"
                / "slides"
                / "_rels"
                / f"slide{slide_num}.xml.rels"
            )
            tree = ET.parse(slide_path)
            root = tree.getroot()
            slide_state[slide_num] = {
                "path": slide_path,
                "rels": _read_relationships(rels_path),
                "root": root,
                "shapes": _top_level_shapes_by_id(root),
                "timing_shape_ids": _timing_shape_ids(root),
                "tokens": _trace_chrome_shape_ids(trace_by_slide.get(slide_num)),
                "tree": tree,
            }

        common_tokens = set.intersection(*(
            set(state["tokens"])
            for state in slide_state.values()
        ))
        common_tokens.difference_update(_PAGE_NUMBER_TOKENS)
        candidates: dict[str, dict[int, str]] = {}
        for token in sorted(common_tokens):
            shape_ids_by_slide: dict[int, str] = {}
            canonical_shapes: set[bytes] = set()
            for slide_num in slide_nums:
                state = slide_state[slide_num]
                shape_ids = state["tokens"].get(token, [])
                if len(shape_ids) != 1:
                    break
                shape_id = shape_ids[0]
                shape = state["shapes"].get(shape_id)
                if shape is None or shape_id in state["timing_shape_ids"]:
                    break
                if not _shape_relationships_supported(shape, state["rels"]):
                    break
                shape_ids_by_slide[slide_num] = shape_id
                canonical_shapes.add(_canonical_shape_xml(shape, state["rels"]))
            if len(shape_ids_by_slide) == len(slide_nums) and len(canonical_shapes) == 1:
                candidates[token] = shape_ids_by_slide

        if not candidates:
            continue

        # Layout shapes render behind slide-local shapes. Keep visual z-order by
        # accepting only the identical leading prefix shared by every family page.
        token_by_shape_id = {
            slide_num: {
                shape_ids[slide_num]: token
                for token, shape_ids in candidates.items()
            }
            for slide_num in slide_nums
        }
        leading_orders: list[list[str]] = []
        for slide_num in slide_nums:
            order: list[str] = []
            for shape_id in slide_state[slide_num]["shapes"]:
                token = token_by_shape_id[slide_num].get(shape_id)
                if token is None:
                    break
                order.append(token)
            leading_orders.append(order)
        safe_tokens = list(leading_orders[0])
        for order in leading_orders[1:]:
            common_length = 0
            for expected, actual in zip(safe_tokens, order):
                if expected != actual:
                    break
                common_length += 1
            safe_tokens = safe_tokens[:common_length]
            if not safe_tokens:
                break
        if not safe_tokens:
            continue

        layout_path = extract_dir / layout_part
        layout_rels_path = _relationships_path_for_part(extract_dir, layout_part)
        for token in safe_tokens:
            shape_ids_by_slide = candidates[token]
            first_slide = slide_nums[0]
            first_state = slide_state[first_slide]
            shape = first_state["shapes"][shape_ids_by_slide[first_slide]]
            layout_shape = _copy_shape_relationships_to_part(
                shape,
                first_state["rels"],
                layout_rels_path,
            )
            _append_shape_to_part(layout_path, layout_shape)
            promoted_roles += 1
            for slide_num, shape_id in shape_ids_by_slide.items():
                state = slide_state[slide_num]
                shape_to_remove = state["shapes"].get(shape_id)
                sp_tree = state["root"].find(
                    f".//{{{PML_NS}}}cSld/{{{PML_NS}}}spTree"
                )
                if sp_tree is not None and shape_to_remove is not None:
                    sp_tree.remove(shape_to_remove)
                    promoted += 1

        for slide_num in slide_nums:
            state = slide_state[slide_num]
            _write_xml_tree(state["path"], state["tree"])

    if verbose and promoted:
        print(
            "  Baseline layout chrome: "
            f"promoted {promoted} slide shape(s) across "
            f"{promoted_roles} shared object(s)"
        )
    return promoted


_TEMPLATE_PLACEHOLDER_PROMPTS = {
    "title": "Click to add title",
    "subtitle": "Click to add subtitle",
    "body": "Click to add text",
    "picture": "Click to add picture",
    "chart": "Click to add chart",
    "table": "Click to add table",
    "object": "Click to add content",
    "media": "Click to add media",
    "date": "Date",
    "footer": "Footer",
}

_PARAGRAPH_BULLET_CHOICE_TAGS = {
    f"{{{DML_NS}}}buNone",
    f"{{{DML_NS}}}buAutoNum",
    f"{{{DML_NS}}}buChar",
    f"{{{DML_NS}}}buBlip",
}
_PARAGRAPH_PROPERTIES_TRAILING_TAGS = {
    f"{{{DML_NS}}}tabLst",
    f"{{{DML_NS}}}defRPr",
    f"{{{DML_NS}}}extLst",
}


def _template_runtime_slides(
    extract_dir: Path,
    specs: list[TemplateSlideSpec],
    conversion_traces: list[dict[str, Any]] | None,
) -> list[_TemplateRuntimeSlide]:
    """Load slide XML state and join it with SVG-to-shape trace ids."""
    if not conversion_traces:
        raise TemplateStructureError(
            "Explicit Layout export requires native conversion traces for every slide"
        )
    trace_by_slide = {
        int(trace.get("slide_num", 0)): trace
        for trace in conversion_traces
        if trace.get("slide_num") is not None
    }
    states: list[_TemplateRuntimeSlide] = []
    for spec in specs:
        trace = trace_by_slide.get(spec.slide_num)
        if trace is None:
            raise TemplateStructureError(
                f"{spec.svg_path.name}: missing native conversion trace"
            )
        slide_path = extract_dir / "ppt" / "slides" / f"slide{spec.slide_num}.xml"
        rels_path = (
            extract_dir
            / "ppt"
            / "slides"
            / "_rels"
            / f"slide{spec.slide_num}.xml.rels"
        )
        tree = ET.parse(slide_path)
        root = tree.getroot()
        states.append(_TemplateRuntimeSlide(
            spec=spec,
            slide_path=slide_path,
            rels_path=rels_path,
            tree=tree,
            root=root,
            rels=_read_relationships(rels_path),
            shapes=_top_level_shapes_by_id(root),
            shape_ids_by_svg_id=_trace_native_shape_ids(trace),
        ))
    return states


def _template_shape_for_item(
    state: _TemplateRuntimeSlide,
    item: TemplateElementSpec,
) -> ET.Element | None:
    """Resolve one metadata item to its generated top-level DrawingML shape."""
    shape_ids = [
        shape_id
        for shape_id in state.shape_ids_by_svg_id.get(item.element_id, [])
        if shape_id in state.shapes
    ]
    if len(shape_ids) == 1:
        return state.shapes[shape_ids[0]]
    if not shape_ids and item.layer and item.order == 0 and item.tag in {"rect", "g"}:
        if _extract_slide_background_xml(
            state.slide_path.read_text(encoding="utf-8")
        ):
            return None
    if not shape_ids:
        text_hint = (
            "; multiline text placeholders require the default paragraph merge "
            "and cannot use --no-merge"
            if item.placeholder and item.tag == "text"
            else ""
        )
        raise TemplateStructureError(
            f"{state.spec.svg_path.name}: metadata element {item.element_id!r} "
            f"did not produce one top-level native shape{text_hint}"
        )
    raise TemplateStructureError(
        f"{state.spec.svg_path.name}: metadata element {item.element_id!r} "
        f"resolved to {len(shape_ids)} top-level shapes; use one direct SVG element"
    )


def _shape_transform(shape: ET.Element) -> ET.Element | None:
    """Return the direct DrawingML transform for one top-level shape."""
    paths = {
        f"{{{PML_NS}}}sp": f"{{{PML_NS}}}spPr/{{{DML_NS}}}xfrm",
        f"{{{PML_NS}}}pic": f"{{{PML_NS}}}spPr/{{{DML_NS}}}xfrm",
        f"{{{PML_NS}}}cxnSp": f"{{{PML_NS}}}spPr/{{{DML_NS}}}xfrm",
        f"{{{PML_NS}}}graphicFrame": f"{{{PML_NS}}}xfrm",
        f"{{{PML_NS}}}grpSp": f"{{{PML_NS}}}grpSpPr/{{{DML_NS}}}xfrm",
    }
    path = paths.get(shape.tag)
    return shape.find(path) if path is not None else None


def _int_attr(elem: ET.Element, name: str, context: str) -> int:
    try:
        return int(elem.attrib[name])
    except (KeyError, ValueError) as exc:
        raise TemplateStructureError(f"{context} has invalid {name!r}") from exc


def _flatten_group_transform(
    group: ET.Element,
    carrier: ET.Element,
    *,
    context: str,
) -> None:
    """Map a single group's child transform into slide coordinates."""
    group_xfrm = _shape_transform(group)
    carrier_xfrm = _shape_transform(carrier)
    if group_xfrm is None or carrier_xfrm is None:
        raise TemplateStructureError(
            f"{context} cannot be unwrapped because its DrawingML transform is missing"
        )
    if any(group_xfrm.get(name) is not None for name in ("rot", "flipH", "flipV")):
        raise TemplateStructureError(
            f"{context} wrapper carries an unsupported group rotation or flip"
        )
    group_off = group_xfrm.find(f"{{{DML_NS}}}off")
    group_ext = group_xfrm.find(f"{{{DML_NS}}}ext")
    child_off = group_xfrm.find(f"{{{DML_NS}}}chOff")
    child_ext = group_xfrm.find(f"{{{DML_NS}}}chExt")
    carrier_off = carrier_xfrm.find(f"{{{DML_NS}}}off")
    carrier_ext = carrier_xfrm.find(f"{{{DML_NS}}}ext")
    if any(value is None for value in (
        group_off,
        group_ext,
        child_off,
        child_ext,
        carrier_off,
        carrier_ext,
    )):
        raise TemplateStructureError(
            f"{context} cannot be unwrapped because its group transform is incomplete"
        )
    child_width = _int_attr(child_ext, "cx", context)
    child_height = _int_attr(child_ext, "cy", context)
    if child_width <= 0 or child_height <= 0:
        raise TemplateStructureError(f"{context} group child extent must be positive")
    scale_x = _int_attr(group_ext, "cx", context) / child_width
    scale_y = _int_attr(group_ext, "cy", context) / child_height
    mapped_x = _int_attr(group_off, "x", context) + round(
        (_int_attr(carrier_off, "x", context) - _int_attr(child_off, "x", context))
        * scale_x
    )
    mapped_y = _int_attr(group_off, "y", context) + round(
        (_int_attr(carrier_off, "y", context) - _int_attr(child_off, "y", context))
        * scale_y
    )
    carrier_off.set("x", str(mapped_x))
    carrier_off.set("y", str(mapped_y))
    carrier_ext.set(
        "cx",
        str(round(_int_attr(carrier_ext, "cx", context) * scale_x)),
    )
    carrier_ext.set(
        "cy",
        str(round(_int_attr(carrier_ext, "cy", context) * scale_y)),
    )


def _unwrap_placeholder_carrier(
    state: _TemplateRuntimeSlide,
    item: TemplateElementSpec,
) -> ET.Element:
    """Remove one SVG-only slot wrapper and return its top-level carrier."""
    wrapper = _template_shape_for_item(state, item)
    if wrapper is None:
        raise TemplateStructureError(
            f"{state.spec.svg_path.name}: placeholder {item.element_id!r} cannot "
            "be a slide background"
        )
    if item.tag != "g" or wrapper.tag != f"{{{PML_NS}}}grpSp":
        return wrapper
    carriers = [
        child for child in wrapper if child.tag in _TOP_LEVEL_SHAPE_TAGS
    ]
    if len(carriers) != 1:
        raise TemplateStructureError(
            f"{state.spec.svg_path.name}: placeholder group {item.element_id!r} "
            f"converted to {len(carriers)} native children; expected one carrier"
        )
    carrier = carriers[0]
    _flatten_group_transform(
        wrapper,
        carrier,
        context=f"{state.spec.svg_path.name} placeholder {item.element_id!r}",
    )
    sp_tree = _slide_sp_tree(state)
    try:
        wrapper_index = list(sp_tree).index(wrapper)
    except ValueError as exc:
        raise TemplateStructureError(
            f"{state.spec.svg_path.name}: placeholder wrapper "
            f"{item.element_id!r} is not top-level"
        ) from exc
    wrapper.remove(carrier)
    sp_tree.remove(wrapper)
    sp_tree.insert(wrapper_index, carrier)
    wrapper_id = _shape_id(wrapper)
    carrier_id = _shape_id(carrier)
    if wrapper_id:
        state.shapes.pop(wrapper_id, None)
    if carrier_id:
        state.shapes[carrier_id] = carrier
    return carrier


def _slide_sp_tree(state: _TemplateRuntimeSlide) -> ET.Element:
    sp_tree = state.root.find(f".//{{{PML_NS}}}cSld/{{{PML_NS}}}spTree")
    if sp_tree is None:
        raise RuntimeError(f"Slide has no p:spTree: {state.slide_path}")
    return sp_tree


def _append_shape_to_runtime_slide(
    state: _TemplateRuntimeSlide,
    elem: ET.Element,
) -> None:
    """Append a generated helper shape to one Slide with fresh object ids."""
    existing_ids = [
        int(cnv.attrib["id"])
        for cnv in state.root.iter(f"{{{PML_NS}}}cNvPr")
        if cnv.attrib.get("id", "").isdigit()
    ]
    clone = ET.fromstring(ET.tostring(elem, encoding="utf-8"))
    _renumber_shape_ids(clone, max(existing_ids, default=1) + 1)
    _slide_sp_tree(state).append(clone)


def _presentation_slide_size_emu(extract_dir: Path) -> tuple[int, int]:
    """Return the package slide size in EMU."""
    presentation_path = extract_dir / "ppt" / "presentation.xml"
    root = ET.parse(presentation_path).getroot()
    slide_size = root.find(f"{{{PML_NS}}}sldSz")
    if slide_size is None:
        raise TemplateStructureError("presentation.xml has no p:sldSz")
    try:
        width = int(slide_size.attrib["cx"])
        height = int(slide_size.attrib["cy"])
    except (KeyError, ValueError) as exc:
        raise TemplateStructureError("presentation.xml has an invalid p:sldSz") from exc
    if width <= 0 or height <= 0:
        raise TemplateStructureError("presentation.xml p:sldSz must be positive")
    return width, height


def _solid_background_xml_from_shape(
    shape: ET.Element,
    slide_size_emu: tuple[int, int],
) -> str | None:
    """Convert one exact full-slide solid rectangle shape into p:bg XML."""
    if shape.tag != f"{{{PML_NS}}}sp":
        return None
    if shape.find(f"{{{PML_NS}}}txBody") is not None:
        return None
    sp_pr = shape.find(f"{{{PML_NS}}}spPr")
    if sp_pr is None:
        return None
    xfrm = sp_pr.find(f"{{{DML_NS}}}xfrm")
    if xfrm is None or xfrm.attrib:
        return None
    off = xfrm.find(f"{{{DML_NS}}}off")
    ext = xfrm.find(f"{{{DML_NS}}}ext")
    if off is None or ext is None:
        return None
    try:
        bounds = (
            int(off.attrib["x"]),
            int(off.attrib["y"]),
            int(ext.attrib["cx"]),
            int(ext.attrib["cy"]),
        )
    except (KeyError, ValueError):
        return None
    if bounds != (0, 0, *slide_size_emu):
        return None

    geometry = sp_pr.find(f"{{{DML_NS}}}prstGeom")
    if geometry is None or geometry.attrib.get("prst") != "rect":
        return None
    solid_fill = sp_pr.find(f"{{{DML_NS}}}solidFill")
    if solid_fill is None:
        return None
    competing_fills = {
        f"{{{DML_NS}}}noFill",
        f"{{{DML_NS}}}gradFill",
        f"{{{DML_NS}}}blipFill",
        f"{{{DML_NS}}}pattFill",
        f"{{{DML_NS}}}grpFill",
    }
    if any(child.tag in competing_fills for child in sp_pr):
        return None
    line = sp_pr.find(f"{{{DML_NS}}}ln")
    if line is not None and line.find(f"{{{DML_NS}}}noFill") is None:
        return None
    for effect_tag in (f"{{{DML_NS}}}effectLst", f"{{{DML_NS}}}effectDag"):
        effect = sp_pr.find(effect_tag)
        if effect is not None and (effect.attrib or list(effect)):
            return None

    background = ET.Element(f"{{{PML_NS}}}bg")
    background_props = ET.SubElement(background, f"{{{PML_NS}}}bgPr")
    background_props.append(
        ET.fromstring(ET.tostring(solid_fill, encoding="utf-8"))
    )
    ET.SubElement(background_props, f"{{{DML_NS}}}effectLst")
    return ET.tostring(background, encoding="unicode")


def _remove_template_shape(
    state: _TemplateRuntimeSlide,
    shape: ET.Element,
) -> None:
    sp_tree = _slide_sp_tree(state)
    if shape not in list(sp_tree):
        raise TemplateStructureError(
            f"{state.spec.svg_path.name}: structure shape is not slide-local"
        )
    sp_tree.remove(shape)


def _move_template_background(
    states: list[_TemplateRuntimeSlide],
    target_path: Path,
) -> str:
    backgrounds = [
        _extract_slide_background_xml(state.slide_path.read_text(encoding="utf-8"))
        for state in states
    ]
    if not backgrounds or any(background is None for background in backgrounds):
        raise TemplateStructureError(
            "Template background metadata must resolve to an explicit background "
            "on every affected slide"
        )
    canonical_backgrounds = set()
    for background in backgrounds:
        if background is None:
            continue
        wrapper = ET.fromstring(
            f'<root xmlns:p="{PML_NS}" xmlns:a="{DML_NS}">{background}</root>'
        )
        canonical_backgrounds.add(
            ET.tostring(list(wrapper)[0], encoding="utf-8")
        )
    if len(canonical_backgrounds) != 1:
        slide_names = ", ".join(state.spec.svg_path.name for state in states)
        raise TemplateStructureError(
            f"Explicit template background differs across slides: {slide_names}"
        )
    background_xml = backgrounds[0]
    if background_xml is None:
        raise TemplateStructureError("Template background is unexpectedly empty")
    target_xml = target_path.read_text(encoding="utf-8")
    updated = _put_background_on_part(target_xml, background_xml)
    if updated is None:
        raise TemplateStructureError(
            f"Cannot install explicit background on {target_path.name}"
        )
    target_path.write_text(updated, encoding="utf-8")
    for state in states:
        c_sld = state.root.find(f"{{{PML_NS}}}cSld")
        background = (
            c_sld.find(f"{{{PML_NS}}}bg") if c_sld is not None else None
        )
        if c_sld is None or background is None:
            raise TemplateStructureError(
                f"{state.spec.svg_path.name}: explicit background disappeared "
                "during explicit Layout structure assembly"
            )
        c_sld.remove(background)
    return background_xml


def _move_template_solid_background_shapes(
    states: list[_TemplateRuntimeSlide],
    shapes: list[ET.Element],
    target_path: Path,
    slide_size_emu: tuple[int, int],
) -> str | None:
    """Move repeated full-slide solid rects into a master/layout p:bg."""
    backgrounds = [
        _solid_background_xml_from_shape(shape, slide_size_emu)
        for shape in shapes
    ]
    if not any(backgrounds):
        return None
    if any(background is None for background in backgrounds):
        raise TemplateStructureError(
            "A template background resolves to a full-slide solid rect on only "
            "some slides sharing the structure"
        )
    canonical = {background for background in backgrounds if background is not None}
    if len(canonical) != 1:
        slide_names = ", ".join(state.spec.svg_path.name for state in states)
        raise TemplateStructureError(
            f"Explicit template solid background differs across slides: {slide_names}"
        )
    background_xml = backgrounds[0]
    if background_xml is None:
        return None
    target_xml = target_path.read_text(encoding="utf-8")
    updated = _put_background_on_part(target_xml, background_xml)
    if updated is None:
        raise TemplateStructureError(
            f"Cannot install explicit solid background on {target_path.name}"
        )
    target_path.write_text(updated, encoding="utf-8")
    for state, shape in zip(states, shapes):
        _remove_template_shape(state, shape)
    return background_xml


def _set_slide_tree_background(
    state: _TemplateRuntimeSlide,
    background_xml: str,
) -> None:
    """Replace the slide tree's p:bg with explicit background XML."""
    c_sld = state.root.find(f"{{{PML_NS}}}cSld")
    if c_sld is None:
        raise TemplateStructureError(
            f"{state.spec.svg_path.name}: slide has no p:cSld"
        )
    for existing in list(c_sld):
        if existing.tag == f"{{{PML_NS}}}bg":
            c_sld.remove(existing)
    background = ET.fromstring(background_xml)
    sp_tree_tag = f"{{{PML_NS}}}spTree"
    insert_at = next(
        (index for index, child in enumerate(c_sld) if child.tag == sp_tree_tag),
        0,
    )
    c_sld.insert(insert_at, background)


def _apply_template_slide_backgrounds(
    states: list[_TemplateRuntimeSlide],
    slide_size_emu: tuple[int, int],
) -> dict[str, str]:
    """Compile one-page solid backgrounds into slide-level p:bg."""
    applied: dict[str, str] = {}
    for state in states:
        items = [
            item for item in state.spec.elements
            if item.layer == "slide" and item.is_background
        ]
        if not items:
            continue
        item = items[0]
        shape = _template_shape_for_item(state, item)
        if shape is None:
            c_sld = state.root.find(f"{{{PML_NS}}}cSld")
            background = (
                c_sld.find(f"{{{PML_NS}}}bg")
                if c_sld is not None
                else None
            )
            if background is None:
                raise TemplateStructureError(
                    f"{state.spec.svg_path.name}: slide background disappeared "
                    "during explicit Layout structure assembly"
                )
            applied[f"ppt/slides/slide{state.spec.slide_num}.xml"] = ET.tostring(
                background,
                encoding="unicode",
            )
            continue
        background_xml = _solid_background_xml_from_shape(shape, slide_size_emu)
        if background_xml is None:
            raise TemplateStructureError(
                f"{state.spec.svg_path.name}: {item.element_id!r} must remain an "
                "exact full-slide solid rectangle"
            )
        _remove_template_shape(state, shape)
        _set_slide_tree_background(state, background_xml)
        applied[f"ppt/slides/slide{state.spec.slide_num}.xml"] = background_xml
    return applied


def _move_template_static_shape(
    states: list[_TemplateRuntimeSlide],
    item: TemplateElementSpec,
    target_path: Path,
    target_rels_path: Path,
    slide_size_emu: tuple[int, int],
) -> str | None:
    shapes = [_template_shape_for_item(state, item) for state in states]
    if any(shape is None for shape in shapes):
        if not all(shape is None for shape in shapes):
            raise TemplateStructureError(
                f"{item.element_id}: structure item is a background on only some slides"
            )
        return _move_template_background(states, target_path)

    resolved_shapes = [shape for shape in shapes if shape is not None]
    if item.is_background:
        background_xml = _move_template_solid_background_shapes(
            states,
            resolved_shapes,
            target_path,
            slide_size_emu,
        )
        if background_xml is None:
            raise TemplateStructureError(
                f"{item.element_id!r} must compile to one exact p:bg payload"
            )
        return background_xml
    canonical = {
        _canonical_shape_xml(shape, state.rels)
        for state, shape in zip(states, resolved_shapes)
    }
    if len(canonical) != 1:
        slide_names = ", ".join(state.spec.svg_path.name for state in states)
        raise TemplateStructureError(
            f"Explicit structure element {item.element_id!r} differs across slides: "
            f"{slide_names}"
        )
    for state, shape in zip(states, resolved_shapes):
        shape_id = _shape_id(shape)
        if shape_id and shape_id in _timing_shape_ids(state.root):
            raise TemplateStructureError(
                f"{state.spec.svg_path.name}: structure element {item.element_id!r} "
                "is referenced by slide timing"
            )
        if not _shape_relationships_supported(shape, state.rels):
            raise TemplateStructureError(
                f"{state.spec.svg_path.name}: structure element {item.element_id!r} "
                "uses a non-image or external relationship"
            )

    prototype_state = states[0]
    prototype_shape = resolved_shapes[0]
    target_shape = _copy_shape_relationships_to_part(
        prototype_shape,
        prototype_state.rels,
        target_rels_path,
    )
    _set_shape_name(target_shape, f"{item.element_id} {item.layer.title()}")
    _append_shape_to_part(target_path, target_shape)
    for state, shape in zip(states, resolved_shapes):
        _remove_template_shape(state, shape)
    return None


def _shape_bounds_emu(
    shape: ET.Element,
    override_px: tuple[float, float, float, float] | None,
) -> tuple[int, int, int, int]:
    if override_px is not None:
        x, y, width, height = override_px
        return tuple(
            round(value * EMU_PER_PX)
            for value in (x, y, width, height)
        )

    xfrm = shape.find(f"{{{PML_NS}}}spPr/{{{DML_NS}}}xfrm")
    if xfrm is None:
        xfrm = shape.find(f"{{{PML_NS}}}xfrm")
    if xfrm is None:
        raise TemplateStructureError(
            "Placeholder shape has no directly readable DrawingML transform; "
            "set data-pptx-placeholder-bounds"
        )
    off = xfrm.find(f"{{{DML_NS}}}off")
    ext = xfrm.find(f"{{{DML_NS}}}ext")
    if off is None or ext is None:
        raise TemplateStructureError("Placeholder transform has no a:off/a:ext")
    try:
        return (
            int(off.attrib["x"]),
            int(off.attrib["y"]),
            int(ext.attrib["cx"]),
            int(ext.attrib["cy"]),
        )
    except (KeyError, ValueError) as exc:
        raise TemplateStructureError("Placeholder transform is invalid") from exc


def _replace_shape_xfrm(
    sp_pr: ET.Element,
    bounds: tuple[int, int, int, int],
) -> None:
    for existing in list(sp_pr):
        if existing.tag == f"{{{DML_NS}}}xfrm":
            sp_pr.remove(existing)
    x, y, width, height = bounds
    xfrm = ET.Element(f"{{{DML_NS}}}xfrm")
    ET.SubElement(xfrm, f"{{{DML_NS}}}off", {"x": str(x), "y": str(y)})
    ET.SubElement(
        xfrm,
        f"{{{DML_NS}}}ext",
        {"cx": str(width), "cy": str(height)},
    )
    sp_pr.insert(0, xfrm)


def _layout_level_one_paragraph_properties(
    list_style: ET.Element,
) -> ET.Element:
    """Return the Layout list style's level-one paragraph properties."""
    level_tag = f"{{{DML_NS}}}lvl1pPr"
    level_props = list_style.find(level_tag)
    if level_props is None:
        level_props = ET.Element(level_tag)
        trailing_tags = {
            f"{{{DML_NS}}}lvl{level}pPr" for level in range(2, 10)
        }
        trailing_tags.add(f"{{{DML_NS}}}extLst")
        insert_at = next(
            (
                index
                for index, child in enumerate(list_style)
                if child.tag in trailing_tags
            ),
            len(list_style),
        )
        list_style.insert(insert_at, level_props)
    return level_props


def _set_layout_level_one_default_size(
    list_style: ET.Element,
    source_run_pr: ET.Element | None,
) -> None:
    """Persist the prototype run size as the Layout's level-one text default."""
    if source_run_pr is None or source_run_pr.get("sz") is None:
        return
    level_props = _layout_level_one_paragraph_properties(list_style)
    default_props = level_props.find(f"{{{DML_NS}}}defRPr")
    if default_props is None:
        default_props = ET.Element(f"{{{DML_NS}}}defRPr")
        ext_tag = f"{{{DML_NS}}}extLst"
        insert_at = next(
            (
                index
                for index, child in enumerate(level_props)
                if child.tag == ext_tag
            ),
            len(level_props),
        )
        level_props.insert(insert_at, default_props)
    default_props.set("sz", source_run_pr.get("sz", ""))


def _set_no_bullet_paragraph_properties(
    paragraph_props: ET.Element,
    *,
    replace_existing: bool = False,
) -> None:
    """Disable inherited bullets and hanging indent for a prose paragraph."""
    if replace_existing:
        for child in list(paragraph_props):
            if child.tag in _PARAGRAPH_BULLET_CHOICE_TAGS:
                paragraph_props.remove(child)

    bullet_choice = next(
        (
            child
            for child in paragraph_props
            if child.tag in _PARAGRAPH_BULLET_CHOICE_TAGS
        ),
        None,
    )
    if bullet_choice is not None and bullet_choice.tag != f"{{{DML_NS}}}buNone":
        return
    if bullet_choice is None:
        insert_at = next(
            (
                index
                for index, child in enumerate(paragraph_props)
                if child.tag in _PARAGRAPH_PROPERTIES_TRAILING_TAGS
            ),
            len(paragraph_props),
        )
        paragraph_props.insert(insert_at, ET.Element(f"{{{DML_NS}}}buNone"))

    paragraph_props.set("marL", "0")
    paragraph_props.set("indent", "0")


def _placeholder_text_body(
    source_shape: ET.Element,
    item: TemplateElementSpec,
) -> ET.Element:
    tx_body = ET.Element(f"{{{PML_NS}}}txBody")
    source_tx_body = source_shape.find(f"{{{PML_NS}}}txBody")
    source_body_pr = (
        source_tx_body.find(f"{{{DML_NS}}}bodyPr")
        if source_tx_body is not None
        else None
    )
    source_lst_style = (
        source_tx_body.find(f"{{{DML_NS}}}lstStyle")
        if source_tx_body is not None
        else None
    )
    source_run_pr = (
        source_tx_body.find(f".//{{{DML_NS}}}rPr")
        if source_tx_body is not None
        else None
    )
    tx_body.append(
        ET.fromstring(ET.tostring(source_body_pr, encoding="utf-8"))
        if source_body_pr is not None
        else ET.Element(f"{{{DML_NS}}}bodyPr")
    )
    list_style = (
        ET.fromstring(ET.tostring(source_lst_style, encoding="utf-8"))
        if source_lst_style is not None
        else ET.Element(f"{{{DML_NS}}}lstStyle")
    )
    _set_layout_level_one_default_size(list_style, source_run_pr)
    if item.placeholder in {"body", "subtitle"}:
        _set_no_bullet_paragraph_properties(
            _layout_level_one_paragraph_properties(list_style),
            replace_existing=item.placeholder == "body",
        )
    tx_body.append(list_style)

    paragraph = ET.SubElement(tx_body, f"{{{DML_NS}}}p")
    if item.placeholder in {"body", "subtitle"}:
        paragraph_props = ET.SubElement(paragraph, f"{{{DML_NS}}}pPr")
        _set_no_bullet_paragraph_properties(paragraph_props)
    if item.placeholder in {"slide-number", "date"}:
        field_type = (
            "slidenum"
            if item.placeholder == "slide-number"
            else "datetimeFigureOut"
        )
        field = ET.SubElement(
            paragraph,
            f"{{{DML_NS}}}fld",
            {"id": f"{{{str(uuid.uuid4()).upper()}}}", "type": field_type},
        )
        if source_run_pr is not None:
            field.append(ET.fromstring(ET.tostring(source_run_pr, encoding="utf-8")))
        source_text = ""
        if source_tx_body is not None:
            source_text = "".join(
                text.text or ""
                for text in source_tx_body.findall(f".//{{{DML_NS}}}t")
            )
        field_text = (
            "‹#›"
            if item.placeholder == "slide-number"
            else source_text or "Date"
        )
        ET.SubElement(field, f"{{{DML_NS}}}t").text = field_text
    else:
        run = ET.SubElement(paragraph, f"{{{DML_NS}}}r")
        if source_run_pr is not None:
            run.append(ET.fromstring(ET.tostring(source_run_pr, encoding="utf-8")))
        ET.SubElement(run, f"{{{DML_NS}}}t").text = _TEMPLATE_PLACEHOLDER_PROMPTS.get(
            item.placeholder or "",
            "Click to add content",
        )
    ET.SubElement(paragraph, f"{{{DML_NS}}}endParaRPr", {"lang": "en-US"})
    return tx_body


def _set_placeholder_no_inherited_bullets(
    shape: ET.Element,
    item: TemplateElementSpec,
) -> None:
    """Keep prose bullet-free while preserving explicit subtitle bullets."""
    if item.placeholder not in {"body", "subtitle"}:
        return
    tx_body = shape.find(f"{{{PML_NS}}}txBody")
    if tx_body is None:
        return
    for paragraph in tx_body.findall(f"{{{DML_NS}}}p"):
        paragraph_props = paragraph.find(f"{{{DML_NS}}}pPr")
        if paragraph_props is None:
            paragraph_props = ET.Element(f"{{{DML_NS}}}pPr")
            paragraph.insert(0, paragraph_props)
        _set_no_bullet_paragraph_properties(
            paragraph_props,
            replace_existing=item.placeholder == "body",
        )


def _set_placeholder_theme_font_role(
    shape: ET.Element,
    item: TemplateElementSpec,
    theme_font_spec: ThemeFontSpec | None,
) -> None:
    """Force semantic text placeholders onto the correct theme font role."""
    if theme_font_spec is None:
        return
    if item.placeholder == "title":
        prefix = "+mj"
    elif item.placeholder in TEMPLATE_PLACEHOLDER_TYPES:
        prefix = "+mn"
    else:
        return
    for props_tag in ("rPr", "defRPr", "endParaRPr"):
        for props in shape.iter(f"{{{DML_NS}}}{props_tag}"):
            for font_tag, suffix in (("latin", "lt"), ("ea", "ea"), ("cs", "cs")):
                font = props.find(f"{{{DML_NS}}}{font_tag}")
                if font is not None:
                    font.set("typeface", f"{prefix}-{suffix}")


def _layout_placeholder_shape(
    source_shape: ET.Element,
    item: TemplateElementSpec,
    placeholder_idx: int | None,
    theme_font_spec: ThemeFontSpec | None = None,
) -> ET.Element:
    """Build one reusable p:sp placeholder from a prototype slide object."""
    placeholder_type = TEMPLATE_PLACEHOLDER_TYPES.get(item.placeholder or "")
    if placeholder_type is None:
        raise TemplateStructureError(
            f"Unsupported placeholder type: {item.placeholder!r}"
        )
    bounds = _shape_bounds_emu(source_shape, item.placeholder_bounds)
    shape = ET.Element(f"{{{PML_NS}}}sp")
    nv_sp_pr = ET.SubElement(shape, f"{{{PML_NS}}}nvSpPr")
    ET.SubElement(
        nv_sp_pr,
        f"{{{PML_NS}}}cNvPr",
        {"id": "2", "name": f"{item.element_id} Placeholder"},
    )
    c_nv_sp_pr = ET.SubElement(nv_sp_pr, f"{{{PML_NS}}}cNvSpPr")
    ET.SubElement(c_nv_sp_pr, f"{{{DML_NS}}}spLocks", {"noGrp": "1"})
    nv_pr = ET.SubElement(nv_sp_pr, f"{{{PML_NS}}}nvPr")
    placeholder_attrs = {"type": placeholder_type}
    if placeholder_idx is not None:
        placeholder_attrs["idx"] = str(placeholder_idx)
    elif item.placeholder != "title":
        raise TemplateStructureError(
            f"Placeholder {item.element_id!r} requires an idx"
        )
    ET.SubElement(nv_pr, f"{{{PML_NS}}}ph", placeholder_attrs)

    source_sp_pr = (
        source_shape.find(f"{{{PML_NS}}}spPr")
        if source_shape.tag == f"{{{PML_NS}}}sp"
        else None
    )
    if source_sp_pr is not None:
        sp_pr = ET.fromstring(ET.tostring(source_sp_pr, encoding="utf-8"))
    else:
        sp_pr = ET.Element(f"{{{PML_NS}}}spPr")
        geometry = ET.SubElement(sp_pr, f"{{{DML_NS}}}prstGeom", {"prst": "rect"})
        ET.SubElement(geometry, f"{{{DML_NS}}}avLst")
        ET.SubElement(sp_pr, f"{{{DML_NS}}}noFill")
        line = ET.SubElement(sp_pr, f"{{{DML_NS}}}ln")
        ET.SubElement(line, f"{{{DML_NS}}}noFill")
    _replace_shape_xfrm(sp_pr, bounds)
    shape.append(sp_pr)
    shape.append(_placeholder_text_body(source_shape, item))
    _set_placeholder_theme_font_role(shape, item, theme_font_spec)
    return shape


def _placeholder_binding_proxy(
    layout_placeholder: ET.Element,
    item: TemplateElementSpec,
) -> ET.Element:
    """Bind a Layout slot invisibly while leaving its visible content ordinary.

    An unbound object placeholder can leak its inherited empty frame into a
    finished Slide in non-PowerPoint renderers. A hidden matching proxy suppresses
    that inheritance. The zero-width transparent run avoids a LibreOffice empty-
    placeholder black fill without adding visible content.
    """
    proxy = ET.fromstring(ET.tostring(layout_placeholder, encoding="utf-8"))
    c_nv_pr = next(proxy.iter(f"{{{PML_NS}}}cNvPr"), None)
    if c_nv_pr is None:
        raise TemplateStructureError(
            f"Cannot create placeholder binding for {item.element_id!r}: "
            "p:cNvPr is missing"
        )
    placeholder = proxy.find(f".//{{{PML_NS}}}ph")
    if placeholder is None:
        raise TemplateStructureError(
            f"Cannot create placeholder binding for {item.element_id!r}: "
            "p:ph is missing"
        )
    c_nv_pr.set(
        "name",
        "Placeholder Binding "
        f"{placeholder.get('type', 'body')} {placeholder.get('idx', '0')}",
    )
    c_nv_pr.set("hidden", "1")
    tx_body = proxy.find(f"{{{PML_NS}}}txBody")
    if tx_body is None:
        raise TemplateStructureError(
            f"Cannot create placeholder binding for {item.element_id!r}: "
            "p:txBody is missing"
        )
    for child in list(tx_body):
        if child.tag == f"{{{DML_NS}}}p":
            tx_body.remove(child)
    paragraph = ET.SubElement(tx_body, f"{{{DML_NS}}}p")
    run = ET.SubElement(paragraph, f"{{{DML_NS}}}r")
    run_props = ET.SubElement(
        run,
        f"{{{DML_NS}}}rPr",
        {"lang": "en-US", "sz": "100"},
    )
    solid_fill = ET.SubElement(run_props, f"{{{DML_NS}}}solidFill")
    color = ET.SubElement(solid_fill, f"{{{DML_NS}}}srgbClr", {"val": "FFFFFF"})
    ET.SubElement(color, f"{{{DML_NS}}}alpha", {"val": "0"})
    ET.SubElement(run, f"{{{DML_NS}}}t").text = "\u200b"
    ET.SubElement(
        paragraph,
        f"{{{DML_NS}}}endParaRPr",
        {"lang": "en-US"},
    )
    return proxy


def _patch_slide_placeholder(
    shape: ET.Element,
    item: TemplateElementSpec,
    placeholder_idx: int | None,
    placeholder_type: str | None = None,
    theme_font_spec: ThemeFontSpec | None = None,
) -> None:
    resolved_type = placeholder_type or TEMPLATE_PLACEHOLDER_TYPES.get(
        item.placeholder or ""
    )
    if resolved_type is None:
        raise TemplateStructureError(
            f"Unsupported placeholder type: {item.placeholder!r}"
        )
    nv_paths = {
        f"{{{PML_NS}}}sp": f"{{{PML_NS}}}nvSpPr/{{{PML_NS}}}nvPr",
        f"{{{PML_NS}}}pic": f"{{{PML_NS}}}nvPicPr/{{{PML_NS}}}nvPr",
        f"{{{PML_NS}}}graphicFrame": (
            f"{{{PML_NS}}}nvGraphicFramePr/{{{PML_NS}}}nvPr"
        ),
    }
    nv_path = nv_paths.get(shape.tag)
    if nv_path is None:
        raise TemplateStructureError(
            f"Placeholder {item.element_id!r} converted to unsupported "
            f"DrawingML element {shape.tag.rsplit('}', 1)[-1]!r}; text/picture/"
            "native chart/table placeholders must remain one top-level object"
        )
    nv_pr = shape.find(nv_path)
    if nv_pr is None:
        raise TemplateStructureError(
            f"Placeholder {item.element_id!r} has no non-visual properties"
        )
    _set_placeholder_no_inherited_bullets(shape, item)
    _set_placeholder_theme_font_role(shape, item, theme_font_spec)
    for existing in list(nv_pr):
        if existing.tag == f"{{{PML_NS}}}ph":
            nv_pr.remove(existing)
    placeholder_attrs: dict[str, str] = {}
    if (
        placeholder_type is not None
        or resolved_type != "obj"
        or placeholder_idx is None
    ):
        placeholder_attrs["type"] = resolved_type
    if placeholder_idx is not None:
        placeholder_attrs["idx"] = str(placeholder_idx)
    ph = ET.Element(f"{{{PML_NS}}}ph", placeholder_attrs)
    ext_tag = f"{{{PML_NS}}}extLst"
    insert_at = next(
        (idx for idx, child in enumerate(nv_pr) if child.tag == ext_tag),
        len(nv_pr),
    )
    nv_pr.insert(insert_at, ph)


def _set_template_layout_header_footer(
    layout_path: Path,
    placeholders: tuple[TemplateElementSpec, ...],
) -> None:
    """Enable declared footer fields for slides newly created from the layout."""
    kinds = {item.placeholder for item in placeholders}
    if not kinds.intersection({"date", "footer", "slide-number"}):
        return
    tree = ET.parse(layout_path)
    root = tree.getroot()
    hf = root.find(f"{{{PML_NS}}}hf")
    if hf is None:
        hf = ET.Element(f"{{{PML_NS}}}hf")
        trailing_tags = {
            f"{{{PML_NS}}}timing",
            f"{{{PML_NS}}}transition",
            f"{{{PML_NS}}}extLst",
        }
        insert_at = next(
            (idx for idx, child in enumerate(root) if child.tag in trailing_tags),
            len(root),
        )
        root.insert(insert_at, hf)
    hf.set("hdr", "0")
    hf.set("dt", "1" if "date" in kinds else "0")
    hf.set("ftr", "1" if "footer" in kinds else "0")
    hf.set("sldNum", "1" if "slide-number" in kinds else "0")
    _write_xml_tree(layout_path, tree)


def _apply_explicit_layout_structure(
    extract_dir: Path,
    structure: PptxStructureContext,
    specs: list[TemplateSlideSpec],
    conversion_traces: list[dict[str, Any]] | None,
    theme_font_spec: ThemeFontSpec | None,
    *,
    verbose: bool = False,
) -> tuple[dict[str, str | None], dict[str, tuple[str, ...]]]:
    """Materialize explicit SVG master/layout/placeholder metadata into OOXML."""
    master_parts_by_key = _assign_structured_masters(
        extract_dir,
        structure,
        specs,
    )
    states = _template_runtime_slides(extract_dir, specs, conversion_traces)
    states_by_slide = {state.spec.slide_num: state for state in states}
    slide_size_emu = _presentation_slide_size_emu(extract_dir)

    expected_backgrounds: dict[str, str | None] = {}
    expected_shape_rosters: dict[str, tuple[str, ...]] = {}
    states_by_master: dict[str, list[_TemplateRuntimeSlide]] = {}
    for state in states:
        master_part = master_parts_by_key[state.spec.master_key]
        states_by_master.setdefault(master_part, []).append(state)
    master_shape_count = 0
    for master_part, master_states in states_by_master.items():
        master_path = extract_dir / master_part
        master_rels_path = _relationships_path_for_part(extract_dir, master_part)
        expected_backgrounds[master_part] = _extract_slide_background_xml(
            master_path.read_text(encoding="utf-8")
        )
        _clear_master_placeholder_shapes(master_path)
        master_items = master_states[0].spec.master_elements
        for item in master_items:
            background_xml = _move_template_static_shape(
                master_states,
                item,
                master_path,
                master_rels_path,
                slide_size_emu,
            )
            if background_xml is not None:
                expected_backgrounds[master_part] = background_xml
            master_shape_count += 1

    specs_by_layout: dict[str, list[TemplateSlideSpec]] = {}
    for spec in specs:
        specs_by_layout.setdefault(spec.layout_key, []).append(spec)
    placeholder_count = 0
    layout_shape_count = 0
    created_layout_parts: set[str] = set()
    for layout_key, layout_specs in specs_by_layout.items():
        layout_states = [states_by_slide[spec.slide_num] for spec in layout_specs]
        master_parts = {
            structure.slide_master_part(spec.slide_num) for spec in layout_specs
        }
        if len(master_parts) != 1:
            raise TemplateStructureError(
                f"Layout {layout_key!r} spans multiple slide masters; use distinct "
                "layout keys per master"
            )
        master_part = next(iter(master_parts))
        prototype = layout_specs[0]
        base_target = structure.slide_layout_target(prototype.slide_num)
        base_layout_part = _resolve_package_target(
            f"ppt/slides/slide{prototype.slide_num}.xml",
            base_target,
        )
        layout_target, layout_part = _create_custom_layout(
            extract_dir,
            master_part,
            base_layout_part,
            prototype.layout_name,
        )
        layout_path = extract_dir / layout_part
        layout_rels_path = _relationships_path_for_part(extract_dir, layout_part)
        created_layout_parts.add(layout_part)

        placeholder_bindings = {
            binding.element.element_id: binding
            for binding in template_placeholder_bindings(prototype)
        }
        for item in prototype.elements:
            if item.layer == "layout":
                background_xml = _move_template_static_shape(
                    layout_states,
                    item,
                    layout_path,
                    layout_rels_path,
                    slide_size_emu,
                )
                if background_xml is not None:
                    expected_backgrounds[layout_part] = background_xml
                layout_shape_count += 1
                continue
            if not item.placeholder:
                continue
            proxy_binding = is_proxy_placeholder(item)
            if proxy_binding:
                placeholder_shapes = [
                    _template_shape_for_item(state, item)
                    for state in layout_states
                ]
                if any(shape is None for shape in placeholder_shapes):
                    raise TemplateStructureError(
                        f"Placeholder {item.element_id!r} cannot be a slide background"
                    )
                resolved_shapes = [
                    shape for shape in placeholder_shapes if shape is not None
                ]
            else:
                resolved_shapes = [
                    _unwrap_placeholder_carrier(state, item)
                    for state in layout_states
                ]
            prototype_shape = resolved_shapes[0]
            binding = placeholder_bindings[item.element_id]
            assigned_idx = binding.assigned_idx
            layout_placeholder = _layout_placeholder_shape(
                prototype_shape,
                item,
                assigned_idx,
                theme_font_spec,
            )
            _append_shape_to_part(layout_path, layout_placeholder)
            for state, shape in zip(layout_states, resolved_shapes):
                if not proxy_binding:
                    _patch_slide_placeholder(
                        shape,
                        item,
                        assigned_idx,
                        theme_font_spec=theme_font_spec,
                    )
                    _set_shape_name(
                        shape,
                        f"{item.element_id} Placeholder Carrier",
                    )
                else:
                    _set_shape_name(
                        shape,
                        f"{item.element_id} Proxy Content",
                    )
                    _append_shape_to_runtime_slide(
                        state,
                        _placeholder_binding_proxy(layout_placeholder, item),
                    )
            placeholder_count += 1

        _set_template_layout_header_footer(layout_path, prototype.placeholders)
        for state in layout_states:
            _set_slide_layout_target(state.rels_path, layout_target)

    slide_backgrounds = _apply_template_slide_backgrounds(
        states,
        slide_size_emu,
    )
    expected_backgrounds.update(slide_backgrounds)
    for state in states:
        _write_xml_tree(state.slide_path, state.tree)
        expected_backgrounds.setdefault(
            f"ppt/slides/slide{state.spec.slide_num}.xml",
            None,
        )
        expected_shape_rosters[
            f"ppt/slides/slide{state.spec.slide_num}.xml"
        ] = _top_level_shape_name_roster(state.root)
    for part in states_by_master:
        expected_backgrounds.setdefault(
            part,
            _extract_slide_background_xml(
                (extract_dir / part).read_text(encoding="utf-8")
            ),
        )
    for part in created_layout_parts:
        expected_backgrounds.setdefault(part, None)
    for part in (*states_by_master, *created_layout_parts):
        expected_shape_rosters[part] = _top_level_shape_name_roster(
            ET.parse(extract_dir / part).getroot()
        )

    if verbose:
        print(
            "  Explicit Layout structure: "
            f"{len(states_by_master)} master(s), "
            f"{len(specs_by_layout)} layout(s), "
            f"{master_shape_count} master element(s), "
            f"{layout_shape_count} layout element(s), "
            f"{len(slide_backgrounds)} slide background(s), "
            f"{placeholder_count} placeholder definition(s)"
        )
    return expected_backgrounds, expected_shape_rosters


def _apply_preserved_structure(
    extract_dir: Path,
    specs: list[TemplateSlideSpec],
    contract: NativeStructureContract,
    conversion_traces: list[dict[str, Any]] | None,
    *,
    verbose: bool = False,
) -> None:
    """Drop preview-only inherited layers and bind content to source placeholders."""
    states = _template_runtime_slides(extract_dir, specs, conversion_traces)
    removed_preview_shapes = 0
    removed_preview_backgrounds = 0
    placeholder_count = 0
    for state in states:
        removed_background = False
        for item in state.spec.elements:
            if item.layer not in {"master", "layout"}:
                continue
            shape = _template_shape_for_item(state, item)
            if shape is not None:
                _remove_template_shape(state, shape)
                removed_preview_shapes += 1
                continue
            if removed_background:
                raise TemplateStructureError(
                    f"{state.spec.svg_path.name}: multiple inherited preview "
                    "backgrounds resolved to one slide background"
                )
            common_slide = state.root.find(f"{{{PML_NS}}}cSld")
            background = (
                common_slide.find(f"{{{PML_NS}}}bg")
                if common_slide is not None
                else None
            )
            if common_slide is None or background is None:
                raise TemplateStructureError(
                    f"{state.spec.svg_path.name}: inherited preview background "
                    "did not produce a removable slide background"
                )
            common_slide.remove(background)
            removed_background = True
            removed_preview_backgrounds += 1

        layout = contract.layout(state.spec.layout_key)
        for item, source_placeholder in match_native_placeholders(state.spec, layout):
            shape = _template_shape_for_item(state, item)
            if shape is None:
                raise TemplateStructureError(
                    f"{state.spec.svg_path.name}: placeholder {item.element_id!r} "
                    "cannot resolve to a slide background"
                )
            _patch_slide_placeholder(
                shape,
                item,
                source_placeholder.idx,
                source_placeholder.placeholder_type,
            )
            placeholder_count += 1
        _write_xml_tree(state.slide_path, state.tree)

    if verbose:
        print(
            "  Preserved structure: "
            f"{len({spec.layout_key for spec in specs})} source layout(s), "
            f"{removed_preview_shapes} preview shape(s) removed, "
            f"{removed_preview_backgrounds} preview background(s) removed, "
            f"{placeholder_count} source placeholder binding(s)"
        )


def _promote_common_chrome_shapes_to_masters(
    extract_dir: Path,
    structure: PptxStructureContext,
    slide_count: int,
    conversion_traces: list[dict[str, Any]] | None,
    *,
    verbose: bool = False,
) -> int:
    """Promote explicit repeated chrome SVG ids to their shared master."""
    if not conversion_traces:
        return 0
    trace_by_slide = {
        int(trace.get("slide_num", 0)): trace
        for trace in conversion_traces
        if trace.get("slide_num") is not None
    }
    if len(trace_by_slide) < slide_count:
        return 0

    slides_by_master: dict[str, list[int]] = {}
    for slide_num in range(1, slide_count + 1):
        master_part = structure.slide_master_part(slide_num)
        slides_by_master.setdefault(master_part, []).append(slide_num)

    promoted = 0
    promoted_roles = 0
    for master_part, slide_nums in slides_by_master.items():
        if len(slide_nums) < 2:
            continue
        slide_state: dict[int, dict[str, Any]] = {}
        for slide_num in slide_nums:
            slide_path = extract_dir / "ppt" / "slides" / f"slide{slide_num}.xml"
            rels_path = extract_dir / "ppt" / "slides" / "_rels" / f"slide{slide_num}.xml.rels"
            tree = ET.parse(slide_path)
            root = tree.getroot()
            slide_state[slide_num] = {
                "path": slide_path,
                "rels": _read_relationships(rels_path),
                "root": root,
                "shapes": _top_level_shapes_by_id(root),
                "timing_shape_ids": _timing_shape_ids(root),
                "tokens": _trace_chrome_shape_ids(trace_by_slide.get(slide_num)),
                "tree": tree,
            }

        # Per token, find the strict-majority identical variant. Slides
        # outside every dominant set become cover-layout minority slides.
        candidate_sets: dict[str, dict[int, str]] = {}
        all_tokens = sorted({
            token
            for state in slide_state.values()
            for token in state["tokens"]
        })
        for token in all_tokens:
            carriers: dict[int, str] = {}
            canonical_by_slide: dict[int, bytes] = {}
            for slide_num in slide_nums:
                state = slide_state[slide_num]
                shape_ids = state["tokens"].get(token, [])
                if len(shape_ids) != 1:
                    continue
                shape_id = shape_ids[0]
                shape = state["shapes"].get(shape_id)
                if shape is None:
                    continue
                if shape_id in state["timing_shape_ids"]:
                    continue
                if not _shape_relationships_supported(shape, state["rels"]):
                    continue
                carriers[slide_num] = shape_id
                canonical_by_slide[slide_num] = _canonical_shape_xml(
                    shape,
                    state["rels"],
                )
            dominant_xml, dominant_slides = _dominant_variant(canonical_by_slide)
            if dominant_xml is None:
                continue
            if not _is_strict_majority(len(dominant_slides), len(slide_nums)):
                continue
            candidate_sets[token] = {
                slide_num: carriers[slide_num] for slide_num in dominant_slides
            }

        if not candidate_sets:
            continue
        content_slides = sorted(set.intersection(
            *(set(slides) for slides in candidate_sets.values())
        ))
        if not _is_strict_majority(len(content_slides), len(slide_nums)):
            continue

        promotions: list[tuple[str, dict[int, str]]] = []
        claimed_shape_ids: dict[int, set[str]] = {
            slide_num: set() for slide_num in content_slides
        }
        for token in sorted(candidate_sets):
            shape_ids_by_slide = {
                slide_num: candidate_sets[token][slide_num]
                for slide_num in content_slides
            }
            # A flattened nested chrome group can emit several semantic trace
            # ids for the same generated DrawingML shape. Claim it once.
            if any(
                shape_ids_by_slide[slide_num] in claimed_shape_ids[slide_num]
                for slide_num in content_slides
            ):
                continue
            for slide_num, shape_id in shape_ids_by_slide.items():
                claimed_shape_ids[slide_num].add(shape_id)
            promotions.append((token, shape_ids_by_slide))

        if not promotions:
            continue

        # Master shapes always render behind slide-local shapes. Preserve the
        # original z-order by promoting only a common leading chrome prefix;
        # overlay headers/footers remain slide-local.
        token_by_shape_id = {
            slide_num: {
                shape_ids[slide_num]: token
                for token, shape_ids in promotions
            }
            for slide_num in content_slides
        }
        leading_token_orders: list[list[str]] = []
        for slide_num in content_slides:
            order: list[str] = []
            for shape_id in slide_state[slide_num]["shapes"]:
                token = token_by_shape_id[slide_num].get(shape_id)
                if token is None:
                    break
                order.append(token)
            leading_token_orders.append(order)

        safe_tokens = list(leading_token_orders[0])
        for order in leading_token_orders[1:]:
            common_length = 0
            for expected, actual in zip(safe_tokens, order):
                if expected != actual:
                    break
                common_length += 1
            safe_tokens = safe_tokens[:common_length]
            if not safe_tokens:
                break
        promotion_by_token = {token: shape_ids for token, shape_ids in promotions}
        promotions = [
            (token, promotion_by_token[token])
            for token in safe_tokens
        ]

        if not promotions:
            continue

        master_path = extract_dir / master_part
        master_rels_path = _relationships_path_for_part(extract_dir, master_part)
        for _token, shape_ids_by_slide in promotions:
            first_slide = content_slides[0]
            first_state = slide_state[first_slide]
            shape = first_state["shapes"][shape_ids_by_slide[first_slide]]
            master_shape = _copy_shape_relationships_to_master(
                shape,
                first_state["rels"],
                master_rels_path,
            )
            _append_shape_to_master(master_path, master_shape)
            promoted_roles += 1

            for slide_num, shape_id in shape_ids_by_slide.items():
                state = slide_state[slide_num]
                shape_to_remove = state["shapes"].get(shape_id)
                sp_tree = state["root"].find(f".//{{{PML_NS}}}cSld/{{{PML_NS}}}spTree")
                if sp_tree is not None and shape_to_remove is not None:
                    sp_tree.remove(shape_to_remove)
                    promoted += 1

        for slide_num in content_slides:
            state = slide_state[slide_num]
            _write_xml_tree(state["path"], state["tree"])

        # Minority slides (covers, section pages) keep every shape
        # slide-local and move to a Cover layout that hides the newly
        # promoted master chrome, so their rendering never changes.
        minority_slides = [
            slide_num for slide_num in slide_nums
            if slide_num not in set(content_slides)
        ]
        if minority_slides:
            first_minority_rels = (
                extract_dir / "ppt" / "slides" / "_rels"
                / f"slide{minority_slides[0]}.xml.rels"
            )
            base_target = _find_relationship_target(
                first_minority_rels, SLIDE_LAYOUT_REL_TYPE
            )
            if not base_target:
                raise RuntimeError(
                    f"Slide {minority_slides[0]} has no slide layout relationship"
                )
            base_layout_part = _resolve_package_target(
                f"ppt/slides/slide{minority_slides[0]}.xml", base_target
            )
            cover_target = _create_cover_layout(
                extract_dir, master_part, base_layout_part
            )
            for slide_num in minority_slides:
                rels_path = (
                    extract_dir / "ppt" / "slides" / "_rels"
                    / f"slide{slide_num}.xml.rels"
                )
                _set_slide_layout_target(rels_path, cover_target)
            if verbose:
                print(
                    "  Baseline cover layout: "
                    f"{len(minority_slides)} slide(s) keep slide-local chrome"
                )

    if verbose and promoted:
        print(
            "  Baseline master chrome: "
            f"promoted {promoted} slide shape(s) across {promoted_roles} shared object(s)"
        )
    return promoted


_PAGE_NUMBER_TOKENS = {"pagenumber", "pagenum", "slidenumber"}


def _first_slide_number(extract_dir: Path) -> int:
    """Read firstSlideNum from presentation.xml (defaults to 1)."""
    presentation_path = extract_dir / "ppt" / "presentation.xml"
    try:
        root = ET.parse(presentation_path).getroot()
    except (OSError, ET.ParseError):
        return 1
    raw = root.attrib.get("firstSlideNum")
    if raw is None:
        return 1
    try:
        return int(raw)
    except ValueError:
        return 1


def _shape_with_id(root: ET.Element, shape_id: str) -> ET.Element | None:
    """Find a p:sp anywhere in the slide tree by its cNvPr id."""
    for shape in root.iter(f"{{{PML_NS}}}sp"):
        cnv = shape.find(f"{{{PML_NS}}}nvSpPr/{{{PML_NS}}}cNvPr")
        if cnv is not None and cnv.attrib.get("id") == shape_id:
            return shape
    return None


def _replace_literal_run_with_slidenum_field(
    shape: ET.Element,
    expected_text: str,
    field_guid: str,
) -> bool:
    """Swap a single literal page-number run for an a:fld slidenum field."""
    tx_body = shape.find(f"{{{PML_NS}}}txBody")
    if tx_body is None:
        return False
    a_t = f"{{{DML_NS}}}t"
    total_text = "".join(t.text or "" for t in tx_body.iter(a_t))
    if total_text.strip() != expected_text:
        return False
    text_runs = [
        (paragraph, run)
        for paragraph in tx_body.iter(f"{{{DML_NS}}}p")
        for run in paragraph.findall(f"{{{DML_NS}}}r")
        if (run.findtext(a_t) or "").strip()
    ]
    if len(text_runs) != 1:
        return False
    paragraph, run = text_runs[0]
    if (run.findtext(a_t) or "").strip() != expected_text:
        return False

    fld = ET.Element(f"{{{DML_NS}}}fld", {"id": field_guid, "type": "slidenum"})
    r_pr = run.find(f"{{{DML_NS}}}rPr")
    if r_pr is not None:
        fld.append(ET.fromstring(ET.tostring(r_pr, encoding="utf-8")))
    fld_text = ET.SubElement(fld, a_t)
    fld_text.text = expected_text
    index = list(paragraph).index(run)
    paragraph.remove(run)
    paragraph.insert(index, fld)
    return True


def _convert_page_number_texts_to_fields(
    extract_dir: Path,
    slide_count: int,
    conversion_traces: list[dict[str, Any]] | None,
    *,
    context: str = "Baseline",
    verbose: bool = False,
) -> int:
    """Replace literal page-number chrome text with auto-updating fields.

    Only converts when the traced pageNumber/slideNumber shape's whole text
    equals the slide's expected display number (honoring firstSlideNum), so
    schemes like content-only numbering keep their literal text untouched.
    """
    if not conversion_traces:
        return 0
    trace_by_slide = {
        int(trace.get("slide_num", 0)): trace
        for trace in conversion_traces
        if trace.get("slide_num") is not None
    }
    first_slide_number = _first_slide_number(extract_dir)
    field_guid = f"{{{str(uuid.uuid4()).upper()}}}"

    converted = 0
    for slide_num in range(1, slide_count + 1):
        tokens = _trace_chrome_shape_ids(trace_by_slide.get(slide_num))
        shape_ids = sorted({
            shape_id
            for token, ids in tokens.items()
            if token in _PAGE_NUMBER_TOKENS
            for shape_id in ids
        })
        if len(shape_ids) != 1:
            continue
        slide_path = extract_dir / "ppt" / "slides" / f"slide{slide_num}.xml"
        tree = ET.parse(slide_path)
        shape = _shape_with_id(tree.getroot(), shape_ids[0])
        if shape is None:
            continue
        expected_text = str(first_slide_number + slide_num - 1)
        if _replace_literal_run_with_slidenum_field(shape, expected_text, field_guid):
            _write_xml_tree(slide_path, tree)
            converted += 1

    if verbose and converted:
        print(
            f"  {context} slide-number fields: "
            f"converted {converted} page number(s)"
        )
    return converted


def _remove_relationship(rels_path: Path, rel_id: str) -> None:
    """Remove one relationship entry by rId."""
    rels_content = rels_path.read_text(encoding="utf-8")
    pattern = re.compile(
        rf'[ \t]*<Relationship\b[^>]*\bId="{re.escape(rel_id)}"[^>]*/>[ \t]*\n?'
    )
    new_content, removed = pattern.subn("", rels_content, count=1)
    if not removed:
        raise RuntimeError(f"Relationship {rel_id} not found in {rels_path}")
    rels_path.write_text(new_content, encoding="utf-8")


def _remove_content_type_override(content_types_path: Path, part_name: str) -> None:
    """Remove the Override content-type entry for a deleted package part."""
    normalized = "/" + part_name.lstrip("/")
    content = content_types_path.read_text(encoding="utf-8")
    pattern = re.compile(
        rf'[ \t]*<Override\b[^>]*\bPartName="{re.escape(normalized)}"[^>]*/>[ \t]*\n?'
    )
    new_content, removed = pattern.subn("", content, count=1)
    if removed:
        content_types_path.write_text(new_content, encoding="utf-8")


def _prune_unused_slide_layouts(
    extract_dir: Path,
    structure: PptxStructureContext,
    slide_count: int,
    *,
    verbose: bool = False,
) -> int:
    """Remove base-template slide layouts no generated slide references.

    The python-pptx base package ships the full Office layout set; unused
    entries only pollute the PowerPoint new-slide picker. Layouts referenced
    by any generated slide are always kept, and a master keeps its layout
    list untouched unless at least one referenced layout remains in it.
    """
    # Read layout references live: earlier baseline passes may have rebound
    # minority slides to a Cover layout that the initial structure context
    # does not know about.
    referenced_layouts: set[str] = set()
    for slide_num in range(1, slide_count + 1):
        rels_path = (
            extract_dir / "ppt" / "slides" / "_rels" / f"slide{slide_num}.xml.rels"
        )
        target = _find_relationship_target(rels_path, SLIDE_LAYOUT_REL_TYPE)
        if not target:
            raise RuntimeError(f"Slide {slide_num} has no slide layout relationship")
        referenced_layouts.add(
            _resolve_package_target(f"ppt/slides/slide{slide_num}.xml", target)
        )

    pruned = 0
    content_types_path = extract_dir / "[Content_Types].xml"
    for master_part in sorted(set(structure.slide_master_parts.values())):
        master_path = extract_dir / master_part
        master_rels_path = _relationships_path_for_part(extract_dir, master_part)
        layout_rels = {
            rel_id: _resolve_package_target(master_part, attrs.get("Target", ""))
            for rel_id, attrs in _read_relationships(master_rels_path).items()
            if attrs.get("Type") == SLIDE_LAYOUT_REL_TYPE
        }
        if not any(part in referenced_layouts for part in layout_rels.values()):
            continue

        master_xml = master_path.read_text(encoding="utf-8")
        for rel_id, layout_part in sorted(layout_rels.items()):
            if layout_part in referenced_layouts:
                continue
            entry_re = re.compile(
                rf'[ \t]*<p:sldLayoutId\b[^>]*\br:id="{re.escape(rel_id)}"[^>]*/>[ \t]*\n?'
            )
            master_xml, removed = entry_re.subn("", master_xml, count=1)
            if not removed:
                raise RuntimeError(
                    f"Slide master {master_part} has no sldLayoutId entry for {rel_id}"
                )
            _remove_relationship(master_rels_path, rel_id)
            (extract_dir / layout_part).unlink()
            layout_rels_path = _relationships_path_for_part(extract_dir, layout_part)
            if layout_rels_path.exists():
                layout_rels_path.unlink()
            _remove_content_type_override(content_types_path, layout_part)
            pruned += 1
        master_path.write_text(master_xml, encoding="utf-8")

    if verbose and pruned:
        print(f"  Layout prune: removed {pruned} unused base layout(s)")
    return pruned


def _append_relationship(
    rels_path: Path,
    rel_type: str,
    target: str,
) -> str:
    """Append a relationship entry with the next available rId."""
    with open(rels_path, 'r', encoding='utf-8') as f:
        rels_content = f.read()

    rid_numbers = [int(match) for match in re.findall(r'Id="rId(\d+)"', rels_content)]
    next_rid = f'rId{max(rid_numbers, default=0) + 1}'
    rel_xml = (
        f'  <Relationship Id="{next_rid}" '
        f'Type="{rel_type}" Target="{target}"/>'
    )
    rels_content = rels_content.replace(
        '</Relationships>', rel_xml + '\n</Relationships>',
    )

    with open(rels_path, 'w', encoding='utf-8') as f:
        f.write(rels_content)

    return next_rid


def _add_default_content_type(content_types: str, extension: str, content_type: str) -> str:
    """Add a Default content type if it is not already present."""
    ext = extension.lstrip(".")
    if f'Extension="{ext}"' in content_types:
        return content_types
    entry = f'  <Default Extension="{ext}" ContentType="{content_type}"/>'
    override_pos = content_types.find('<Override ')
    if override_pos >= 0:
        return content_types[:override_pos] + entry + '\n' + content_types[override_pos:]
    return content_types.replace('</Types>', entry + '\n</Types>')


def _add_content_type_override(content_types: str, part_name: str, content_type: str) -> str:
    """Add an Override content type if it is not already present."""
    normalized = '/' + part_name.lstrip('/')
    if f'PartName="{normalized}"' in content_types:
        return content_types
    entry = f'  <Override PartName="{normalized}" ContentType="{content_type}"/>'
    return content_types.replace('</Types>', entry + '\n</Types>')


_IMAGE_CONTENT_TYPES = {
    'png': 'image/png',
    'jpg': 'image/jpeg',
    'jpeg': 'image/jpeg',
    'gif': 'image/gif',
    'webp': 'image/webp',
    'svg': 'image/svg+xml',
    'bmp': 'image/bmp',
    'emf': 'image/x-emf',
    'tif': 'image/tiff',
    'tiff': 'image/tiff',
    'wmf': 'image/x-wmf',
}


def _content_type_for_extension(ext: str) -> str:
    clean = ext.lower().lstrip('.')
    content_type = _IMAGE_CONTENT_TYPES.get(clean) or mimetypes.guess_type(f'x.{clean}')[0]
    if not content_type:
        raise ValueError(f"Unknown media content type for extension: {ext}")
    return content_type


def _as_dict(value: Any) -> dict[str, Any]:
    return value if isinstance(value, dict) else {}


def _create_writable_work_dir(output_path: Path) -> Path:
    """Create a real writable work directory for PPTX assembly."""
    parents = [output_path.parent, Path.cwd(), Path(tempfile.gettempdir())]
    seen: set[str] = set()
    errors: list[str] = []

    for parent in parents:
        parent = parent if str(parent) else Path(".")
        try:
            key = str(parent.resolve())
        except OSError:
            key = str(parent.absolute())
        if key in seen:
            continue
        seen.add(key)

        try:
            parent.mkdir(parents=True, exist_ok=True)
        except OSError as exc:
            errors.append(f"{parent}: cannot create parent ({exc})")
            continue

        for _ in range(3):
            work_dir = parent / f".pptx-build-{os.getpid()}-{uuid.uuid4().hex}"
            try:
                work_dir.mkdir(mode=0o700)
                probe_path = work_dir / ".write-probe"
                probe_path.write_text("ok", encoding="utf-8")
                probe_path.unlink()
                return work_dir
            except OSError as exc:
                errors.append(f"{work_dir}: {exc}")
                shutil.rmtree(work_dir, ignore_errors=True)

    details = "\n  - ".join(errors) if errors else "no candidate directories available"
    raise PermissionError(
        "Unable to create a writable PPTX work directory. "
        "Set the output path to a writable project directory or adjust sandbox permissions. "
        f"Tried:\n  - {details}"
    )


def _relax_output_permissions(output_path: Path) -> list[str]:
    """Make exported files readable outside the sandbox owner where possible."""
    warnings: list[str] = []

    try:
        current_mode = output_path.stat().st_mode
        readable_mode = (
            current_mode
            | stat.S_IRUSR
            | stat.S_IWUSR
            | stat.S_IRGRP
            | stat.S_IROTH
        )
        os.chmod(output_path, readable_mode)
    except OSError as exc:
        warnings.append(f"chmod skipped for {output_path}: {exc}")

    if os.name != 'nt':
        return warnings

    # Windows ACLs can remain sandbox-only even when the file mode looks sane.
    # Grant the built-in Users SID read access; the SID avoids localization
    # issues on non-English Windows installations.
    try:
        result = subprocess.run(
            ['icacls', str(output_path), '/grant', '*S-1-5-32-545:R'],
            capture_output=True,
            text=True,
            check=False,
        )
    except OSError as exc:
        warnings.append(f"icacls skipped for {output_path}: {exc}")
    else:
        if result.returncode != 0:
            message = (result.stderr or result.stdout or '').strip()
            details = f": {message}" if message else ''
            warnings.append(f"icacls failed for {output_path}{details}")

    return warnings


_NOTES_MASTER_REL_TYPE = (
    'http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesMaster'
)


def _ensure_notes_master(extract_dir: Path) -> None:
    """Create notesMaster parts and wire them into the presentation package."""
    ppt_dir = extract_dir / 'ppt'
    notes_masters_dir = ppt_dir / 'notesMasters'
    notes_masters_dir.mkdir(exist_ok=True)

    notes_master_path = notes_masters_dir / 'notesMaster1.xml'
    if not notes_master_path.exists():
        notes_master_path.write_text(create_notes_master_xml(), encoding='utf-8')

    theme_dir = ppt_dir / 'theme'
    theme_dir.mkdir(exist_ok=True)
    theme1_path = theme_dir / 'theme1.xml'
    theme2_path = theme_dir / 'theme2.xml'
    if not theme2_path.exists():
        if theme1_path.exists():
            shutil.copy2(theme1_path, theme2_path)
        else:
            raise RuntimeError('Cannot create notes theme: ppt/theme/theme1.xml is missing')

    notes_master_rels_dir = notes_masters_dir / '_rels'
    notes_master_rels_dir.mkdir(exist_ok=True)
    notes_master_rels_path = notes_master_rels_dir / 'notesMaster1.xml.rels'
    if not notes_master_rels_path.exists():
        notes_master_rels_path.write_text(
            create_notes_master_rels_xml(),
            encoding='utf-8',
        )

    presentation_rels_path = ppt_dir / '_rels' / 'presentation.xml.rels'
    notes_master_rid = _find_relationship_id(
        presentation_rels_path,
        _NOTES_MASTER_REL_TYPE,
        'notesMasters/notesMaster1.xml',
    )
    if notes_master_rid is None:
        notes_master_rid = _append_relationship(
            presentation_rels_path,
            _NOTES_MASTER_REL_TYPE,
            'notesMasters/notesMaster1.xml',
        )

    presentation_path = ppt_dir / 'presentation.xml'
    presentation_xml = presentation_path.read_text(encoding='utf-8')
    if '<p:notesMasterIdLst>' in presentation_xml:
        return
    notes_master_lst = (
        f'<p:notesMasterIdLst><p:notesMasterId r:id="{notes_master_rid}"/>'
        '</p:notesMasterIdLst>'
    )
    if '</p:sldMasterIdLst>' not in presentation_xml:
        raise RuntimeError('presentation.xml is missing p:sldMasterIdLst')
    presentation_xml = presentation_xml.replace(
        '</p:sldMasterIdLst>',
        '</p:sldMasterIdLst>' + notes_master_lst,
        1,
    )
    presentation_path.write_text(presentation_xml, encoding='utf-8')


def _slide_config(animation_config: dict[str, Any] | None, svg_stem: str) -> dict[str, Any]:
    if not animation_config:
        return {}
    slides_value = animation_config.get('slides', {})
    if not isinstance(slides_value, dict):
        raise ValueError('animations.json field "slides" must be an object')
    slide_value = slides_value.get(svg_stem, {})
    if not isinstance(slide_value, dict):
        raise ValueError(
            f'animations.json slide "{svg_stem}" must be an object'
        )
    return slide_value


def _slide_transition_settings(
    slide_cfg: dict[str, Any],
    transition: str | None,
    duration: float,
    auto_advance: float | None,
    cli_overrides: dict[str, bool],
) -> tuple[str | None, float, float | None]:
    trans_value = slide_cfg.get('transition', {})
    if not isinstance(trans_value, dict):
        raise ValueError('animations.json slide transition must be an object')
    trans_cfg = trans_value
    effect = transition
    if not cli_overrides.get('transition') and 'effect' in trans_cfg:
        raw_effect = trans_cfg['effect']
        if not isinstance(raw_effect, str):
            raise ValueError('animations.json transition effect must be a string')
        cfg_effect = normalize_transition_effect(raw_effect)
        effect = cfg_effect
    if not cli_overrides.get('transition_duration'):
        if 'duration' in trans_cfg:
            duration = validate_seconds(
                trans_cfg.get('duration'),
                "transition duration",
                allow_zero=effect is None,
            )
    if not cli_overrides.get('auto_advance') and 'auto_advance' in trans_cfg:
        auto_advance = validate_seconds(
            trans_cfg.get('auto_advance'),
            "transition auto_advance",
            allow_zero=True,
        )
    return effect, duration, auto_advance


def _slide_animation_settings(
    slide_cfg: dict[str, Any],
    animation: str | None,
    duration: float,
    stagger: float,
    trigger: str,
    cli_overrides: dict[str, bool],
) -> tuple[str | None, float, float, str]:
    anim_value = slide_cfg.get('animation', {})
    if not isinstance(anim_value, dict):
        raise ValueError('animations.json slide animation must be an object')
    anim_cfg = anim_value
    effect = normalize_animation_effect(
        animation,
        allow_none=True,
        allow_modes=True,
    )
    if not cli_overrides.get('animation') and 'effect' in anim_cfg:
        effect = normalize_animation_effect(
            anim_cfg.get('effect'),
            allow_none=True,
            allow_modes=True,
        )
    if not cli_overrides.get('animation_duration'):
        duration = validate_seconds(
            anim_cfg.get('duration', duration),
            'animation duration',
            allow_zero=False,
        )
    if not cli_overrides.get('animation_stagger'):
        stagger = validate_seconds(
            anim_cfg.get('stagger', stagger),
            'animation stagger',
            allow_zero=True,
        )
    if not cli_overrides.get('animation_trigger') and 'trigger' in anim_cfg:
        trigger = normalize_animation_trigger(anim_cfg.get('trigger'))
    else:
        trigger = normalize_animation_trigger(trigger)
    animation_seconds_to_milliseconds(
        duration,
        'animation duration',
        allow_zero=False,
    )
    animation_seconds_to_milliseconds(
        stagger,
        'animation stagger',
        allow_zero=True,
    )
    return effect, duration, stagger, trigger


def _build_sequence_targets(
    anim_targets: list[tuple[int, str]],
    slide_cfg: dict[str, Any],
    animation: str | None,
    duration: float,
    stagger: float,
    mixed_animation_offset: int,
    animation_rng: random.Random,
) -> tuple[list[tuple[int, int, str, float]], int]:
    groups_value = slide_cfg.get('groups', {})
    if not isinstance(groups_value, dict):
        raise ValueError('animations.json slide groups must be an object')
    groups_cfg = groups_value
    ordered: list[tuple[int, int, str, dict[str, Any]]] = []
    for idx, (sid, svg_id) in enumerate(anim_targets):
        group_value = groups_cfg.get(svg_id, {})
        if not isinstance(group_value, dict):
            raise ValueError(
                f'animations.json group "{svg_id}" must be an object'
            )
        group_cfg = group_value
        raw_effect = group_cfg.get('effect')
        if raw_effect is not None:
            normalized_effect = normalize_animation_effect(
                raw_effect,
                allow_none=True,
                allow_modes=True,
            )
        else:
            normalized_effect = None
        if 'effect' in group_cfg and normalized_effect is None:
            continue
        if animation is None and normalized_effect is None:
            continue
        order_value = group_cfg.get('order')
        order = order_value if order_value is not None else idx + 1
        if isinstance(order, bool) or not isinstance(order, int) or order <= 0:
            raise ValueError(
                f'animations.json group "{svg_id}" order must be a positive integer'
            )
        group_entry = dict(group_cfg)
        group_entry['_shape_id'] = sid
        group_entry['_effect'] = normalized_effect
        ordered.append((order, idx, svg_id, group_entry))

    ordered.sort(key=lambda item: (item[0], item[1]))

    seq_targets: list[tuple[int, int, str, float]] = []
    resolved_group_modes: list[str | None] = []
    for seq_idx, (_order, _original_idx, _svg_id, group_cfg) in enumerate(ordered):
        shape_id = int(group_cfg['_shape_id'])
        raw_effect = group_cfg.get('_effect')
        resolved_group_modes.append(
            raw_effect if raw_effect in ('auto', 'mixed', 'random') else None
        )
        if raw_effect in ('auto', 'mixed', 'random'):
            effect = pick_animation_effect(
                str(raw_effect), seq_idx, mixed_animation_offset, group_id=_svg_id,
                rng=animation_rng,
            )
        else:
            effect = str(raw_effect or pick_animation_effect(
                animation, seq_idx, mixed_animation_offset, group_id=_svg_id,
                rng=animation_rng,
            ))
        item_duration = validate_seconds(
            group_cfg.get('duration', duration),
            f'animation duration for group "{_svg_id}"',
            allow_zero=False,
        )
        delay_seconds = validate_seconds(
            group_cfg.get('delay', 0 if seq_idx == 0 else stagger),
            f'animation delay for group "{_svg_id}"',
            allow_zero=True,
        )
        delay_ms = animation_seconds_to_milliseconds(
            delay_seconds,
            f'animation delay for group "{_svg_id}"',
            allow_zero=True,
        )
        seq_targets.append((shape_id, delay_ms, effect, item_duration))

    mixed_count = 0
    if animation == 'mixed':
        mixed_count = sum(1 for _target in seq_targets[1:])
    elif animation == 'auto':
        # 'auto' accumulates a cross-slide offset so the image pool and the
        # unmatched-id fallback rotate as the deck advances. Single-effect
        # semantic matches (title→fade, chart→wipe etc.) are unaffected
        # because they ignore the offset.
        mixed_count = len(seq_targets)
    else:
        mixed_count = sum(
            1
            for seq_idx, mode in enumerate(resolved_group_modes)
            if mode == 'auto' or (mode == 'mixed' and seq_idx > 0)
        )
    return seq_targets, mixed_count


def _prerender_legacy_pngs(
    svg_files: list[Path],
    media_dir: Path,
    pixel_width: int,
    pixel_height: int,
    cache_dir: Path | None,
    workers: int,
    verbose: bool,
) -> dict[int, bool]:
    """Render every SVG→PNG into media_dir in parallel.

    Returns {1-based slide index: success}. Falls back to sequential when
    workers<=1 or len(svg_files)<=2.
    """
    results: dict[int, bool] = {}
    targets: list[tuple[int, Path, Path]] = [
        (i, svg, media_dir / f'image{i}.png')
        for i, svg in enumerate(svg_files, 1)
    ]

    if workers <= 1 or len(targets) <= 2:
        for i, svg, png in targets:
            ok = convert_svg_to_png_cached(svg, png, pixel_width, pixel_height, cache_dir)
            results[i] = ok
            if verbose:
                tag = 'cached/ok' if ok else 'failed'
                print(f"  [PNG {i}/{len(targets)}] {svg.name} - {tag}")
        return results

    with ProcessPoolExecutor(max_workers=workers) as pool:
        future_map = {
            pool.submit(
                convert_svg_to_png_cached,
                svg, png, pixel_width, pixel_height, cache_dir,
            ): (i, svg)
            for i, svg, png in targets
        }
        done = 0
        for future in as_completed(future_map):
            i, svg = future_map[future]
            try:
                ok = future.result()
            except Exception as exc:
                ok = False
                if verbose:
                    print(f"  [PNG] {svg.name} - worker error: {exc}")
            results[i] = ok
            done += 1
            if verbose:
                tag = 'cached/ok' if ok else 'failed'
                print(f"  [PNG {done}/{len(targets)}] {svg.name} - {tag}")

    return results


_OPC_UNRESERVED = frozenset(
    "ABCDEFGHIJKLMNOPQRSTUVWXYZabcdefghijklmnopqrstuvwxyz0123456789-._~"
)
_ASCII_LOWER_TRANSLATION = str.maketrans(
    "ABCDEFGHIJKLMNOPQRSTUVWXYZ",
    "abcdefghijklmnopqrstuvwxyz",
)


def _canonical_opc_part_path(path: str) -> str | None:
    """Return an OPC-equivalent package path key, or None when invalid."""
    if (
        not path
        or "\\" in path
        or path.endswith("/")
        or "//" in path
        or any(ord(char) <= 0x20 for char in path)
    ):
        return None
    output: list[str] = []
    index = 0
    while index < len(path):
        char = path[index]
        if char != "%":
            output.append(char)
            index += 1
            continue
        if index + 2 >= len(path) or not re.fullmatch(r"[0-9A-Fa-f]{2}", path[index + 1:index + 3]):
            return None
        value = int(path[index + 1:index + 3], 16)
        decoded = chr(value)
        if value in {0, ord("/"), ord("\\")}:
            return None
        output.append(decoded if decoded in _OPC_UNRESERVED else f"%{value:02X}")
        index += 3

    decoded_path = "".join(output)
    if decoded_path.rsplit("/", 1)[-1] in {".", ".."}:
        return None
    normalized = posixpath.normpath(decoded_path)
    if (
        not normalized
        or normalized in {".", ".."}
        or normalized.startswith("/")
        or normalized.startswith("../")
    ):
        return None
    return normalized.translate(_ASCII_LOWER_TRANSLATION)


def _source_part_for_rels(rels_path: str) -> str | None:
    """Return the source part path represented by a relationship sidecar."""
    filename = posixpath.basename(rels_path)
    if filename == ".rels" or not filename.endswith(".rels"):
        return None
    source_dir = posixpath.dirname(posixpath.dirname(rels_path))
    source_name = filename[:-len(".rels")]
    return posixpath.join(source_dir, source_name) if source_dir else source_name


def _resolve_internal_opc_target(rels_path: str, target: str) -> str | None:
    """Resolve one valid internal OPC Target to its canonical package key."""
    target_path_query = target.split("#", 1)[0]
    if (
        "\\" in target
        or "?" in target_path_query
        or any(ord(char) <= 0x20 for char in target)
    ):
        return None
    try:
        parsed = urlsplit(target)
    except ValueError:
        return None
    if parsed.scheme or parsed.netloc or parsed.query:
        return None

    source_part = _source_part_for_rels(rels_path)
    if parsed.path.startswith("/"):
        resolved = parsed.path[1:]
    elif parsed.path:
        base_dir = posixpath.dirname(source_part) if source_part else ""
        resolved = posixpath.join(base_dir, parsed.path) if base_dir else parsed.path
    elif source_part and "#" in target:
        resolved = source_part
    else:
        return None
    return _canonical_opc_part_path(resolved)


def _verify_internal_rels_targets(extract_dir: Path) -> list[str]:
    """Return a list of dangling internal Targets across every .rels in the package.

    Each entry is formatted as "<rels-path> -> <missing-target>". An empty list
    means every internal Target resolves to a real file in the package.
    """
    package_parts: set[str] = set()
    for path in extract_dir.rglob("*"):
        if not path.is_file():
            continue
        key = _canonical_opc_part_path(path.relative_to(extract_dir).as_posix())
        if key is not None:
            package_parts.add(key)
    problems: list[str] = []
    for rels_path in extract_dir.rglob('*.rels'):
        rels_rel = rels_path.relative_to(extract_dir).as_posix()
        try:
            root = ET.parse(rels_path).getroot()
        except ET.ParseError as exc:
            problems.append(f'{rels_rel} -> <invalid relationships XML: {exc}>')
            continue
        for elem in root:
            attrs = _relationship_attrs(elem)
            if attrs.get('TargetMode', '').lower() == 'external':
                continue
            target = attrs.get('Target')
            if not target:
                problems.append(f'{rels_rel} -> <missing Target>')
                continue
            resolved = _resolve_internal_opc_target(rels_rel, target)
            if resolved is None:
                problems.append(f'{rels_rel} -> <invalid Target {target!r}>')
                continue
            if resolved not in package_parts:
                problems.append(f'{rels_rel} -> {resolved}')
    return problems


def _presentation_format(width: float, height: float) -> str:
    """Map the slide aspect ratio to PowerPoint's PresentationFormat label.
    Non-standard ratios (square, portrait, banner crops) report 'Custom'.
    """
    if width <= 0 or height <= 0:
        return 'Custom'
    ratio = width / height
    for target, label in (
        (4 / 3, 'On-screen Show (4:3)'),
        (16 / 9, 'On-screen Show (16:9)'),
        (16 / 10, 'On-screen Show (16:10)'),
    ):
        if abs(ratio - target) < 0.02:
            return label
    return 'Custom'


def _stamp_docprops(
    extract_dir: Path,
    slide_count: int,
    pres_format: str,
    meta: dict[str, Any] | None = None,
) -> None:
    """Overwrite the misleading python-pptx default metadata with accurate
    values. Factual fields (slide count, export timestamp, presentation format,
    application) are always machine-derived. Authored fields — including the
    title — come solely from an optional per-project ``metadata.json``
    (``meta``); whatever it omits stays blank. ``lastModifiedBy`` follows
    ``creator`` rather than ever carrying the base template's author or a tool
    name. No field is guessed from slide content: a blank title is preferable
    to an unreliable heuristic pick.
    """
    meta = meta or {}

    def field(key: str, default: str = '') -> str:
        value = meta.get(key)
        return value.strip() if isinstance(value, str) and value.strip() else default

    title = field('title')
    creator = field('creator')

    now = datetime.now(timezone.utc).strftime('%Y-%m-%dT%H:%M:%SZ')

    core_path = extract_dir / 'docProps' / 'core.xml'
    if core_path.exists():
        core_path.write_text(
            "<?xml version='1.0' encoding='UTF-8' standalone='yes'?>\n"
            '<cp:coreProperties '
            'xmlns:cp="http://schemas.openxmlformats.org/package/2006/metadata/core-properties" '
            'xmlns:dc="http://purl.org/dc/elements/1.1/" '
            'xmlns:dcterms="http://purl.org/dc/terms/" '
            'xmlns:dcmitype="http://purl.org/dc/dcmitype/" '
            'xmlns:xsi="http://www.w3.org/2001/XMLSchema-instance">'
            f'<dc:title>{escape(title)}</dc:title>'
            f'<dc:subject>{escape(field("subject"))}</dc:subject>'
            f'<dc:creator>{escape(creator)}</dc:creator>'
            f'<cp:keywords>{escape(field("keywords"))}</cp:keywords>'
            f'<dc:description>{escape(field("description"))}</dc:description>'
            f'<dc:language>{escape(field("language"))}</dc:language>'
            f'<cp:lastModifiedBy>{escape(creator)}</cp:lastModifiedBy>'
            '<cp:revision>1</cp:revision>'
            f'<dcterms:created xsi:type="dcterms:W3CDTF">{now}</dcterms:created>'
            f'<dcterms:modified xsi:type="dcterms:W3CDTF">{now}</dcterms:modified>'
            f'<cp:category>{escape(field("category"))}</cp:category>'
            f'<cp:contentStatus>{escape(field("contentStatus"))}</cp:contentStatus>'
            '</cp:coreProperties>',
            encoding='utf-8',
        )

    app_path = extract_dir / 'docProps' / 'app.xml'
    if app_path.exists():
        app = app_path.read_text(encoding='utf-8')
        app = re.sub(r'<Slides>.*?</Slides>', f'<Slides>{slide_count}</Slides>', app)
        app = re.sub(
            r'<Company>.*?</Company>',
            f'<Company>{escape(field("company"))}</Company>',
            app,
        )
        app = re.sub(
            r'<Manager>.*?</Manager>',
            f'<Manager>{escape(field("manager"))}</Manager>',
            app,
        )
        app = re.sub(
            r'<Application>.*?</Application>',
            '<Application>Microsoft Office PowerPoint</Application>',
            app,
        )
        app = re.sub(
            r'<PresentationFormat>.*?</PresentationFormat>',
            f'<PresentationFormat>{escape(pres_format)}</PresentationFormat>',
            app,
        )
        app_path.write_text(app, encoding='utf-8')


def _create_preserved_base_pptx(
    contract: NativeStructureContract,
    specs: list[TemplateSlideSpec],
    output_path: Path,
    slide_size_emu: tuple[int, int],
) -> None:
    """Create empty slides bound to the source package's original layouts."""
    presentation = Presentation(str(contract.source_template))
    actual_size = (int(presentation.slide_width), int(presentation.slide_height))
    if actual_size != contract.slide_size_emu:
        raise TemplateStructureError(
            f"{contract.source_template.name} slide size does not match "
            f"{contract.contract_path.name}"
        )
    if actual_size != slide_size_emu:
        raise TemplateStructureError(
            "Generated SVG canvas does not match the preserved source template size"
        )

    layouts_by_part = {
        str(layout.part.partname).lstrip("/"): layout
        for master in presentation.slide_masters
        for layout in master.slide_layouts
    }
    slide_ids = presentation.slides._sldIdLst
    for slide_id in list(slide_ids):
        presentation.part.drop_rel(slide_id.rId)
        slide_ids.remove(slide_id)

    for spec in specs:
        layout_contract = contract.layout(spec.layout_key)
        layout = layouts_by_part.get(layout_contract.package_part)
        if layout is None:
            raise TemplateStructureError(
                f"Preserved source package did not load layout part "
                f"{layout_contract.package_part!r}"
            )
        presentation.slides.add_slide(layout)
    presentation.save(str(output_path))


def _clear_preserved_slide_collections(extract_dir: Path) -> None:
    """Remove source slide-order metadata that cannot apply to generated pages."""
    presentation_path = extract_dir / "ppt" / "presentation.xml"
    tree = ET.parse(presentation_path)
    root = tree.getroot()
    custom_shows = root.find(f"{{{PML_NS}}}custShowLst")
    if custom_shows is not None:
        root.remove(custom_shows)
    for extension_list in root.findall(f".//{{{PML_NS}}}extLst"):
        for extension in list(extension_list):
            if any(
                child.tag.rsplit("}", 1)[-1] == "sectionLst"
                for child in extension.iter()
                if isinstance(child.tag, str)
            ):
                extension_list.remove(extension)
    _write_xml_tree(presentation_path, tree)


def create_pptx_with_native_svg(
    svg_files: list[Path],
    output_path: Path,
    canvas_format: str | None = None,
    verbose: bool = True,
    transition: str | None = 'fade',
    transition_duration: float = 0.5,
    auto_advance: float | None = None,
    use_compat_mode: bool = True,
    notes: dict[str, str] | None = None,
    enable_notes: bool = True,
    use_native_shapes: bool = True,
    animation: str | None = None,
    animation_duration: float = 0.4,
    animation_stagger: float = 0.5,
    animation_trigger: str = 'after-previous',
    animation_config: dict[str, Any] | None = None,
    animation_cli_overrides: dict[str, bool] | None = None,
    narration_audio: dict[str, Path] | None = None,
    use_narration_timings: bool = False,
    narration_padding: float = 0.5,
    cache_dir: Path | None = None,
    workers: int | None = None,
    merge_paragraphs: bool = True,
    image_optimize: bool = True,
    image_max_dimension: int | None = 2560,
    image_sizing: str = 'cap',
    image_scale: float = 2.0,
    image_quality: int = 85,
    native_objects: bool = False,
    conversion_trace_path: Path | None = None,
    doc_metadata: dict[str, Any] | None = None,
    pptx_structure: str = "structured",
    native_structure_contract: NativeStructureContract | None = None,
    theme_font_spec: ThemeFontSpec | None = None,
    master_text_style_spec: MasterTextStyleSpec | None = None,
    theme_color_spec: ThemeColorSpec | None = None,
    structured_baseline: bool = False,
    baseline_layout_specs: list[TemplateSlideSpec] | None = None,
) -> bool:
    """Create a PPTX file with native DrawingML shapes.

    Args:
        svg_files: List of SVG files.
        output_path: Output PPTX path.
        canvas_format: Canvas format key.
        verbose: Whether to output detailed information.
        transition: Transition effect name.
        transition_duration: Transition duration in seconds.
        auto_advance: Auto-advance interval in seconds.
        use_compat_mode: Retained for API compatibility; ignored in native mode.
        notes: Notes dict, key is SVG stem, value is notes content.
        enable_notes: Whether to enable notes embedding.
        use_native_shapes: Must remain true; SVG-image PPTX export is unsupported.
        animation: Per-element entrance animation mode (single effect name,
            'mixed', 'random', or None to disable). Native shapes mode only.
        animation_duration: Per-element entrance duration in seconds.
        animation_stagger: Delay between elements in ``after-previous``
            trigger mode (seconds). Ignored otherwise.
        animation_trigger: PowerPoint Start mode — ``'after-previous'`` (default),
            ``'on-click'``, or ``'with-previous'``.
        animation_config: Optional sidecar overrides loaded from animations.json.
        animation_cli_overrides: Flags indicating explicit CLI overrides.
        narration_audio: Optional dict mapping SVG stem to narration audio file.
        use_narration_timings: Whether to set slide auto-advance from audio duration.
        narration_padding: Extra seconds added after each narration before advancing.
        image_optimize: Whether native export downscales oversized raster images.
        image_max_dimension: Maximum optimized image dimension in pixels.
        image_sizing: ``cap`` only limits source dimensions; ``display`` sizes
            from rendered SVG boxes.
        image_scale: Target image pixels per SVG display pixel.
        image_quality: JPEG quality used for opaque optimized rasters.
        native_objects: Convert explicit ``data-pptx-native`` table/chart
            markers to native PowerPoint objects. Default off.
        conversion_trace_path: Optional JSON path for native conversion diagnostics.
        pptx_structure: PPTX structure strategy. ``baseline`` promotes safe
            shared native backgrounds and leading chrome to slide masters,
            then extracts semantic page-role layout families and exact
            family-wide structurally marked leading chrome; marker-free legacy
            SVGs retain filename/id fallback;
            ``structured`` consumes explicit SVG master/layout/placeholder
            metadata; ``preserve`` reuses an imported source PPTX package;
            ``flat`` keeps generated structure slide-local.
        native_structure_contract: Validated source package contract for
            ``preserve`` mode.
        theme_font_spec: Locked project major/minor fonts for baseline/structured
            theme inheritance. Preserve and flat modes ignore this value.
        master_text_style_spec: Required locked title/body sizes for structured
            slide-master text styles. Other structure
            routes ignore this value.
        theme_color_spec: Locked project color scheme for context-aware
            baseline/structured theme inheritance. Preserve and flat modes
            ignore this value.
        structured_baseline: Obsolete compatibility argument; must remain false.
        baseline_layout_specs: Obsolete compatibility argument; must remain None.

    Returns:
        Whether all slides were successfully created.
    """
    if not use_native_shapes:
        raise ValueError(
            "SVG-image PPTX export is no longer supported; use svg_final/ "
            "directly for preview and native DrawingML PPTX for delivery"
        )
    if not svg_files:
        print("Error: No SVG files found")
        return False

    use_compat_mode = False
    if pptx_structure not in {"baseline", "structured", "preserve", "flat"}:
        raise ValueError(f"Unsupported pptx_structure: {pptx_structure}")
    if structured_baseline:
        raise ValueError(
            "structured_baseline is obsolete; use pptx_structure='structured'"
        )
    if baseline_layout_specs is not None:
        raise ValueError(
            "baseline_layout_specs is obsolete; structured export parses SVG metadata"
        )
    if pptx_structure == "structured" and master_text_style_spec is None:
        raise ValueError(
            "Structured export requires locked typography title/body sizes "
            "in master_text_style_spec"
        )
    if use_native_shapes and pptx_structure == "structured":
        template_specs = parse_template_slides(svg_files)
    elif use_native_shapes and pptx_structure == "preserve":
        if native_structure_contract is None:
            raise TemplateStructureError(
                "Preserve export requires a validated native structure contract"
            )
        template_specs = parse_preserve_slides(svg_files)
        for spec in template_specs:
            native_structure_contract.layout(spec.layout_key)
    else:
        template_specs = None
    template_background_expectations: dict[str, str | None] | None = None
    template_shape_roster_expectations: (
        dict[str, tuple[str, ...]] | None
    ) = None
    if template_specs is not None and not native_objects:
        native_placeholders = sorted({
            item.placeholder
            for spec in template_specs
            for item in spec.placeholders
            if item.placeholder in {"chart", "table"}
        })
        if native_placeholders:
            kinds = ", ".join(str(kind) for kind in native_placeholders)
            context = (
                pptx_structure.capitalize()
            )
            raise TemplateStructureError(
                f"{context} {kinds} placeholder(s) require --native-objects so each "
                "marker becomes one native PowerPoint object"
            )

    # Check compatibility mode dependencies
    renderer_name, renderer_status, renderer_hint = get_png_renderer_info()
    if not use_native_shapes and use_compat_mode and PNG_RENDERER is None:
        print("Warning: No PNG rendering library installed, cannot use compatibility mode")
        print(f"  {renderer_hint}")
        print("  Will use pure SVG mode (may not display in Office LTSC 2021 and similar versions)")
        use_compat_mode = False

    # Auto-detect canvas format or get dimensions from viewBox
    custom_pixels: tuple[int, int] | None = None
    if canvas_format is None:
        canvas_format = detect_format_from_svg(svg_files[0])
        if canvas_format and verbose:
            format_name = CANVAS_FORMATS.get(canvas_format, {}).get('name', canvas_format)
            print(f"  Detected canvas format: {format_name}")

    if canvas_format is None:
        custom_pixels = get_viewbox_dimensions(svg_files[0])
        if custom_pixels and verbose:
            print(f"  Using SVG viewBox dimensions: {custom_pixels[0]} x {custom_pixels[1]} px")

    if canvas_format is None and custom_pixels is None:
        canvas_format = 'ppt169'
        if verbose:
            print(f"  Using default format: PPT 16:9")

    width_emu, height_emu = get_slide_dimensions(canvas_format or 'ppt169', custom_pixels)
    pixel_width, pixel_height = get_pixel_dimensions(canvas_format or 'ppt169', custom_pixels)

    if verbose:
        print(f"  Slide dimensions: {pixel_width} x {pixel_height} px")
        print(f"  SVG file count: {len(svg_files)}")
        if use_native_shapes:
            print(f"  Mode: Native DrawingML shapes (directly editable)")
            print(
                "  Native table/chart objects: "
                f"{'Enabled' if native_objects else 'Disabled'}"
            )
            print(f"  PPTX structure: {pptx_structure}")
            if image_optimize:
                if image_sizing == 'display':
                    image_mode = (
                        f"display scale {image_scale:g}, "
                        f"max {image_max_dimension or 'unlimited'} px"
                    )
                else:
                    image_mode = f"cap max {image_max_dimension or 'unlimited'} px"
                print(
                    "  Image optimization: Enabled "
                    f"({image_mode}, JPEG q{image_quality})"
                )
            else:
                print("  Image optimization: Disabled")
        elif use_compat_mode:
            print(f"  Compatibility mode: Enabled (PNG + SVG dual format)")
            print(f"  PNG renderer: {renderer_name} {renderer_status}")
        else:
            print(f"  Compatibility mode: Disabled (pure SVG)")
        if transition:
            trans_name = TRANSITIONS.get(transition, {}).get('name', transition) if TRANSITIONS else transition
            print(f"  Transition effect: {trans_name}")
        if enable_notes and notes:
            print(f"  Speaker notes: {len(notes)} page(s)")
        elif enable_notes:
            print(f"  Speaker notes: Enabled (no notes files found)")
        else:
            print(f"  Speaker notes: Disabled")
        print()

    animation_cli_overrides = animation_cli_overrides or {}

    temp_dir = _create_writable_work_dir(output_path)

    try:
        base_pptx = temp_dir / 'base.pptx'
        if (
            use_native_shapes
            and pptx_structure == "preserve"
            and native_structure_contract is not None
            and template_specs is not None
        ):
            _create_preserved_base_pptx(
                native_structure_contract,
                template_specs,
                base_pptx,
                (width_emu, height_emu),
            )
        else:
            # Create the standard base PPTX with python-pptx.
            prs = Presentation()
            prs.slide_width = width_emu
            prs.slide_height = height_emu

            blank_layout = prs.slide_layouts[6]
            for _ in svg_files:
                prs.slides.add_slide(blank_layout)
            prs.save(str(base_pptx))

        # Extract PPTX
        extract_dir = temp_dir / 'pptx_content'
        with zipfile.ZipFile(base_pptx, 'r') as zf:
            zf.extractall(extract_dir)
        if use_native_shapes and pptx_structure == "preserve":
            _clear_preserved_slide_collections(extract_dir)
        active_theme_font_spec = (
            theme_font_spec
            if use_native_shapes and pptx_structure in {"baseline", "structured"}
            else None
        )
        if active_theme_font_spec is not None:
            apply_theme_font_spec(extract_dir, active_theme_font_spec)
        active_theme_color_spec = (
            theme_color_spec
            if use_native_shapes and pptx_structure in {"baseline", "structured"}
            else None
        )
        if active_theme_color_spec is not None:
            apply_theme_color_spec(extract_dir, active_theme_color_spec)
        structure = _read_slide_layout_targets(extract_dir, len(svg_files))

        media_dir = extract_dir / 'ppt' / 'media'
        media_dir.mkdir(exist_ok=True)

        prerender_results: dict[int, bool] | None = None
        if not use_native_shapes and use_compat_mode and PNG_RENDERER is not None:
            if workers is None:
                resolved_workers = min(os.cpu_count() or 2, len(svg_files), 8)
            else:
                resolved_workers = max(0, workers)
            if verbose:
                cache_label = str(cache_dir) if cache_dir else 'disabled'
                mode = f'parallel x{resolved_workers}' if resolved_workers > 1 else 'sequential'
                print(f"  Pre-rendering PNGs ({mode}, cache: {cache_label})")
            prerender_results = _prerender_legacy_pngs(
                svg_files, media_dir, pixel_width, pixel_height,
                cache_dir, resolved_workers, verbose,
            )
            if verbose:
                print()

        success_count = 0
        has_any_image = False
        media_cache: dict[tuple[str, str], str] = {}
        image_exts_used: set[str] = set()
        package_exts_used: set[str] = set()
        package_content_overrides: dict[str, str] = {}
        notes_slides_created: set[int] = set()
        narration_slides_created: set[int] = set()
        audio_exts_used: set[str] = set()
        package_uses_timings = False
        mixed_animation_offset = 0
        animation_seed = json.dumps(
            {
                'animation': animation,
                'config': animation_config,
                'slides': [path.name for path in svg_files],
            },
            ensure_ascii=False,
            sort_keys=True,
            separators=(',', ':'),
        )
        animation_rng = random.Random(animation_seed)
        conversion_trace: list[dict[str, Any]] | None = [] if conversion_trace_path else None
        structure_trace: list[dict[str, Any]] | None = (
            []
            if use_native_shapes and pptx_structure in {"baseline", "structured", "preserve"}
            else None
        )

        for i, svg_path in enumerate(svg_files, 1):
            slide_num = i
            expected_animation_targets: list[tuple[int, int, str, float]] = []
            expected_animation_duration = animation_duration
            expected_animation_trigger = normalize_animation_trigger(animation_trigger)

            try:
                # ---- Native shapes mode ----
                if use_native_shapes:
                    slide_cfg = _slide_config(animation_config, svg_path.stem)
                    slide_transition, slide_transition_duration, slide_auto_advance = (
                        _slide_transition_settings(
                            slide_cfg,
                            transition,
                            transition_duration,
                            auto_advance,
                            animation_cli_overrides,
                        )
                    )
                    (
                        slide_animation,
                        slide_animation_duration,
                        slide_animation_stagger,
                        slide_animation_trigger,
                    ) = _slide_animation_settings(
                        slide_cfg,
                        animation,
                        animation_duration,
                        animation_stagger,
                        animation_trigger,
                        animation_cli_overrides,
                    )
                    groups_value = slide_cfg.get('groups', {})
                    if not isinstance(groups_value, dict):
                        raise ValueError(
                            'animations.json slide groups must be an object'
                        )
                    animation_hard_disabled = (
                        animation_cli_overrides.get('animation', False)
                        and animation is None
                    )
                    explicit_animation_groups = (
                        frozenset(
                            str(group_id)
                            for group_id, group_cfg in groups_value.items()
                            if isinstance(group_cfg, dict)
                            and group_cfg.get('effect') != 'none'
                            and (
                                slide_animation is not None
                                or 'effect' in group_cfg
                            )
                        )
                        if not animation_hard_disabled
                        else frozenset()
                    )
                    (
                        slide_xml,
                        media_files_dict,
                        rel_entries,
                        anim_targets,
                        package_files_dict,
                        content_type_overrides,
                    ) = (
                        convert_svg_to_slide_shapes(
                            svg_path, slide_num=slide_num, verbose=verbose,
                            merge_paragraphs=merge_paragraphs,
                            image_optimize=image_optimize,
                            image_max_dimension=image_max_dimension,
                            image_sizing=image_sizing,
                            image_scale=image_scale,
                            image_quality=image_quality,
                            native_objects=native_objects,
                            animation_group_overrides=explicit_animation_groups,
                            theme_font_spec=active_theme_font_spec,
                            theme_color_spec=active_theme_color_spec,
                            trace_out=conversion_trace
                            if conversion_trace is not None
                            else structure_trace,
                        )
                    )
                    # Order matters: OOXML schema requires <p:transition>
                    # to precede <p:timing> inside <p:sld>. Both use the same
                    # </p:sld> string-replace anchor, so transition must be
                    # injected first and timing second.
                    if slide_transition is not None or slide_auto_advance is not None:
                        transition_fragment = create_transition_xml(
                            effect=slide_transition,
                            duration=slide_transition_duration,
                            advance_after=slide_auto_advance,
                        )
                        if transition_fragment:
                            slide_xml = slide_xml.replace(
                                '</p:sld>',
                                '\n' + transition_fragment + '\n</p:sld>',
                            )
                        if slide_auto_advance is not None:
                            package_uses_timings = True

                    expected_animation_duration = slide_animation_duration
                    expected_animation_trigger = slide_animation_trigger
                    if (
                        not animation_hard_disabled
                        and (slide_animation or explicit_animation_groups)
                        and anim_targets
                    ):
                        seq_targets, mixed_count = _build_sequence_targets(
                            anim_targets,
                            slide_cfg,
                            slide_animation,
                            slide_animation_duration,
                            slide_animation_stagger,
                            mixed_animation_offset,
                            animation_rng,
                        )
                        expected_animation_targets = seq_targets
                        if mixed_count:
                            mixed_animation_offset += mixed_count
                        timing_xml = '\n' + create_sequence_timing_xml(
                            seq_targets, duration=slide_animation_duration,
                            trigger=slide_animation_trigger,
                        )
                        slide_xml = slide_xml.replace(
                            '</p:sld>',
                            timing_xml + '\n</p:sld>',
                        )

                    # Write slide XML
                    slide_xml_path = extract_dir / 'ppt' / 'slides' / f'slide{slide_num}.xml'
                    with open(slide_xml_path, 'w', encoding='utf-8') as f:
                        f.write(slide_xml)

                    # Write media files
                    media_name_map: dict[str, str] = {}
                    for media_name, media_data in media_files_dict.items():
                        ext = media_name.rsplit('.', 1)[-1].lower()
                        media_hash = hashlib.sha256(media_data).hexdigest()
                        cache_key = (ext, media_hash)
                        cached_name = media_cache.get(cache_key)

                        if cached_name is None:
                            cached_name = f'image_{media_hash[:16]}.{ext}'
                            media_cache[cache_key] = cached_name
                            with open(media_dir / cached_name, 'wb') as f:
                                f.write(media_data)

                        media_name_map[media_name] = cached_name

                    for rel in rel_entries:
                        target = rel.get('target', '')
                        if not target.startswith('../media/'):
                            continue
                        media_name = target.split('../media/', 1)[1]
                        mapped_name = media_name_map.get(media_name)
                        if mapped_name:
                            rel['target'] = f'../media/{mapped_name}'

                    # Write non-media OOXML package parts produced by native
                    # object converters, e.g. chart XML, chart rels, and
                    # embedded workbooks.
                    for part_name, part_data in package_files_dict.items():
                        if (
                            part_name.startswith('ppt/charts/')
                            and part_name.endswith('.xml')
                        ):
                            part_data = rewrite_chart_accent_colors(
                                part_data,
                                active_theme_color_spec,
                            )
                        package_path = extract_dir / part_name
                        package_path.parent.mkdir(parents=True, exist_ok=True)
                        with open(package_path, 'wb') as f:
                            f.write(part_data)
                        suffix = package_path.suffix.lstrip('.').lower()
                        if suffix:
                            package_exts_used.add(suffix)
                    package_content_overrides.update(content_type_overrides)

                    # Build relationships XML
                    rels_dir = extract_dir / 'ppt' / 'slides' / '_rels'
                    rels_dir.mkdir(exist_ok=True)
                    rels_path = rels_dir / f'slide{slide_num}.xml.rels'

                    extra_rels = ''
                    for rel in rel_entries:
                        extra_rels += (
                            f'\n  <Relationship Id="{rel["id"]}" '
                            f'Type="{rel["type"]}" Target="{rel["target"]}"/>'
                        )

                    rels_xml = f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1"
                Type="{SLIDE_LAYOUT_REL_TYPE}"
                Target="{structure.slide_layout_target(slide_num)}"/>{extra_rels}
</Relationships>'''
                    with open(rels_path, 'w', encoding='utf-8') as f:
                        f.write(rels_xml)

                    # Track image formats for Content_Types
                    for media_name in media_name_map.values():
                        ext = media_name.rsplit('.', 1)[-1].lower()
                        _content_type_for_extension(ext)
                        image_exts_used.add(ext)
                        has_any_image = True

                # ---- Legacy SVG embedding mode ----
                else:
                    slide_cfg = _slide_config(animation_config, svg_path.stem)
                    slide_transition, slide_transition_duration, slide_auto_advance = (
                        _slide_transition_settings(
                            slide_cfg,
                            transition,
                            transition_duration,
                            auto_advance,
                            animation_cli_overrides,
                        )
                    )
                    svg_filename = f'image{i}.svg'
                    png_filename = f'image{i}.png'
                    png_rid = 'rId2'
                    svg_rid = 'rId3' if use_compat_mode else 'rId2'

                    shutil.copy(svg_path, media_dir / svg_filename)

                    slide_has_png = False
                    if use_compat_mode:
                        if prerender_results is not None:
                            png_success = prerender_results.get(i, False)
                        else:
                            png_path = media_dir / png_filename
                            png_success = convert_svg_to_png(
                                svg_path, png_path,
                                width=pixel_width, height=pixel_height,
                            )
                        if png_success:
                            slide_has_png = True
                            has_any_image = True
                            image_exts_used.add('png')
                        else:
                            if verbose:
                                print(
                                    f"  [{i}/{len(svg_files)}] {svg_path.name} - "
                                    "PNG generation failed, using pure SVG"
                                )
                            svg_rid = 'rId2'

                    slide_xml_path = extract_dir / 'ppt' / 'slides' / f'slide{slide_num}.xml'
                    slide_xml = create_slide_xml_with_svg(
                        slide_num,
                        png_rid=png_rid, svg_rid=svg_rid,
                        width_emu=width_emu, height_emu=height_emu,
                        transition=slide_transition,
                        transition_duration=slide_transition_duration,
                        auto_advance=slide_auto_advance,
                        use_compat_mode=(use_compat_mode and slide_has_png),
                    )
                    with open(slide_xml_path, 'w', encoding='utf-8') as f:
                        f.write(slide_xml)

                    rels_dir = extract_dir / 'ppt' / 'slides' / '_rels'
                    rels_dir.mkdir(exist_ok=True)
                    rels_path = rels_dir / f'slide{slide_num}.xml.rels'
                    rels_xml = create_slide_rels_xml(
                        png_rid=png_rid, png_filename=png_filename,
                        svg_rid=svg_rid, svg_filename=svg_filename,
                        use_compat_mode=(use_compat_mode and slide_has_png),
                        slide_layout_target=structure.slide_layout_target(slide_num),
                    )
                    with open(rels_path, 'w', encoding='utf-8') as f:
                        f.write(rels_xml)

                resolved_advance_after = slide_auto_advance
                resolved_advance_on_click = True

                # --- Process notes (shared between native and legacy mode) ---
                notes_content = ''
                if enable_notes:
                    svg_stem = svg_path.stem
                    notes_content = notes.get(svg_stem, '') if notes else ''
                    notes_text = markdown_to_plain_text(notes_content) if notes_content else ''
                    if notes_text:
                        _ensure_notes_master(extract_dir)

                        notes_slides_dir = extract_dir / 'ppt' / 'notesSlides'
                        notes_slides_dir.mkdir(exist_ok=True)

                        notes_xml_path = notes_slides_dir / f'notesSlide{slide_num}.xml'
                        notes_xml = create_notes_slide_xml(slide_num, notes_text)
                        with open(notes_xml_path, 'w', encoding='utf-8') as f:
                            f.write(notes_xml)

                        notes_rels_dir = notes_slides_dir / '_rels'
                        notes_rels_dir.mkdir(exist_ok=True)
                        notes_rels_path = notes_rels_dir / f'notesSlide{slide_num}.xml.rels'
                        notes_rels_xml = create_notes_slide_rels_xml(slide_num)
                        with open(notes_rels_path, 'w', encoding='utf-8') as f:
                            f.write(notes_rels_xml)

                        _append_relationship(
                            rels_path,
                            'http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesSlide',
                            f'../notesSlides/notesSlide{slide_num}.xml',
                        )
                        notes_slides_created.add(slide_num)

                # --- Process narration audio (shared between native and legacy mode) ---
                svg_stem = svg_path.stem
                audio_path = narration_audio.get(svg_stem) if narration_audio else None
                if audio_path:
                    slide_xml_path = extract_dir / 'ppt' / 'slides' / f'slide{slide_num}.xml'
                    rels_path = extract_dir / 'ppt' / 'slides' / '_rels' / f'slide{slide_num}.xml.rels'

                    ext = audio_path.suffix.lower()
                    media_name = f'narration{slide_num}{ext}'
                    shutil.copy2(audio_path, media_dir / media_name)
                    audio_exts_used.add(ext)

                    poster_name = 'narration_poster.png'
                    poster_path = media_dir / poster_name
                    if not poster_path.exists():
                        poster_path.write_bytes(AUDIO_MARKER_PNG_BYTES)
                    has_any_image = True
                    image_exts_used.add('png')

                    media_rid = _append_relationship(
                        rels_path,
                        MEDIA_REL_TYPE,
                        f'../media/{media_name}',
                    )
                    audio_rid = _append_relationship(
                        rels_path,
                        AUDIO_REL_TYPE,
                        f'../media/{media_name}',
                    )
                    poster_rid = _append_relationship(
                        rels_path,
                        IMAGE_REL_TYPE,
                        f'../media/{poster_name}',
                    )

                    slide_xml = slide_xml_path.read_text(encoding='utf-8')
                    narration_shape_id = next_shape_id(slide_xml)
                    slide_xml = inject_narration(
                        slide_xml,
                        shape_id=narration_shape_id,
                        shape_name=media_name,
                        audio_rid=audio_rid,
                        media_rid=media_rid,
                        poster_rid=poster_rid,
                    )

                    if use_narration_timings:
                        duration = probe_audio_duration(audio_path)
                        if duration is None:
                            raise RuntimeError(
                                f"Unable to read narration duration with ffprobe: {audio_path}"
                            )
                        slide_xml = apply_recorded_timing(
                            slide_xml,
                            advance_after=duration + narration_padding,
                            transition_duration=slide_transition_duration,
                            transition_effect=slide_transition,
                        )
                        resolved_advance_after = duration + narration_padding
                        resolved_advance_on_click = False
                        package_uses_timings = True
                    slide_xml_path.write_text(slide_xml, encoding='utf-8')
                    narration_slides_created.add(slide_num)

                final_slide_xml = slide_xml_path.read_text(encoding='utf-8')
                try:
                    resolved_motion = validate_generated_transition_xml(
                        final_slide_xml,
                        effect=slide_transition,
                        duration=slide_transition_duration,
                        advance_on_click=resolved_advance_on_click,
                        advance_after=resolved_advance_after,
                    )
                except ValueError as exc:
                    raise RuntimeError(
                        f'Slide {slide_num} transition validation failed: {exc}'
                    ) from exc
                try:
                    resolved_animation = validate_generated_animation_xml(
                        final_slide_xml,
                        expected_animation_targets,
                        duration=expected_animation_duration,
                        trigger=expected_animation_trigger,
                    )
                except ValueError as exc:
                    raise RuntimeError(
                        f'Slide {slide_num} animation validation failed: {exc}'
                    ) from exc

                if conversion_trace is not None:
                    motion_summary = asdict(resolved_motion)
                    for trace_entry in reversed(conversion_trace):
                        if trace_entry.get('slide_num') == slide_num:
                            trace_entry['motion'] = motion_summary
                            trace_entry['animation'] = asdict(resolved_animation)
                            break

                if verbose:
                    if use_native_shapes:
                        mode_str = " (Native)"
                    elif use_compat_mode and not use_native_shapes:
                        mode_str = " (PNG+SVG)" if has_any_image else " (SVG)"
                    else:
                        mode_str = " (SVG)"
                    has_notes = slide_num in notes_slides_created
                    notes_str = " +notes" if has_notes else ""
                    narration_str = " +narration" if slide_num in narration_slides_created else ""
                    print(f"  [{i}/{len(svg_files)}] {svg_path.name}{mode_str}{notes_str}{narration_str}")

                success_count += 1

            except Exception as e:
                if verbose:
                    print(f"  [{i}/{len(svg_files)}] {svg_path.name} - Error: {e}")
                if use_native_shapes:
                    raise

        if (
            use_native_shapes
            and pptx_structure == "baseline"
            and success_count == len(svg_files)
        ):
            _convert_page_number_texts_to_fields(
                extract_dir,
                len(svg_files),
                conversion_trace if conversion_trace is not None else structure_trace,
                context="Baseline",
                verbose=verbose,
            )
            _promote_common_slide_backgrounds_to_masters(
                extract_dir,
                structure,
                len(svg_files),
                verbose=verbose,
            )
            _promote_common_chrome_shapes_to_masters(
                extract_dir,
                structure,
                len(svg_files),
                conversion_trace if conversion_trace is not None else structure_trace,
                verbose=verbose,
            )
            _extract_baseline_layout_families(
                extract_dir,
                structure,
                svg_files,
                verbose=verbose,
            )
            _promote_common_chrome_shapes_to_layouts(
                extract_dir,
                len(svg_files),
                conversion_trace if conversion_trace is not None else structure_trace,
                verbose=verbose,
            )
            _prune_unused_slide_layouts(
                extract_dir,
                structure,
                len(svg_files),
                verbose=verbose,
            )

        if (
            use_native_shapes
            and pptx_structure == "structured"
            and success_count == len(svg_files)
        ):
            _convert_page_number_texts_to_fields(
                extract_dir,
                len(svg_files),
                conversion_trace if conversion_trace is not None else structure_trace,
                context="Structured",
                verbose=verbose,
            )
            if template_specs is None:
                raise TemplateStructureError(
                    "Structured metadata was not parsed before export"
                )
            (
                template_background_expectations,
                template_shape_roster_expectations,
            ) = _apply_explicit_layout_structure(
                extract_dir,
                structure,
                template_specs,
                conversion_trace if conversion_trace is not None else structure_trace,
                active_theme_font_spec,
                verbose=verbose,
            )
            master_count = apply_master_text_style_spec(
                extract_dir,
                master_text_style_spec,
            )
            if verbose:
                print(
                    "  Structured master text styles: "
                    f"{master_count} master(s), "
                    f"title {master_text_style_spec.title_hpt / 100:g}pt, "
                    f"body {master_text_style_spec.body_hpt / 100:g}pt"
                )
            _prune_unused_slide_layouts(
                extract_dir,
                structure,
                len(svg_files),
                verbose=verbose,
            )

        if (
            use_native_shapes
            and pptx_structure == "preserve"
            and success_count == len(svg_files)
        ):
            _convert_page_number_texts_to_fields(
                extract_dir,
                len(svg_files),
                conversion_trace if conversion_trace is not None else structure_trace,
                context="Preserve",
                verbose=verbose,
            )
            if template_specs is None or native_structure_contract is None:
                raise TemplateStructureError(
                    "Preserved structure metadata was not parsed before export"
                )
            _apply_preserved_structure(
                extract_dir,
                template_specs,
                native_structure_contract,
                conversion_trace if conversion_trace is not None else structure_trace,
                verbose=verbose,
            )

        # Update [Content_Types].xml
        content_types_path = extract_dir / '[Content_Types].xml'
        with open(content_types_path, 'r', encoding='utf-8') as f:
            content_types = f.read()

        if not use_native_shapes:
            content_types = _add_default_content_type(content_types, 'svg', 'image/svg+xml')
        for ext in sorted(image_exts_used):
            content_types = _add_default_content_type(
                content_types,
                ext,
                _content_type_for_extension(ext),
            )
        if 'xlsx' in package_exts_used:
            content_types = _add_default_content_type(
                content_types,
                'xlsx',
                'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
            )
        for part_name, content_type in sorted(package_content_overrides.items()):
            content_types = _add_content_type_override(content_types, part_name, content_type)
        with open(content_types_path, 'w', encoding='utf-8') as f:
            f.write(content_types)

        if audio_exts_used:
            for ext in sorted(audio_exts_used):
                content_type = AUDIO_CONTENT_TYPES.get(ext)
                if content_type:
                    content_types = _add_default_content_type(content_types, ext, content_type)
            if 'Extension="png"' not in content_types:
                content_types = _add_default_content_type(content_types, 'png', 'image/png')
            with open(content_types_path, 'w', encoding='utf-8') as f:
                f.write(content_types)

        # Add notes master / slides content types
        if enable_notes and notes_slides_created:
            notes_theme_override = (
                '  <Override PartName="/ppt/theme/theme2.xml" '
                'ContentType="application/vnd.openxmlformats-officedocument.theme+xml"/>'
            )
            if notes_theme_override not in content_types:
                content_types = content_types.replace(
                    '</Types>',
                    notes_theme_override + '\n</Types>',
                )
            notes_master_override = (
                '  <Override PartName="/ppt/notesMasters/notesMaster1.xml" '
                'ContentType="application/vnd.openxmlformats-officedocument.presentationml.notesMaster+xml"/>'
            )
            if notes_master_override not in content_types:
                content_types = content_types.replace(
                    '</Types>',
                    notes_master_override + '\n</Types>',
                )
            for i in sorted(notes_slides_created):
                override = (
                    f'  <Override PartName="/ppt/notesSlides/notesSlide{i}.xml" '
                    f'ContentType="application/vnd.openxmlformats-officedocument.presentationml.notesSlide+xml"/>'
                )
                if override not in content_types:
                    content_types = content_types.replace('</Types>', override + '\n</Types>')
            with open(content_types_path, 'w', encoding='utf-8') as f:
                f.write(content_types)

        if package_uses_timings:
            set_directory_use_timings(extract_dir)

        rels_problems = _verify_internal_rels_targets(extract_dir)
        if rels_problems:
            details = '\n'.join(f'  - {p}' for p in rels_problems)
            raise RuntimeError(
                'PPTX package contains dangling internal relationship targets; '
                'PowerPoint will report the file as corrupt:\n' + details
            )

        # Replace the python-pptx base-template metadata (stale "Steve Canny"
        # author, 2013 dates, "generated using python-pptx", Slides=0) with
        # accurate, tool-neutral document properties.
        pres_format = _presentation_format(width_emu, height_emu)
        _stamp_docprops(extract_dir, len(svg_files), pres_format, doc_metadata)

        # Repackage PPTX to a temporary file first. The public output path is
        # replaced only after every slide and relationship has succeeded.
        temp_output_path = temp_dir / 'result.pptx'
        with zipfile.ZipFile(temp_output_path, 'w', zipfile.ZIP_DEFLATED) as zf:
            for file_path in extract_dir.rglob('*'):
                if file_path.is_file():
                    arcname = file_path.relative_to(extract_dir)
                    zf.write(file_path, arcname)
        if (
            use_native_shapes
            and pptx_structure == "structured"
            and success_count == len(svg_files)
        ):
            if template_specs is None:
                raise TemplateStructureError(
                    "Explicit Layout metadata was not parsed before validation"
                )
            try:
                validate_pptx_template_package(
                    temp_output_path,
                    template_specs,
                    expected_backgrounds=template_background_expectations,
                    expected_shape_rosters=template_shape_roster_expectations,
                )
            except ValueError as exc:
                raise TemplateStructureError(
                    f"PPTX structured package validation failed: {exc}"
                ) from exc
        try:
            validate_pptx_transition_package(
                temp_output_path,
                require_use_timings=package_uses_timings,
            )
        except ValueError as exc:
            raise RuntimeError(
                f'PPTX transition package validation failed: {exc}'
            ) from exc
        try:
            validate_pptx_animation_package(
                temp_output_path,
                require_supported_effects=True,
            )
        except ValueError as exc:
            raise RuntimeError(
                f'PPTX animation package validation failed: {exc}'
            ) from exc
        shutil.move(str(temp_output_path), str(output_path))
        permission_warnings = _relax_output_permissions(output_path)

        if conversion_trace_path and conversion_trace is not None:
            conversion_trace_path.parent.mkdir(parents=True, exist_ok=True)
            payload = {
                'output': str(output_path),
                'slide_count': len(svg_files),
                'slides': conversion_trace,
            }
            conversion_trace_path.write_text(
                json.dumps(payload, ensure_ascii=False, indent=2),
                encoding='utf-8',
            )

        if verbose:
            print()
            print(f"[Done] Saved: {output_path}")
            for warning in permission_warnings:
                print(f"  [warn] {warning}")
            if conversion_trace_path and conversion_trace is not None:
                print(f"  Trace: {conversion_trace_path}")
            print(f"  Succeeded: {success_count}, Failed: {len(svg_files) - success_count}")
            if use_compat_mode and has_any_image:
                print(f"  Mode: Office compatibility mode (supports all Office versions)")
                if PNG_RENDERER == 'svglib' and renderer_hint:
                    print(f"  [Tip] {renderer_hint}")

        return success_count == len(svg_files)

    finally:
        shutil.rmtree(temp_dir, ignore_errors=True)
