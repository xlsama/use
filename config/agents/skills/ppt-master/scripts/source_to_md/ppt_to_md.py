#!/usr/bin/env python3
"""
PowerPoint to Markdown Converter

Extracts slide text, tables, SmartArt node structure, speaker notes, and
embedded pictures from Open XML PowerPoint files into Markdown.

Primary use case: PPTX source decks -> Markdown for PPT generation input.

Hyperlinks present in the source deck are preserved: run-level external URLs
and slide-internal jumps are emitted as ``[text](url)`` / ``[text](#slide-N)``,
with a shape-level ``click_action`` fallback.

Dependency:
    pip install python-pptx

API stability note:
    Detecting slide-internal jumps (``ppaction://hlinksldjump``) reads
    ``run._r`` (the CT_TextRun lxml element) because python-pptx exposes no
    public API to distinguish an internal jump from an external URL. XY chart
    extraction likewise reads ``series._element`` for X values and bubble sizes,
    which the public chart API does not expose. Keep these private accesses
    localized here and covered by conversion smoke tests.
"""

from __future__ import annotations

import argparse
import hashlib
import json
import re
import shutil
import sys
import zipfile
from dataclasses import dataclass
from io import BytesIO
from pathlib import Path
from urllib.parse import quote
from xml.etree import ElementTree as ET

_SCRIPTS_DIR = Path(__file__).resolve().parents[1]
if str(_SCRIPTS_DIR) not in sys.path:
    sys.path.insert(0, str(_SCRIPTS_DIR))

from console_encoding import configure_utf8_stdio  # noqa: E402
from _batch import run_path_batch  # noqa: E402
from _conversion_profile import write_conversion_profile_best_effort  # noqa: E402
from template_fill_pptx.diagram_read import (  # noqa: E402
    read_smartart_diagrams,
    smartart_to_markdown,
)

from pptx import Presentation
from pptx.enum.action import PP_ACTION
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

configure_utf8_stdio()


EMU_PER_INCH = 914400
DRAWINGML_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
PRESENTATIONML_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
RELATIONSHIP_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
CHARTEX_URI = "http://schemas.microsoft.com/office/drawing/2014/chartex"
OFFICE_VECTOR_EXTENSIONS = {"emf", "wmf"}
IMAGE_EXT_BY_CONTENT_TYPE = {
    "image/bmp": "bmp",
    "image/gif": "gif",
    "image/jpeg": "jpg",
    "image/jpg": "jpg",
    "image/png": "png",
    "image/svg+xml": "svg",
    "image/tiff": "tiff",
    "image/x-emf": "emf",
    "image/x-wmf": "wmf",
}
LEGACY_GENERATED_IMAGE_RE = re.compile(r"^slide_\d{2}_image_\d{2}\.[A-Za-z0-9]+$")

# Hyperlink schemes dropped during extraction (a blacklist of known-dangerous
# schemes). PowerPoint also rejects unrecognized schemes at open time, so the
# residual risk from schemes not listed here is low.
UNSUPPORTED_URL_SCHEMES = ("javascript:", "data:", "vbscript:", "file:")


SUPPORTED_FORMATS = {
    ".pptx": "PowerPoint Presentation",
    ".pptm": "Macro-enabled PowerPoint Presentation",
    ".ppsx": "PowerPoint Slide Show",
    ".ppsm": "Macro-enabled PowerPoint Slide Show",
    ".potx": "PowerPoint Template",
    ".potm": "Macro-enabled PowerPoint Template",
}


@dataclass
class LeafShape:
    """Flattened leaf shape with stable position ordering."""

    shape: object
    top: int
    left: int


@dataclass
class SavedPicture:
    """Extracted image asset plus manifest metadata."""

    filename: str
    manifest_entry: dict[str, object]
    is_new_asset: bool


def normalize_text(value: str) -> str:
    """Collapse whitespace while preserving paragraph boundaries elsewhere."""
    value = value.replace("\r\n", "\n").replace("\r", "\n")
    lines = [re.sub(r"\s+", " ", line).strip() for line in value.split("\n")]
    lines = [line for line in lines if line]
    return "\n".join(lines)


def normalize_ext(ext: str | None, content_type: str | None = None) -> str:
    """Return a lowercase extension without a leading dot."""
    if ext:
        ext = ext.lower().lstrip(".")
        if ext == "jpeg":
            return "jpg"
        return ext
    if content_type:
        return IMAGE_EXT_BY_CONTENT_TYPE.get(content_type.lower(), "bin")
    return "bin"


def sanitize_filename(value: str) -> str:
    """Return a filesystem-safe basename."""
    value = re.sub(r"[^\w.\-]+", "_", value, flags=re.UNICODE)
    return value.strip("._") or "asset"


def escape_table_cell(value: str) -> str:
    """Escape Markdown table syntax inside a cell."""
    normalized = value.replace("\r\n", "\n").replace("\r", "\n")
    lines = [re.sub(r"\s+", " ", line).strip() for line in normalized.split("\n")]
    with_breaks = "<br>".join(lines)
    return with_breaks.replace("|", r"\|") or " "


def _safe_position(shape: object, attr: str) -> int:
    """Read a shape's ``top`` / ``left`` EMU, tolerating broken inheritance.

    A placeholder with no explicit position resolves it by walking up to its
    master. A deck that ships notesSlides without a notesMaster (or any other
    partial inheritance chain) makes python-pptx raise on that lookup, so treat
    an unresolvable position as 0 rather than aborting the whole conversion.
    """
    try:
        return int(getattr(shape, attr, 0) or 0)
    except Exception:
        return 0


def iter_leaf_shapes(shapes: object) -> list[LeafShape]:
    """Return a flattened, reading-order list of shapes."""
    items: list[LeafShape] = []
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            items.extend(iter_leaf_shapes(shape.shapes))
            continue
        items.append(
            LeafShape(
                shape=shape,
                top=_safe_position(shape, "top"),
                left=_safe_position(shape, "left"),
            )
        )
    items.sort(key=lambda item: (item.top, item.left))
    return items


def _is_supported_url(url: str) -> bool:
    """Reject empty URLs and the known-dangerous schemes."""
    return bool(url) and not any(
        url.lower().startswith(scheme) for scheme in UNSUPPORTED_URL_SCHEMES
    )


def _escape_md_link_text(text: str) -> str:
    """Backslash-escape characters that would break a Markdown link label.

    A stray ``]`` in anchor text would otherwise close the ``[...]`` early.
    """
    return text.replace("\\", "\\\\").replace("[", "\\[").replace("]", "\\]")


def _encode_md_url(url: str) -> str:
    """Percent-encode a URL so Markdown link syntax stays unambiguous.

    Notably encodes ``(`` / ``)`` to ``%28`` / ``%29`` so a parenthesised URL
    does not terminate the ``[text](url)`` form early.
    """
    return quote(url, safe="/:?=&%#@!$'*+,;")


def _resolve_internal_jump(run: object, shape: object) -> str | None:
    """Return ``#slide-N`` for a run carrying a slide-internal jump, else None.

    Reads ``run._r`` (private python-pptx API) because the public
    ``run.hyperlink.address`` cannot tell an internal jump apart from an
    external URL — see the module docstring's API stability note.
    """
    r_id = ""
    try:
        rpr = run._r.find(qn("a:rPr"))
        if rpr is None:
            return None
        hlink = rpr.find(qn("a:hlinkClick"))
        if hlink is None or "hlinksldjump" not in (hlink.get("action", "") or ""):
            return None
        r_id = hlink.get(qn("r:id"), "")
        if not r_id:
            return None
        target_slide = shape.part.related_part(r_id).slide
        prs = shape.part.slide.part.package.presentation_part.presentation
        return f"#slide-{list(prs.slides).index(target_slide) + 1}"
    except (KeyError, ValueError, AttributeError):
        print(f"[WARN] ppt_to_md: could not resolve slide jump rId={r_id}", file=sys.stderr)
        return None


def _run_url(run: object, shape: object) -> str | None:
    """Resolve a run's hyperlink target to a markdown-ready URL, or None."""
    if shape is not None:
        internal = _resolve_internal_jump(run, shape)
        if internal:
            return internal
    try:
        addr = run.hyperlink.address
    except AttributeError:
        return None
    if _is_supported_url(addr or ""):
        return _encode_md_url(addr)
    return None


def _paragraph_to_markdown(paragraph: object, shape: object) -> str:
    """Render one paragraph, merging consecutive runs that share a URL.

    Run text is concatenated verbatim — including the spaces between runs — and
    normalized only once over the assembled paragraph, so a link in the middle
    of a sentence does not swallow its surrounding spaces. A link group's own
    leading / trailing whitespace is kept outside the ``[...]`` so it separates
    words rather than padding the anchor text.
    """
    parts = []
    current_url = None
    current_text = ""

    def flush():
        if not current_text:
            return
        if current_url is None:
            parts.append(current_text)
            return
        lead = current_text[: len(current_text) - len(current_text.lstrip())]
        trail = current_text[len(current_text.rstrip()):]
        core = current_text.strip()
        display = _escape_md_link_text(core) if core else current_url
        parts.append(f"{lead}[{display}]({current_url}){trail}")

    has_run_hyperlink = False
    for run in paragraph.runs:
        url = _run_url(run, shape)
        if url:
            has_run_hyperlink = True
        if url != current_url:
            flush()
            current_text = ""
            current_url = url
        current_text += run.text or ""
    flush()

    text = normalize_text("".join(parts))

    # Shape-level click_action only matters when no run carried its own link.
    if not has_run_hyperlink and shape is not None:
        text = _apply_shape_click_action(text, shape)
    return text


def _apply_shape_click_action(text: str, shape: object) -> str:
    """Wrap paragraph text in a link from the shape's click_action, if any."""
    try:
        action = shape.click_action
        if action.action == PP_ACTION.HYPERLINK:
            url = action.hyperlink.address or ""
            if _is_supported_url(url):
                return f"[{_escape_md_link_text(text)}]({_encode_md_url(url)})"
        elif action.action == PP_ACTION.NAMED_SLIDE:
            target = action.target_slide
            if target is not None:
                prs = shape.part.slide.part.package.presentation_part.presentation
                idx = list(prs.slides).index(target) + 1
                return f"[{_escape_md_link_text(text)}](#slide-{idx})"
    except (AttributeError, ValueError):
        print("[WARN] ppt_to_md: could not process shape click_action", file=sys.stderr)
    return text


def _paragraph_has_hyperlink(paragraph: object) -> bool:
    """True if any run carries an external URL or an internal slide jump."""
    for run in paragraph.runs:
        try:
            if run.hyperlink.address:
                return True
        except AttributeError:
            pass
        try:
            rpr = run._r.find(qn("a:rPr"))
            if rpr is not None and rpr.find(qn("a:hlinkClick")) is not None:
                return True
        except AttributeError:
            continue
    return False


def text_frame_to_markdown(text_frame: object, shape: object = None) -> str:
    """Convert a PowerPoint text frame into Markdown, preserving hyperlinks.

    Run-level external URLs and slide-internal jumps are emitted as
    ``[text](url)`` / ``[text](#slide-N)``; consecutive runs sharing a URL are
    merged. When no run carries a link, the shape's ``click_action`` is used as
    a paragraph-level fallback. Pass ``shape`` to enable hyperlink extraction;
    without it the frame degrades to plain text.
    """
    visible_paragraphs = [
        paragraph for paragraph in text_frame.paragraphs
        if normalize_text(paragraph.text) or _paragraph_has_hyperlink(paragraph)
    ]
    if not visible_paragraphs:
        return ""

    list_like = any(paragraph.level > 0 for paragraph in visible_paragraphs)
    if not list_like:
        list_like = len(visible_paragraphs) > 1

    paragraphs = []
    for paragraph in visible_paragraphs:
        text = _paragraph_to_markdown(paragraph, shape)
        if not text:
            continue
        if list_like:
            indent = "  " * max(paragraph.level, 0)
            paragraphs.append(f"{indent}- {text}")
        else:
            paragraphs.append(text)

    if list_like:
        return "\n".join(paragraphs)
    return "\n\n".join(paragraphs)


def table_to_markdown(table: object) -> str:
    """Convert a PowerPoint table to a Markdown table."""
    rows = []
    for row in table.rows:
        cells = [escape_table_cell(cell.text) for cell in row.cells]
        rows.append(cells)

    if not rows:
        return ""

    column_count = max(len(row) for row in rows)
    normalized_rows = [row + [" "] * (column_count - len(row)) for row in rows]
    header = normalized_rows[0]
    separator = ["---"] * column_count
    body = normalized_rows[1:]

    lines = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join(separator) + " |",
    ]
    for row in body:
        lines.append("| " + " | ".join(row) + " |")
    return "\n".join(lines)


def _format_chart_value(value: object) -> str:
    """Render a chart data point, trimming whole-number floats."""
    if value is None:
        return ""
    if isinstance(value, float) and value.is_integer():
        return str(int(value))
    return str(value)


def _chart_header(chart: object, name: str) -> tuple[str, str]:
    """Return the Markdown chart header and its best-effort type label."""
    try:
        chart_type = str(chart.chart_type)
    except (ValueError, AttributeError, KeyError):
        chart_type = ""
    raw_name = "" if name is None else str(name)
    chart_name = normalize_text(raw_name).replace("\n", " ") or "Chart"
    header = f"> [Chart] {chart_name}" + (f" — {chart_type}" if chart_type else "")
    return header, chart_type


def _chart_warning_lines(warnings: list[str]) -> list[str]:
    """Return stable, de-duplicated Markdown warning blocks."""
    lines: list[str] = []
    seen: set[str] = set()
    for warning in warnings:
        normalized = normalize_text(warning).replace("\n", " ") or "unknown warning"
        if normalized in seen:
            continue
        seen.add(normalized)
        lines.append(f"> [Chart data warning: {normalized}]")
    return lines


def _chart_data_unavailable(
    header: str,
    reason: str,
    *,
    warnings: list[str] | None = None,
) -> str:
    """Attach an explicit data-read failure to a chart heading."""
    normalized_reason = normalize_text(reason).replace("\n", " ") or "unknown reason"
    lines = [header]
    lines.extend(_chart_warning_lines(warnings or []))
    lines.append(f"> [Chart data unavailable: {normalized_reason}]")
    return "\n".join(lines)


def _chart_value_cell(value: object) -> str:
    """Return one chart table cell while keeping missing values visibly empty."""
    rendered = _format_chart_value(value)
    return escape_table_cell(rendered) if rendered else ""


def _chart_series_element(series: object) -> object | None:
    """Return python-pptx's series XML carrier when its public API is insufficient."""
    element = getattr(series, "_element", None)
    if element is not None:
        return element
    return getattr(series, "_ser", None)


def _chart_series_name(series: object, index: int) -> str:
    """Return a stable, Markdown-safe series name."""
    try:
        raw_name = series.name
    except (ValueError, AttributeError, KeyError):
        raw_name = None
    label = str(raw_name) if raw_name not in (None, "") else f"Series {index}"
    return escape_table_cell(label)


def _chart_numeric_cache_values(
    parent: object | None,
) -> tuple[list[object | None] | None, str | None]:
    """Read one XY display cache through python-pptx's OOXML value helpers."""
    if parent is None:
        return None, "missing numeric value container"
    try:
        point_count = int(parent.ptCount_val)
        values = [parent.pt_v(index) for index in range(point_count)]
    except (AttributeError, IndexError, TypeError, ValueError):
        return None, "invalid or unavailable numeric display cache"
    if point_count <= 0 or all(value is None for value in values):
        return None, "numeric display cache contains no values"
    return values, None


def _chart_family(chart_type: str, series: list[object]) -> str:
    """Classify category, scatter, and bubble charts without misreading XY as category."""
    type_key = chart_type.upper()
    if "BUBBLE" in type_key:
        return "bubble"
    if "SCATTER" in type_key:
        return "scatter"
    for item in series:
        element = _chart_series_element(item)
        if element is None:
            continue
        if element.find(qn("c:bubbleSize")) is not None:
            return "bubble"
        if element.find(qn("c:xVal")) is not None or element.find(qn("c:yVal")) is not None:
            return "scatter"
    return "category"


def _xy_chart_to_markdown(
    series: list[object],
    *,
    family: str,
    header: str,
) -> str:
    """Render scatter/bubble series as typed per-point X/Y[/size] rows."""
    table_header = ["Series", "Point", "X", "Y"]
    if family == "bubble":
        table_header.append("Size")
    rows: list[list[str]] = []
    warnings: list[str] = []

    for series_index, item in enumerate(series, start=1):
        series_name = _chart_series_name(item, series_index)
        element = _chart_series_element(item)
        if element is None:
            x_values = None
            warnings.append(f"{series_name}: series XML is unavailable for X data")
        else:
            x_values, x_error = _chart_numeric_cache_values(
                element.find(qn("c:xVal"))
            )
            if x_error:
                warnings.append(f"{series_name}: X data {x_error}")
        try:
            y_values = list(item.values)
        except (ValueError, TypeError, AttributeError, KeyError):
            y_values = []
            warnings.append(f"{series_name}: Y values are unavailable")

        size_values: list[object | None] | None = None
        if family == "bubble":
            if element is None:
                size_error = "missing series XML"
            else:
                size_values, size_error = _chart_numeric_cache_values(
                    element.find(qn("c:bubbleSize"))
                )
            if size_error:
                warnings.append(f"{series_name}: bubble sizes {size_error}")

        x_values = x_values or []
        point_count = max(
            len(x_values),
            len(y_values),
            len(size_values or []),
        )
        if point_count == 0:
            warnings.append(f"{series_name}: no readable points")
            continue

        point_counts = {len(x_values), len(y_values)}
        if family == "bubble":
            point_counts.add(len(size_values or []))
        if len(point_counts) > 1:
            dimensions = "X/Y/size" if family == "bubble" else "X/Y"
            warnings.append(
                f"{series_name}: {dimensions} point counts differ; "
                "missing cells are blank"
            )

        for point_index in range(point_count):
            x_value = x_values[point_index] if point_index < len(x_values) else None
            y_value = y_values[point_index] if point_index < len(y_values) else None
            row = [
                series_name,
                str(point_index + 1),
                _chart_value_cell(x_value),
                _chart_value_cell(y_value),
            ]
            if family == "bubble":
                size_value = (
                    size_values[point_index]
                    if size_values is not None and point_index < len(size_values)
                    else None
                )
                row.append(_chart_value_cell(size_value))
            rows.append(row)

    if not rows:
        return _chart_data_unavailable(
            header,
            "chart has no readable XY points",
            warnings=warnings,
        )
    lines = [header]
    lines.extend(_chart_warning_lines(warnings))
    lines.extend([
        "",
        "| " + " | ".join(table_header) + " |",
        "| " + " | ".join(["---"] * len(table_header)) + " |",
    ])
    lines.extend("| " + " | ".join(row) + " |" for row in rows)
    return "\n".join(lines)


def _category_chart_to_markdown(chart: object, series: list[object], header: str) -> str:
    """Render a conventional category chart through python-pptx's public API."""
    categories: list[str] = []
    warnings: list[str] = []
    has_category_xml = any(
        (element := _chart_series_element(item)) is not None
        and element.find(qn("c:cat")) is not None
        for item in series
    )
    try:
        plots = list(chart.plots)
        if plots:
            categories = [
                escape_table_cell(str(category)) if category is not None else ""
                for category in plots[0].categories
            ]
    except (ValueError, TypeError, IndexError, AttributeError, KeyError):
        if has_category_xml:
            warnings.append("chart categories are unavailable; using point numbers")
    if not has_category_xml and not categories:
        warnings.append("chart categories are missing; using point numbers")
    elif has_category_xml and not categories:
        warnings.append("chart categories are empty; using point numbers")

    series_data: list[tuple[str, list[object]]] = []
    for index, item in enumerate(series, start=1):
        series_name = _chart_series_name(item, index)
        try:
            values = list(item.values)
        except (ValueError, TypeError, AttributeError, KeyError):
            warnings.append(f"{series_name}: series values are unavailable")
            continue
        series_data.append((series_name, values))

    row_count = max(
        len(categories),
        max((len(values) for _, values in series_data), default=0),
    )
    if not series_data or row_count == 0:
        return _chart_data_unavailable(
            header,
            "chart has no readable category-series data",
            warnings=warnings,
        )
    point_counts = {len(values) for _, values in series_data}
    if categories:
        point_counts.add(len(categories))
    if len(point_counts) > 1:
        warnings.append("category/series point counts differ; missing cells are blank")

    table_header = (["Category"] if categories else ["#"]) + [
        series_name for series_name, _ in series_data
    ]
    lines = [header]
    lines.extend(_chart_warning_lines(warnings))
    lines.extend([
        "",
        "| " + " | ".join(table_header) + " |",
        "| " + " | ".join(["---"] * len(table_header)) + " |",
    ])
    for row_index in range(row_count):
        if categories:
            label = categories[row_index] if row_index < len(categories) else ""
        else:
            label = str(row_index + 1)
        cells = [label]
        for _, values in series_data:
            value = values[row_index] if row_index < len(values) else None
            cells.append(_chart_value_cell(value))
        lines.append("| " + " | ".join(cells) + " |")
    return "\n".join(lines)


def chart_to_markdown(chart: object, name: str) -> str:
    """Render category, scatter, and bubble data without flattening chart semantics.

    A native PowerPoint chart stores its data in embedded XML, not in any text
    frame. Public python-pptx APIs cover category values and scatter/bubble Y
    values, but not XY X coordinates or bubble sizes. Read only those missing
    display caches from the series XML. Preserve every readable value and emit
    explicit warnings for missing dimensions or series rather than discarding
    the chart's remaining content.
    """
    header, chart_type = _chart_header(chart, name)
    try:
        series = list(chart.series)
    except (ValueError, TypeError, AttributeError, KeyError):
        return _chart_data_unavailable(header, "chart series are unavailable")
    if not series:
        return _chart_data_unavailable(header, "chart has no readable series")

    family = _chart_family(chart_type, series)
    if family in {"scatter", "bubble"}:
        return _xy_chart_to_markdown(series, family=family, header=header)
    return _category_chart_to_markdown(chart, series, header)


def _chart_reference_id(element: object) -> str | None:
    """Return the first chart relationship id carried by an OOXML shape subtree."""
    for descendant in element.iter():
        if descendant.tag.rsplit("}", 1)[-1] != "chart":
            continue
        relationship_id = descendant.get(f"{{{RELATIONSHIP_NS}}}id")
        if relationship_id:
            return relationship_id
    return None


def _unexposed_chartex_markdown(
    slide: object,
    emitted_relationship_ids: set[str],
) -> list[str]:
    """Report ChartEx objects omitted from ``slide.shapes`` by python-pptx."""
    blocks: list[str] = []
    seen_relationship_ids: set[str] = set()
    slide_element = getattr(slide, "element", None)
    if slide_element is None:
        slide_element = getattr(slide, "_element", None)
    if slide_element is None:
        return blocks
    for graphic_data in slide_element.iter(f"{{{DRAWINGML_NS}}}graphicData"):
        if graphic_data.get("uri") != CHARTEX_URI:
            continue
        relationship_id = _chart_reference_id(graphic_data)
        if relationship_id and (
            relationship_id in emitted_relationship_ids
            or relationship_id in seen_relationship_ids
        ):
            continue
        if relationship_id:
            seen_relationship_ids.add(relationship_id)

        name = "ChartEx chart"
        current = graphic_data
        while current is not None:
            name_element = current.find(f".//{{{PRESENTATIONML_NS}}}cNvPr")
            if name_element is not None and name_element.get("name"):
                name = name_element.get("name")
                break
            current = current.getparent() if hasattr(current, "getparent") else None
        chart_name = normalize_text(str(name)).replace("\n", " ") or "ChartEx chart"
        header = f"> [Chart] {chart_name} — ChartEx"
        blocks.append(_chart_data_unavailable(header, "unsupported ChartEx data model"))
    return blocks


def _image_part_for_shape(shape: object) -> object | None:
    """Return the first embedded image part referenced by a shape."""
    element = getattr(shape, "element", None)
    part = getattr(shape, "part", None)
    if element is None or part is None:
        return None

    try:
        blips = element.xpath(".//a:blip")
    except Exception:
        return None

    for blip in blips:
        rel_id = blip.get(qn("r:embed")) or blip.get(qn("r:link"))
        if not rel_id:
            continue
        try:
            return part.related_part(rel_id)
        except Exception:
            continue
    return None


def _image_size_from_bytes(blob: bytes) -> tuple[int | None, int | None]:
    """Return bitmap dimensions when Pillow can decode the bytes."""
    try:
        from PIL import Image
    except ImportError:
        return None, None
    try:
        with Image.open(BytesIO(blob)) as img:
            return img.width, img.height
    except (OSError, ValueError):
        return None, None


def _shape_emu(shape: object, attr: str) -> int:
    value = getattr(shape, attr, 0) or 0
    return int(value)


def _shape_occurrence(
    shape: object,
    slide_index: int,
) -> dict[str, object]:
    """Return slide-specific image placement metadata."""
    display_width_emu = _shape_emu(shape, "width")
    display_height_emu = _shape_emu(shape, "height")
    display_ratio = (
        display_width_emu / display_height_emu
        if display_width_emu > 0 and display_height_emu > 0
        else None
    )
    return {
        "slide_index": slide_index,
        "shape_name": str(getattr(shape, "name", "")),
        "display_left_emu": _shape_emu(shape, "left"),
        "display_top_emu": _shape_emu(shape, "top"),
        "display_width_emu": display_width_emu,
        "display_height_emu": display_height_emu,
        "display_width_in": round(display_width_emu / EMU_PER_INCH, 4) if display_width_emu else None,
        "display_height_in": round(display_height_emu / EMU_PER_INCH, 4) if display_height_emu else None,
        "display_ratio": round(display_ratio, 6) if display_ratio else None,
    }


def _update_manifest_usage(entry: dict[str, object]) -> None:
    """Refresh aggregate fields after adding an occurrence."""
    occurrences = entry.get("occurrences")
    if not isinstance(occurrences, list):
        occurrences = []
    entry["usage_count"] = len(occurrences)
    ratios = sorted({
        occurrence.get("display_ratio")
        for occurrence in occurrences
        if isinstance(occurrence, dict)
        and isinstance(occurrence.get("display_ratio"), (int, float))
    })
    if ratios:
        entry["display_ratio_variants"] = ratios
        if entry.get("display_ratio") is None:
            entry["display_ratio"] = ratios[0]


def _manifest_entry(
    *,
    index: int,
    filename: str,
    image_part: object,
    ext: str,
    blob: bytes,
    occurrence: dict[str, object],
) -> dict[str, object]:
    """Build image_manifest.json metadata for one unique PowerPoint media part."""
    pixel_width, pixel_height = _image_size_from_bytes(blob)
    pixel_ratio = (
        pixel_width / pixel_height
        if pixel_width and pixel_height
        else None
    )
    is_office_vector = ext in OFFICE_VECTOR_EXTENSIONS
    partname = str(getattr(image_part, "partname", ""))
    content_type = str(getattr(image_part, "content_type", ""))

    entry: dict[str, object] = {
        "index": index,
        "filename": filename,
        "original_filename": filename,
        "asset_kind": "office_vector" if is_office_vector else "bitmap",
        "svg_renderable": not is_office_vector,
        "pptx_native_supported": True,
        "source_kind": "pptx_picture",
        "source_ext": f".{ext}",
        "source_target": partname.lstrip("/"),
        "content_type": content_type,
        "display_left_emu": occurrence.get("display_left_emu"),
        "display_top_emu": occurrence.get("display_top_emu"),
        "display_width_emu": occurrence.get("display_width_emu"),
        "display_height_emu": occurrence.get("display_height_emu"),
        "display_width_in": occurrence.get("display_width_in"),
        "display_height_in": occurrence.get("display_height_in"),
        "display_ratio": occurrence.get("display_ratio"),
        "pixel_width": pixel_width,
        "pixel_height": pixel_height,
        "pixel_ratio": round(pixel_ratio, 6) if pixel_ratio else None,
        "occurrences": [occurrence],
    }
    if entry["display_ratio"] is None and pixel_ratio:
        entry["display_ratio"] = round(pixel_ratio, 6)
    _update_manifest_usage(entry)
    return entry


def _asset_cache_key(image_part: object, blob: bytes) -> str:
    """Return a stable key for deduplicating repeated PPTX media references."""
    partname = str(getattr(image_part, "partname", ""))
    if partname:
        return partname
    return hashlib.sha256(blob).hexdigest()


def _asset_filename(
    image_part: object,
    ext: str,
    asset_index: int,
    used_filenames: set[str],
) -> str:
    """Return a unique asset filename, preferring the PPTX media basename."""
    partname = str(getattr(image_part, "partname", ""))
    base = sanitize_filename(Path(partname).name) if partname else f"image_{asset_index:03d}.{ext}"
    if "." not in base:
        base = f"{base}.{ext}"
    if base not in used_filenames:
        used_filenames.add(base)
        return base

    path = Path(base)
    stem = path.stem
    suffix = path.suffix or f".{ext}"
    counter = 2
    while True:
        candidate = f"{stem}_{counter}{suffix}"
        if candidate not in used_filenames:
            used_filenames.add(candidate)
            return candidate
        counter += 1


def save_picture(
    shape: object,
    asset_dir: Path,
    slide_index: int,
    asset_index: int,
    asset_cache: dict[str, SavedPicture],
    used_filenames: set[str],
) -> SavedPicture | None:
    """Persist a shape image to the output asset directory."""
    image_part = _image_part_for_shape(shape)
    if image_part is None:
        return None

    content_type = getattr(image_part, "content_type", None)
    part_ext = getattr(getattr(image_part, "partname", None), "ext", None)
    ext = normalize_ext(part_ext, content_type)
    blob = bytes(getattr(image_part, "blob", b""))
    if not blob:
        return None

    occurrence = _shape_occurrence(shape, slide_index)
    cache_key = _asset_cache_key(image_part, blob)
    cached = asset_cache.get(cache_key)
    if cached is not None:
        occurrences = cached.manifest_entry.setdefault("occurrences", [])
        if isinstance(occurrences, list):
            occurrences.append(occurrence)
        _update_manifest_usage(cached.manifest_entry)
        return SavedPicture(
            filename=cached.filename,
            manifest_entry=cached.manifest_entry,
            is_new_asset=False,
        )

    filename = _asset_filename(image_part, ext, asset_index, used_filenames)
    output_path = asset_dir / filename
    output_path.write_bytes(blob)
    saved = SavedPicture(
        filename=filename,
        manifest_entry=_manifest_entry(
            index=asset_index,
            filename=filename,
            image_part=image_part,
            ext=ext,
            blob=blob,
            occurrence=occurrence,
        ),
        is_new_asset=True,
    )
    asset_cache[cache_key] = saved
    return saved


def _reset_generated_asset_dir(asset_dir: Path) -> None:
    """Remove a previously generated asset directory."""
    if not asset_dir.exists():
        return
    if not (asset_dir / "image_manifest.json").is_file():
        for path in asset_dir.iterdir():
            if path.is_file() and LEGACY_GENERATED_IMAGE_RE.match(path.name):
                path.unlink()
        return
    shutil.rmtree(asset_dir)


def extract_notes(slide: object) -> str:
    """Extract speaker notes text from a slide, if available."""
    try:
        notes_slide = slide.notes_slide
    except Exception:
        return ""

    blocks = []
    for item in iter_leaf_shapes(notes_slide.shapes):
        shape = item.shape
        if not getattr(shape, "has_text_frame", False):
            continue
        text = text_frame_to_markdown(shape.text_frame, shape)
        if text:
            blocks.append(text)

    return "\n\n".join(blocks).strip()


def convert_presentation_to_markdown(
    input_path: str,
    output_path: str | None = None,
) -> str:
    """Convert a supported PowerPoint file to Markdown."""
    input_file = Path(input_path)
    if not input_file.exists():
        print(f"[ERROR] File not found: {input_path}")
        return ""

    suffix = input_file.suffix.lower()
    if suffix not in SUPPORTED_FORMATS:
        supported = ", ".join(sorted(SUPPORTED_FORMATS.keys()))
        print(f"[ERROR] Unsupported format: {suffix}")
        print(f"   Supported: {supported}")
        print("   Legacy .ppt files should be resaved as .pptx or exported to PDF first.")
        return ""

    print(f"[INFO] Converting {SUPPORTED_FORMATS[suffix]}: {input_file.name}")

    if output_path:
        out_file = Path(output_path)
    else:
        out_file = input_file.with_suffix(".md")

    out_file.parent.mkdir(parents=True, exist_ok=True)
    asset_dir = out_file.parent / f"{out_file.stem}_files"
    _reset_generated_asset_dir(asset_dir)

    presentation = Presentation(str(input_file))
    conversion_warnings: list[str] = []
    diagrams_by_slide: dict[int, list[dict[str, object]]] = {}
    diagram_scan_failures: dict[int, str] = {}
    try:
        with zipfile.ZipFile(input_file) as package:
            for slide_index, slide in enumerate(presentation.slides, 1):
                slide_part = str(slide.part.partname).lstrip("/")
                try:
                    diagrams = read_smartart_diagrams(package, slide_part, slide_index)
                except (OSError, RuntimeError, zipfile.BadZipFile, ET.ParseError) as exc:
                    diagrams_by_slide[slide_index] = []
                    diagram_scan_failures[slide_index] = str(exc)
                    conversion_warnings.append(
                        f"Slide {slide_index}: SmartArt scan failed: {exc}"
                    )
                    continue
                diagrams_by_slide[slide_index] = diagrams
                for diagram in diagrams:
                    issues = [str(item) for item in diagram.get("warnings", []) if item]
                    if diagram.get("status") != "ok":
                        issues.insert(0, f"status={diagram.get('status')}")
                    if not issues:
                        continue
                    conversion_warnings.append(
                        f"Slide {slide_index}, {diagram.get('shape_name') or diagram.get('diagram_id')}: "
                        f"SmartArt content {'; '.join(issues)}"
                    )
    except (OSError, RuntimeError, zipfile.BadZipFile, ET.ParseError) as exc:
        conversion_warnings.append(f"SmartArt package scan failed: {exc}")
        for slide_index in range(1, len(presentation.slides) + 1):
            diagrams_by_slide.setdefault(slide_index, [])
            diagram_scan_failures.setdefault(slide_index, str(exc))

    lines = [
        f"# {input_file.stem}",
        "",
        f"- Source: `{input_file.name}`",
        f"- Total slides: {len(presentation.slides)}",
        "",
    ]

    image_count = 0
    image_ref_count = 0
    asset_dir_used = False
    image_manifest: list[dict[str, object]] = []
    asset_cache: dict[str, SavedPicture] = {}
    used_filenames: set[str] = set()

    for slide_index, slide in enumerate(presentation.slides, 1):
        lines.append(f"## Slide {slide_index}")
        lines.append("")

        blocks = []
        slide_diagrams = diagrams_by_slide.get(slide_index, [])
        diagrams_by_shape_id = {
            str(diagram.get("shape_id")): diagram
            for diagram in slide_diagrams
            if diagram.get("shape_id") is not None
        }
        emitted_diagram_ids: set[str] = set()
        emitted_chart_relationship_ids: set[str] = set()
        for item in iter_leaf_shapes(slide.shapes):
            shape = item.shape

            if getattr(shape, "has_table", False):
                table_md = table_to_markdown(shape.table)
                if table_md:
                    blocks.append(table_md)
                continue

            shape_id = str(getattr(shape, "shape_id", ""))
            diagram = diagrams_by_shape_id.get(shape_id)
            if diagram is not None:
                blocks.append(smartart_to_markdown(diagram))
                emitted_diagram_ids.add(str(diagram.get("diagram_id")))
                continue

            is_picture_shape = shape.shape_type in {
                MSO_SHAPE_TYPE.PICTURE,
                MSO_SHAPE_TYPE.LINKED_PICTURE,
            }
            has_shape_image = is_picture_shape or _image_part_for_shape(shape) is not None
            if has_shape_image:
                image_ref_count += 1
                next_image_index = image_count + 1
                asset_dir.mkdir(parents=True, exist_ok=True)
                saved_picture = save_picture(
                    shape,
                    asset_dir,
                    slide_index,
                    next_image_index,
                    asset_cache,
                    used_filenames,
                )
                if saved_picture is None:
                    if is_picture_shape:
                        blocks.append(f"> [Image] {getattr(shape, 'name', 'Picture')}")
                        continue
                else:
                    if saved_picture.is_new_asset:
                        image_count = next_image_index
                        image_manifest.append(saved_picture.manifest_entry)
                    asset_dir_used = True
                    blocks.append(
                        f"![Slide {slide_index} Image {image_ref_count}]"
                        f"({asset_dir.name}/{saved_picture.filename})"
                    )
                    if is_picture_shape:
                        continue

            if getattr(shape, "has_text_frame", False):
                text_md = text_frame_to_markdown(shape.text_frame, shape)
                if text_md:
                    blocks.append(text_md)
                    continue

            if getattr(shape, "has_chart", False):
                shape_element = getattr(shape, "element", None)
                if shape_element is None:
                    shape_element = getattr(shape, "_element", None)
                relationship_id = (
                    _chart_reference_id(shape_element)
                    if shape_element is not None
                    else None
                )
                if relationship_id:
                    emitted_chart_relationship_ids.add(relationship_id)
                try:
                    blocks.append(chart_to_markdown(shape.chart, getattr(shape, "name", "Chart")))
                except (ValueError, TypeError, AttributeError, KeyError) as exc:
                    raw_name = getattr(shape, "name", "Chart")
                    name = normalize_text("" if raw_name is None else str(raw_name)) or "Chart"
                    blocks.append(
                        _chart_data_unavailable(
                            f"> [Chart] {name}",
                            f"chart read failed ({type(exc).__name__})",
                        )
                    )

        blocks.extend(
            _unexposed_chartex_markdown(slide, emitted_chart_relationship_ids)
        )

        for diagram in slide_diagrams:
            if str(diagram.get("diagram_id")) in emitted_diagram_ids:
                continue
            blocks.append(smartart_to_markdown(diagram))
        if slide_index in diagram_scan_failures:
            blocks.append(
                f"> [SmartArt scan unavailable: {diagram_scan_failures[slide_index]}]"
            )

        if blocks:
            lines.append("\n\n".join(blocks))
            lines.append("")
        else:
            lines.append("_No extractable text content._")
            lines.append("")

        notes_md = extract_notes(slide)
        if notes_md:
            lines.append("### Speaker Notes")
            lines.append("")
            lines.append(notes_md)
            lines.append("")

    markdown_content = "\n".join(lines).strip() + "\n"
    out_file.write_text(markdown_content, encoding="utf-8")
    if image_manifest:
        (asset_dir / "image_manifest.json").write_text(
            json.dumps(image_manifest, ensure_ascii=False, indent=2) + "\n",
            encoding="utf-8",
        )
    profile_path = write_conversion_profile_best_effort(
        input_path=str(input_file),
        markdown_path=out_file,
        converter="ppt_to_md.py",
        conversion_type=suffix.lstrip("."),
        asset_dir=asset_dir,
        warnings=conversion_warnings,
    )

    print(f"[OK] Saved Markdown to: {out_file}")
    if profile_path:
        print(f"   Wrote conversion profile -> {profile_path}")
    if asset_dir_used:
        media_files = [
            path for path in asset_dir.iterdir()
            if path.is_file() and path.name != "image_manifest.json"
        ]
        print(f"   Extracted {len(media_files)} image file(s) -> {asset_dir}")
        if image_ref_count != len(media_files):
            print(
                f"   Deduplicated {image_ref_count} image reference(s) "
                f"into {len(media_files)} asset file(s)"
            )
        print(f"   Wrote image manifest -> {asset_dir / 'image_manifest.json'}")

    return markdown_content


def main() -> int:
    """Run the CLI entry point."""
    parser = argparse.ArgumentParser(
        description="Convert PowerPoint files to Markdown",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  python ppt_to_md.py slides.pptx
  python ppt_to_md.py slides.pptx appendix.pptx
  python ppt_to_md.py ./decks -o ./markdown
  python ppt_to_md.py slides.pptx -o output.md
  python ppt_to_md.py deck.ppsx -o notes/deck.md

Supported formats:
  .pptx  .pptm  .ppsx  .ppsm  .potx  .potm

Legacy .ppt is not parsed directly. Resave it as .pptx or export it to PDF first.
        """,
    )
    parser.add_argument("inputs", nargs="+", help="Input PowerPoint file(s) or directories")
    parser.add_argument(
        "-o",
        "--output",
        help="Output Markdown file for one input, or output directory for multiple inputs/directories",
    )

    args = parser.parse_args()

    return run_path_batch(
        args.inputs,
        set(SUPPORTED_FORMATS),
        args.output,
        lambda source, output: bool(convert_presentation_to_markdown(str(source), str(output))),
    )


if __name__ == "__main__":
    raise SystemExit(main())
